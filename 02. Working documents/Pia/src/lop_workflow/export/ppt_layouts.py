"""Per-section PowerPoint layout helpers."""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt

from lop_workflow.models.section import SectionContent, SectionId
from lop_workflow.orchestrator.state import OrchestratorState


def _layout_blank(prs: Presentation):
    """Prefer blank layout when available."""
    if len(prs.slide_layouts) > 6:
        return prs.slide_layouts[6]
    return prs.slide_layouts[1]


def _layout_title_content(prs: Presentation):
    return prs.slide_layouts[1]


def _bullets_from_section(s: SectionContent) -> list[str]:
    if s.outline_bullets:
        return [b[:500] for b in s.outline_bullets[:8]]
    lines = [ln.strip() for ln in s.body_markdown.splitlines() if ln.strip()]
    out: list[str] = []
    for ln in lines:
        if ln.startswith(("- ", "* ")):
            out.append(ln[2:].strip()[:500])
        elif ln.startswith("#"):
            continue
        elif len(out) < 8 and len(ln) < 400:
            out.append(ln)
    return out[:8] or [s.body_markdown[:400]]


def add_title_slide(prs: Presentation, st: OrchestratorState) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    t = slide.shapes.title
    if t is not None:
        t.text = st.lop.project_name if st.lop and st.lop.project_name else "LOP (draft)"
    try:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1 and st.brief:
                ph.text = f"{st.brief.client_name or ''} — Letter of Proposal"
                break
    except (ValueError, AttributeError):
        pass


def add_context_slide(prs: Presentation, s: SectionContent) -> None:
    slide = prs.slides.add_slide(_layout_title_content(prs))
    if slide.shapes.title:
        slide.shapes.title.text = "Context and objectives"
    bullets = _bullets_from_section(s)
    body_set = False
    for shp in slide.shapes:
        if not body_set and shp != slide.shapes.title and shp.has_text_frame:
            tf = shp.text_frame
            tf.text = bullets[0] if bullets else ""
            for b in bullets[1:]:
                p = tf.add_paragraph()
                p.text = b
                p.level = 0
                p.font.size = Pt(14)
            body_set = True
    if not body_set:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        box.text_frame.text = "\n".join(bullets)


def add_why_mckinsey_slide(prs: Presentation, s: SectionContent) -> None:
    slide = prs.slides.add_slide(_layout_title_content(prs))
    if slide.shapes.title:
        slide.shapes.title.text = "Why McKinsey"
    left = Inches(0.5)
    top = Inches(1.4)
    w = Inches(4.3)
    h = Inches(5)
    box1 = slide.shapes.add_textbox(left, top, w, h)
    box1.text_frame.text = "Value proposition\n" + "\n".join(_bullets_from_section(s)[:4])
    box2 = slide.shapes.add_textbox(left + w + Inches(0.2), top, w, h)
    proof = "Proof points (manifest-cited only):\n" + "\n".join(_bullets_from_section(s)[4:8] or ["— TBD / evidence required"])
    box2.text_frame.text = proof


def add_timeline_slide(prs: Presentation, s: SectionContent) -> None:
    slide = prs.slides.add_slide(_layout_blank(prs))
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_box.text_frame.text = "Timeline and team — indicative cadence"
    rows, cols = 4, 3
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(4))
    table = table_shape.table
    hdr = ("Phase / workstream", "Window", "Milestone / output")
    for c, h in enumerate(hdr):
        table.cell(0, c).text = h
    for r in range(1, rows):
        table.cell(r, 0).text = f"Phase {r}"
        table.cell(r, 1).text = "TBD"
        table.cell(r, 2).text = "TBD"
    notes = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    notes.text_frame.text = "Narrative:\n" + s.body_markdown[:1200]


def add_team_or_credentials_slide(prs: Presentation, s: SectionContent, title: str) -> None:
    slide = prs.slides.add_slide(_layout_title_content(prs))
    if slide.shapes.title:
        slide.shapes.title.text = title
    body_set = False
    for shp in slide.shapes:
        if not body_set and shp != slide.shapes.title and shp.has_text_frame:
            tf = shp.text_frame
            tf.text = "\n".join(_bullets_from_section(s))
            for p in tf.paragraphs:
                p.font.size = Pt(13)
            body_set = True


def add_market_slide(prs: Presentation, s: SectionContent) -> None:
    add_team_or_credentials_slide(prs, s, "Market trends")


def add_approach_slide(prs: Presentation, s: SectionContent) -> None:
    slide = prs.slides.add_slide(_layout_title_content(prs))
    if slide.shapes.title:
        slide.shapes.title.text = "Approach"
    bullets = _bullets_from_section(s)
    body_set = False
    for shp in slide.shapes:
        if not body_set and shp != slide.shapes.title and shp.has_text_frame:
            tf = shp.text_frame
            tf.text = bullets[0] if bullets else ""
            for i, b in enumerate(bullets[1:], start=1):
                p = tf.add_paragraph()
                p.text = b
                p.level = 0 if i < 4 else 1
                p.font.size = Pt(13)
            body_set = True


def add_fees_slide(prs: Presentation, s: SectionContent) -> None:
    slide = prs.slides.add_slide(_layout_blank(prs))
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tb.text_frame.text = "Fees — structure (numbers only if sourced)"
    ft = s.extra.get("fees_table") if isinstance(s.extra, dict) else None
    rows_data = [["Line item", "Unit", "Amount", "Notes"]]
    if ft and isinstance(ft, dict):
        rows_data.extend(list(map(str, row)) for row in ft.get("rows", []))
    else:
        rows_data.extend(
            [
                ["Program management", "month", "TBD", "BD/partner"],
                ["Workstreams", "month", "TBD", "Scope"],
            ]
        )
    rlen = min(len(rows_data), 8)
    clen = 4
    tbl = slide.shapes.add_table(rlen, clen, Inches(0.5), Inches(1.0), Inches(9), Inches(3.5)).table
    for ri in range(rlen):
        for ci in range(clen):
            val = rows_data[ri][ci] if ci < len(rows_data[ri]) else ""
            tbl.cell(ri, ci).text = str(val)

    nb = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(2))
    nb.text_frame.text = s.body_markdown[:2000]


def add_appendix_slide(prs: Presentation, s: SectionContent, title: str) -> None:
    slide = prs.slides.add_slide(_layout_title_content(prs))
    if slide.shapes.title:
        slide.shapes.title.text = title
    body_set = False
    for shp in slide.shapes:
        if not body_set and shp != slide.shapes.title and shp.has_text_frame:
            shp.text_frame.text = s.body_markdown[:4500]
            body_set = True


def add_generic_slide(prs: Presentation, s: SectionContent) -> None:
    slide = prs.slides.add_slide(_layout_title_content(prs))
    if slide.shapes.title is not None:
        slide.shapes.title.text = s.section_id.value.replace("_", " ").title()
    body_set = False
    for shp in slide.shapes:
        if not body_set and shp != slide.shapes.title and shp.has_text_frame:
            shp.text_frame.text = s.body_markdown[:4000]
            for p in shp.text_frame.paragraphs:
                p.font.size = Pt(14)
            body_set = True
    if not body_set:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        box.text_frame.text = s.body_markdown[:4000]


def add_section_slides(prs: Presentation, st: OrchestratorState, s: SectionContent) -> None:
    sid = s.section_id
    if sid == SectionId.CONTEXT:
        add_context_slide(prs, s)
    elif sid == SectionId.WHY_MCKINSEY:
        add_why_mckinsey_slide(prs, s)
    elif sid == SectionId.TIMELINE_TEAM:
        add_timeline_slide(prs, s)
    elif sid == SectionId.TEAM:
        add_team_or_credentials_slide(prs, s, "Team")
    elif sid == SectionId.CREDENTIALS:
        add_team_or_credentials_slide(prs, s, "Credentials")
    elif sid == SectionId.MARKET:
        add_market_slide(prs, s)
    elif sid == SectionId.APPROACH:
        add_approach_slide(prs, s)
    elif sid == SectionId.FEES:
        add_fees_slide(prs, s)
    elif sid == SectionId.APPENDIX:
        add_appendix_slide(prs, s, "Appendix")
    elif sid == SectionId.REFERENCES:
        add_appendix_slide(prs, s, "References")
    elif sid == SectionId.TEAM_CVS:
        add_appendix_slide(prs, s, "Team CVs")
    else:
        add_generic_slide(prs, s)
