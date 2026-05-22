"""Import an existing PPTX deck to editable HTML slides.

Reads each slide's shapes with python-pptx, classifies them as chrome
(from the slide master/layout) or content (slide-specific), and generates
one HTML file per slide using the project's tagging conventions.

Usage:
    python pptx_to_html.py <input.pptx> [-o <output_dir>]

    <input.pptx>   Path to the PPTX file to import.
    -o             Output directory for HTML files and extracted images.
                   Default: <input_name>_html/ next to the source file.

Output structure:
    <output_dir>/
        slide_001.html
        slide_002.html
        ...
        images/
            slide_001_image_1.png
            ...

Each HTML file links to ../css/slide.css (or an appropriate relative path)
and uses the data-pptx tagging convention so the slides can be edited and
re-compiled with html_to_pptx.py.
"""

from __future__ import annotations

import argparse
import os
import sys
import zipfile
from html import escape
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from lxml import etree
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE

from shared.constants import (
    CONTENT_BOTTOM_PX,
    CSS_PATH,
    FOOTNOTE_BOTTOM_PX,
    LAYOUT_NAMES,
    LEFT_MARGIN_PX,
    PPTX_TAGS,
    TITLE_BOTTOM_PX,
)
from shared.emu import VIEWPORT_H, emu_to_pt, emu_to_px_x, emu_to_px_y, pt_to_css_px

try:
    from pptx.shapes.connector import Connector as _ConnectorCls
except ImportError:
    _ConnectorCls = None


# ---------------------------------------------------------------------------
# Theme color resolution
# ---------------------------------------------------------------------------

_THEME_CLR_MAP = {
    MSO_THEME_COLOR.DARK_1: "dk1",
    MSO_THEME_COLOR.LIGHT_1: "lt1",
    MSO_THEME_COLOR.DARK_2: "dk2",
    MSO_THEME_COLOR.LIGHT_2: "lt2",
    MSO_THEME_COLOR.ACCENT_1: "accent1",
    MSO_THEME_COLOR.ACCENT_2: "accent2",
    MSO_THEME_COLOR.ACCENT_3: "accent3",
    MSO_THEME_COLOR.ACCENT_4: "accent4",
    MSO_THEME_COLOR.ACCENT_5: "accent5",
    MSO_THEME_COLOR.ACCENT_6: "accent6",
    MSO_THEME_COLOR.HYPERLINK: "hlink",
    MSO_THEME_COLOR.FOLLOWED_HYPERLINK: "folHlink",
}
# TEXT_1/2, BACKGROUND_1/2 alias DARK/LIGHT
_THEME_CLR_MAP[13] = "dk1"   # TEXT_1
_THEME_CLR_MAP[14] = "dk2"   # TEXT_2
_THEME_CLR_MAP[15] = "lt1"   # BACKGROUND_1
_THEME_CLR_MAP[16] = "lt2"   # BACKGROUND_2


def _extract_theme_colors(pptx_path: str) -> dict[str, str]:
    """Extract the theme color scheme from a PPTX as {name: '#RRGGBB'}."""
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    colors: dict[str, str] = {}
    try:
        zf = zipfile.ZipFile(pptx_path)
        for name in zf.namelist():
            if "theme1.xml" in name:
                root = etree.fromstring(zf.read(name))
                clrScheme = root.find(".//a:clrScheme", ns)
                if clrScheme is not None:
                    for child in clrScheme:
                        tag = child.tag.split("}")[-1]
                        for sub in child:
                            val = sub.get("val") or sub.get("lastClr")
                            if val:
                                colors[tag] = f"#{val}"
                break
    except Exception:
        pass
    return colors


def _extract_default_fonts(pptx_path: str) -> dict[str, tuple[str, float, bool]]:
    """Read txStyles from the slide master to get default font name/size/bold per style.

    Returns a dict with keys 'title', 'body', 'other' -> (font_name, size_pt, bold).
    """
    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    ns = {"p": ns_p, "a": ns_a}

    defaults: dict[str, tuple[str, float, bool]] = {
        "title": ("Georgia", 25.0, True),
        "body": ("Arial", 16.0, False),
        "other": ("Arial", 18.0, False),
    }

    # +mj-lt / +mn-lt resolution
    theme_fonts: dict[str, str] = {"major": "Georgia", "minor": "Arial"}
    try:
        zf = zipfile.ZipFile(pptx_path)
        for name in zf.namelist():
            if "theme1.xml" in name:
                root = etree.fromstring(zf.read(name))
                for tag, key in [("majorFont", "major"), ("minorFont", "minor")]:
                    el = root.find(f".//a:{tag}/a:latin", ns)
                    if el is not None:
                        tf = el.get("typeface")
                        if tf:
                            theme_fonts[key] = tf
                break
    except Exception:
        pass

    def _resolve_typeface(raw: str | None) -> str:
        if raw == "+mj-lt":
            return theme_fonts["major"]
        if raw == "+mn-lt":
            return theme_fonts["minor"]
        return raw or theme_fonts["minor"]

    try:
        zf = zipfile.ZipFile(pptx_path)
        for name in zf.namelist():
            if "slideMaster1.xml" in name:
                root = etree.fromstring(zf.read(name))
                txStyles = root.find("p:txStyles", ns)
                if txStyles is None:
                    break
                for style_key, xml_tag in [
                    ("title", "titleStyle"),
                    ("body", "bodyStyle"),
                    ("other", "otherStyle"),
                ]:
                    style_el = txStyles.find(f"p:{xml_tag}", ns)
                    if style_el is None:
                        continue
                    pPr = style_el.find("a:lvl1pPr", ns)
                    if pPr is None:
                        pPr = style_el.find("a:defPPr", ns)
                    if pPr is None:
                        continue
                    defRPr = pPr.find("a:defRPr", ns)
                    if defRPr is None:
                        continue
                    sz = defRPr.get("sz")
                    latin = defRPr.find("a:latin", ns)
                    typeface = latin.get("typeface") if latin is not None else None
                    font_name = _resolve_typeface(typeface)
                    font_size = float(sz) / 100.0 if sz else defaults[style_key][1]
                    bold = defRPr.get("b") == "1"
                    defaults[style_key] = (font_name, font_size, bold)
                break
    except Exception:
        pass

    return defaults


def _resolve_theme_color(theme_color_enum, theme_colors: dict[str, str]) -> str | None:
    """Resolve a MSO_THEME_COLOR enum to a hex string using the theme palette."""
    try:
        key = _THEME_CLR_MAP.get(int(theme_color_enum))
    except (TypeError, ValueError):
        key = _THEME_CLR_MAP.get(theme_color_enum)
    if key and key in theme_colors:
        return theme_colors[key]
    return None


# ---------------------------------------------------------------------------
# Zone lookup
# ---------------------------------------------------------------------------

_ZONE_TABLE: list[tuple[str, float, float, str]] = [
    ("title",    0,                  TITLE_BOTTOM_PX,    "title"),
    ("content",  TITLE_BOTTOM_PX,    CONTENT_BOTTOM_PX,  "content"),
    ("footnote", CONTENT_BOTTOM_PX,  FOOTNOTE_BOTTOM_PX, "footnote"),
    ("source",   FOOTNOTE_BOTTOM_PX, float(VIEWPORT_H),  "source"),
]

_ZONE_LEFT = LEFT_MARGIN_PX


def _classify_zone(top_px: float, height_px: float) -> tuple[str, float, str]:
    """Return (zone_name, zone_top_px, css_class) for a shape's midpoint."""
    mid_y = top_px + height_px / 2
    for name, ztop, zbot, css_cls in _ZONE_TABLE:
        if ztop <= mid_y < zbot:
            return name, ztop, css_cls
    return "content", TITLE_BOTTOM_PX, "content"


# ---------------------------------------------------------------------------
# Shape classification
# ---------------------------------------------------------------------------

def _is_connector(shape) -> bool:
    if _ConnectorCls is not None and isinstance(shape, _ConnectorCls):
        return True
    try:
        return shape._element.tag.endswith("}cxnSp")
    except Exception:
        return False


def _get_pptx_tag(shape) -> str:
    """Map a python-pptx shape to its data-pptx tag value."""
    if _is_connector(shape):
        return "line"
    try:
        ast = shape.auto_shape_type
        if ast == MSO_SHAPE.RECTANGLE:
            return "rect"
        if ast == MSO_SHAPE.OVAL:
            return "ellipse"
        return "rect"
    except Exception:
        pass
    if shape.has_text_frame:
        return "textbox"
    try:
        _ = shape.image
        return "image"
    except Exception:
        pass
    return "rect"


def _is_chrome(
    shape,
    layout_shapes: dict[str, str],
    master_shapes: dict[str, str],
) -> bool:
    """A shape is chrome if it exists on the layout/master AND its text hasn't changed."""
    if shape.is_placeholder:
        return False
    name = shape.name
    shape_text = shape.text_frame.text.strip() if shape.has_text_frame else ""

    # Check layout shapes
    if name in layout_shapes:
        if not shape_text or shape_text == layout_shapes[name]:
            return True
    # Check master shapes
    if name in master_shapes:
        if not shape_text or shape_text == master_shapes[name]:
            return True
    return False


# ---------------------------------------------------------------------------
# Property extraction
# ---------------------------------------------------------------------------

def _fill_hex(shape, theme_colors: dict[str, str]) -> str | None:
    try:
        fill = shape.fill
        if fill.type is None:
            return None
        fc = fill.fore_color
        # Try direct RGB first
        try:
            return f"#{fc.rgb}"
        except AttributeError:
            pass
        # Resolve theme color
        try:
            tc = fc.theme_color
            resolved = _resolve_theme_color(tc, theme_colors)
            if resolved:
                return resolved
        except Exception:
            pass
    except Exception:
        pass
    return None


def _font_color_hex(font, theme_colors: dict[str, str]) -> str | None:
    """Extract font color, resolving theme colors if needed."""
    try:
        fc = font.color
        try:
            rgb = fc.rgb
            if rgb is not None:
                return f"#{rgb}"
        except AttributeError:
            pass
        try:
            tc = fc.theme_color
            resolved = _resolve_theme_color(tc, theme_colors)
            if resolved:
                return resolved
        except Exception:
            pass
    except Exception:
        pass
    return None


def _stroke_props(
    shape, theme_colors: dict[str, str],
) -> tuple[str | None, float, str]:
    """Return (color_hex, width_pt, dash_css).

    *dash_css* is a CSS border-style keyword: 'solid', 'dashed', or 'dotted'.
    Returns (None, ..., ...) when no visible outline.
    """
    color: str | None = None
    width = 1.0
    dash_css = "solid"
    try:
        line = shape.line
        if line.fill.type is not None and str(line.fill.type) != "BACKGROUND (5)":
            try:
                rgb = line.color.rgb
                if rgb is not None:
                    color = f"#{rgb}"
            except (AttributeError, TypeError):
                try:
                    tc = line.color.theme_color
                    color = _resolve_theme_color(tc, theme_colors)
                except Exception:
                    pass
            try:
                ds = line.dash_style
                if ds is not None and ds != 1:
                    from pptx.enum.dml import MSO_LINE_DASH_STYLE
                    if ds in (
                        MSO_LINE_DASH_STYLE.SQUARE_DOT,
                        MSO_LINE_DASH_STYLE.ROUND_DOT,
                        MSO_LINE_DASH_STYLE.DASH_DOT,
                    ):
                        dash_css = "dashed"
                    elif ds == MSO_LINE_DASH_STYLE.DASH:
                        dash_css = "dashed"
                    elif ds in (
                        MSO_LINE_DASH_STYLE.LONG_DASH,
                        MSO_LINE_DASH_STYLE.LONG_DASH_DOT,
                        MSO_LINE_DASH_STYLE.LONG_DASH_DOT_DOT,
                    ):
                        dash_css = "dashed"
            except Exception:
                pass
    except Exception:
        pass
    try:
        w = shape.line.width
        if w:
            width = emu_to_pt(w)
    except Exception:
        pass
    return color, width, dash_css


_NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# Wingdings / symbol font chars that need mapping to Unicode for HTML
_WINGDINGS_TO_UNICODE: dict[str, str] = {
    "\x9f": "\u2022",  # Wingdings 159 -> bullet •
    "\xfc": "\u2714",  # Wingdings 252 -> checkmark ✔
    "\xd8": "\u25cf",  # Wingdings 216 -> black circle ●
}


def _get_lstStyle_props(shape, para_level: int) -> dict:
    """Read per-shape txBody/lstStyle defaults for the given paragraph level.

    Returns dict with possible keys:
        sz_pt         - font size in points (float)
        bold          - True/False
        italic        - True/False
        font_name     - typeface string
        bu_char       - resolved bullet character (Unicode)
        line_spacing   - fraction, e.g. 1.0 for 100%
    All values may be absent if not defined in lstStyle.
    """
    result: dict = {}
    try:
        lstStyle = shape._element.find(f".//{_NS_A}lstStyle")
    except Exception:
        return result
    if lstStyle is None or len(lstStyle) == 0:
        return result

    # OOXML: lvl1pPr = paragraph level 0, lvl2pPr = level 1, etc.
    tag = f"{_NS_A}lvl{para_level + 1}pPr"
    pPr = lstStyle.find(tag)
    if pPr is None:
        pPr = lstStyle.find(f"{_NS_A}defPPr")
    if pPr is None:
        return result

    defRPr = pPr.find(f"{_NS_A}defRPr")
    if defRPr is not None:
        sz = defRPr.get("sz")
        if sz:
            result["sz_pt"] = float(sz) / 100.0
        b = defRPr.get("b")
        if b is not None:
            result["bold"] = b == "1"
        i = defRPr.get("i")
        if i is not None:
            result["italic"] = i == "1"
        latin = defRPr.find(f"{_NS_A}latin")
        if latin is not None:
            tf = latin.get("typeface")
            if tf and not tf.startswith("+"):
                result["font_name"] = tf

    # Bullet character
    buChar = pPr.find(f"{_NS_A}buChar")
    buNone = pPr.find(f"{_NS_A}buNone")
    if buChar is not None and buNone is None:
        raw_char = buChar.get("char", "")
        buFont = pPr.find(f"{_NS_A}buFont")
        bu_typeface = buFont.get("typeface", "") if buFont is not None else ""
        if bu_typeface.lower() == "wingdings":
            raw_char = _WINGDINGS_TO_UNICODE.get(raw_char, raw_char)
        if raw_char and raw_char != "\u200b":  # skip zero-width space
            result["bu_char"] = raw_char

    # Line spacing
    lnSpc = pPr.find(f"{_NS_A}lnSpc")
    if lnSpc is not None:
        spcPct = lnSpc.find(f"{_NS_A}spcPct")
        if spcPct is not None:
            val = spcPct.get("val")
            if val:
                result["line_spacing"] = float(val) / 100000.0

    return result


def _detect_bullet(para, lst: dict) -> str | None:
    """Return the bullet character for *para*, or None if not a bullet.

    Checks the paragraph's own pPr first, then falls back to lstStyle.
    """
    pf = para._pPr
    if pf is not None:
        buNone = pf.find(f"{_NS_A}buNone")
        buChar = pf.find(f"{_NS_A}buChar")
        buAutoNum = pf.find(f"{_NS_A}buAutoNum")
        if buChar is not None:
            return buChar.get("char", "\u2022")
        if buAutoNum is not None:
            return "\u2022"
        if buNone is not None:
            return None
        # No explicit bullet info -- check lstStyle
        lst_bu = lst.get("bu_char")
        if lst_bu:
            return lst_bu
        if para.level > 0:
            return "\u2013"
        return None
    # No pPr at all -- check lstStyle
    lst_bu = lst.get("bu_char")
    return lst_bu or None


def _render_runs(para, lst: dict, default_font, theme_colors) -> list[str]:
    """Build a list of <span> HTML strings for the runs of *para*."""
    spans: list[str] = []
    for run in para.runs:
        css: list[str] = []

        font_name = run.font.name or lst.get("font_name") or default_font[0]
        css.append(f"font-family:{font_name}")

        if run.font.size:
            css.append(f"font-size:{pt_to_css_px(emu_to_pt(run.font.size)):.1f}px")
        else:
            sz_pt = lst.get("sz_pt") or default_font[1]
            css.append(f"font-size:{pt_to_css_px(sz_pt):.1f}px")

        bold = run.font.bold
        if bold is None:
            bold = lst.get("bold")
        if bold is None:
            bold = default_font[2] if len(default_font) > 2 else False
        if bold:
            css.append("font-weight:bold")

        if run.font.italic:
            css.append("font-style:italic")
        color_hex = _font_color_hex(run.font, theme_colors)
        if color_hex:
            css.append(f"color:{color_hex}")

        style_attr = f' style="{"; ".join(css)}"' if css else ""

        rPr = run._r.find(f"{_NS_A}rPr")
        baseline = rPr.get("baseline") if rPr is not None else None
        text_escaped = escape(run.text)
        if baseline and baseline != "0":
            bl = int(baseline)
            if bl > 0:
                text_escaped = f"<sup>{text_escaped}</sup>"
            else:
                text_escaped = f"<sub>{text_escaped}</sub>"

        spans.append(f"<span{style_attr}>{text_escaped}</span>")
    if not spans and para.text.strip():
        spans.append(f"<span>{escape(para.text)}</span>")
    return spans


def _text_to_html(
    text_frame,
    theme_colors: dict[str, str],
    default_font: tuple[str, float, bool] = ("Arial", 18.0, False),
    shape=None,
) -> str:
    """Render a TextFrame as HTML markup.

    Bullet paragraphs are grouped into <ul>/<li> elements.
    Non-bullet paragraphs are emitted as <p>.
    """
    parts: list[str] = []
    in_ul = False

    for para in text_frame.paragraphs:
        lst = _get_lstStyle_props(shape, para.level) if shape else {}
        bullet_char = _detect_bullet(para, lst)

        # Determine the paragraph's effective font size (for the CSS strut).
        # Use the first run's explicit size, else lstStyle, else template default.
        para_sz_pt: float | None = None
        for run in para.runs:
            if run.font.size:
                para_sz_pt = emu_to_pt(run.font.size)
                break
        if para_sz_pt is None:
            para_sz_pt = lst.get("sz_pt") or default_font[1]

        p_style_parts: list[str] = [f"font-size:{pt_to_css_px(para_sz_pt):.1f}px"]

        # OOXML spcPct 100% = single spacing ≈ CSS line-height 1.2.
        # The .slide default is 1.2, so only emit an override for non-single.
        ls = lst.get("line_spacing")  # raw ratio: 1.0 = single, 1.5 = 1.5×
        if ls is not None and ls != 1.0:
            p_style_parts.append(f"line-height:{ls * 1.2:.2f}")

        spans = _render_runs(para, lst, default_font, theme_colors)
        if not spans:
            if in_ul:
                parts.append("</ul>")
                in_ul = False
            continue

        inner = "".join(spans)
        style_attr = f' style="{"; ".join(p_style_parts)}"'

        if bullet_char:
            if not in_ul:
                bc_escaped = escape(bullet_char)
                parts.append(f'<ul data-bullet-char="{bc_escaped}">')
                in_ul = True
            parts.append(f"<li{style_attr}>{inner}</li>")
        else:
            if in_ul:
                parts.append("</ul>")
                in_ul = False
            parts.append(f"<p{style_attr}>{inner}</p>")

    if in_ul:
        parts.append("</ul>")
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# Shape → HTML
# ---------------------------------------------------------------------------

def _connector_html(shape, zone_top: float, zone_left: float,
                    theme_colors: dict[str, str]) -> str:
    try:
        x1 = emu_to_px_x(shape.begin_x) - zone_left
        y1 = emu_to_px_y(shape.begin_y) - zone_top
        x2 = emu_to_px_x(shape.end_x) - zone_left
        y2 = emu_to_px_y(shape.end_y) - zone_top
    except Exception:
        x1 = emu_to_px_x(shape.left) - zone_left
        y1 = emu_to_px_y(shape.top) - zone_top
        x2 = x1 + emu_to_px_x(shape.width)
        y2 = y1 + emu_to_px_y(shape.height)

    stroke_color, stroke_w_pt, _dash = _stroke_props(shape, theme_colors)
    stroke = stroke_color or "#000000"
    stroke_w_px = pt_to_css_px(stroke_w_pt)

    svg_l = min(x1, x2)
    svg_t = min(y1, y2)
    svg_w = max(abs(x2 - x1), 1)
    svg_h = max(abs(y2 - y1), stroke_w_px + 2)

    lx1, ly1 = x1 - svg_l, y1 - svg_t
    lx2, ly2 = x2 - svg_l, y2 - svg_t

    return (
        f'<svg style="position:absolute; left:{svg_l:.1f}px; top:{svg_t:.1f}px; '
        f'width:{svg_w:.1f}px; height:{svg_h:.1f}px; overflow:visible">'
        f'<line data-pptx="line" x1="{lx1:.1f}" y1="{ly1:.1f}" '
        f'x2="{lx2:.1f}" y2="{ly2:.1f}" '
        f'stroke="{stroke}" stroke-width="{stroke_w_px:.1f}"/>'
        f"</svg>"
    )


def _picture_html(
    shape, slide_num: int, images_dir: Path,
    left: float, top: float, w: float, h: float,
) -> str:
    images_dir.mkdir(parents=True, exist_ok=True)
    try:
        blob = shape.image.blob
        ext = shape.image.ext or "png"
    except Exception:
        return ""
    fname = f"slide_{slide_num:03d}_image_{shape.shape_id}.{ext}"
    (images_dir / fname).write_bytes(blob)
    style = (
        f"position:absolute; left:{left:.1f}px; top:{top:.1f}px; "
        f"width:{w:.1f}px; height:{h:.1f}px"
    )
    return f'<img data-pptx="image" src="images/{fname}" style="{style}"/>'


def _body_anchor(shape) -> str | None:
    """Read the vertical text anchor from bodyPr: 't', 'ctr', or 'b'."""
    try:
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        body_pr = shape._element.find(".//a:bodyPr", ns)
        if body_pr is not None:
            return body_pr.get("anchor")
    except Exception:
        pass
    return None


def _shape_html(
    shape, tag: str, slide_num: int, images_dir: Path,
    zone_top: float, zone_left: float,
    theme_colors: dict[str, str],
    default_fonts: dict[str, tuple[str, float]] | None = None,
) -> str:
    """Convert a single shape to an HTML string with zone-relative positioning."""
    left = emu_to_px_x(shape.left) - zone_left
    top = emu_to_px_y(shape.top) - zone_top
    w = emu_to_px_x(shape.width)
    h = emu_to_px_y(shape.height)

    if tag == "line":
        return _connector_html(shape, zone_top, zone_left, theme_colors)
    if tag == "image":
        return _picture_html(shape, slide_num, images_dir, left, top, w, h)

    style_parts = [
        "position:absolute",
        f"left:{left:.1f}px",
        f"top:{top:.1f}px",
        f"width:{w:.1f}px",
        f"height:{h:.1f}px",
    ]
    fill = _fill_hex(shape, theme_colors)
    if fill:
        style_parts.append(f"background:{fill}")
    if tag == "ellipse":
        style_parts.append("border-radius:50%")

    stroke_color, stroke_w_pt, stroke_dash = _stroke_props(shape, theme_colors)
    if stroke_color:
        style_parts.append(f"border:{pt_to_css_px(stroke_w_pt):.1f}px {stroke_dash} {stroke_color}")
        style_parts.append("box-sizing:border-box")

    if shape.has_text_frame:
        tf = shape.text_frame
        pad_t = emu_to_px_y(tf.margin_top) if tf.margin_top else 0
        pad_r = emu_to_px_x(tf.margin_right) if tf.margin_right else 0
        pad_b = emu_to_px_y(tf.margin_bottom) if tf.margin_bottom else 0
        pad_l = emu_to_px_x(tf.margin_left) if tf.margin_left else 0
        if pad_t or pad_r or pad_b or pad_l:
            style_parts.append(
                f"padding:{pad_t:.1f}px {pad_r:.1f}px {pad_b:.1f}px {pad_l:.1f}px"
            )

    attrs = f'data-pptx="{tag}"'
    if tag == "placeholder" and shape.is_placeholder:
        attrs += f' data-ph-idx="{shape.placeholder_format.idx}"'

    anchor = _body_anchor(shape)
    if anchor == "b":
        style_parts.extend([
            "display:flex",
            "flex-direction:column",
            "justify-content:flex-end",
        ])
    elif anchor == "ctr":
        style_parts.extend([
            "display:flex",
            "flex-direction:column",
            "justify-content:center",
        ])

    inner = ""
    if shape.has_text_frame:
        df = default_fonts or {}
        if tag == "placeholder" and shape.is_placeholder:
            idx = shape.placeholder_format.idx
            font_key = "title" if idx in (0, 1) else "body"
        else:
            font_key = "other"
        dfont = df.get(font_key, ("Arial", 18.0, False))
        inner = _text_to_html(shape.text_frame, theme_colors, default_font=dfont, shape=shape)

    return f'<div {attrs} style="{"; ".join(style_parts)}">\n{inner}\n</div>'


# ---------------------------------------------------------------------------
# Page assembly
# ---------------------------------------------------------------------------

def _build_page(
    layout_name: str,
    css_href: str,
    zone_shapes: dict[str, list[str]],
    chrome_parts: list[str],
) -> str:
    zone_order = [
        ("title", "title"),
        ("content", "content"),
        ("footnote", "footnote"),
        ("source", "source"),
    ]
    zone_sections: list[str] = []
    for zname, zcls in zone_order:
        inner = "\n".join(zone_shapes.get(zname, []))
        zone_sections.append(f'  <div class="{zcls}">\n{inner}\n  </div>')

    chrome_section = ""
    if chrome_parts:
        chrome_inner = "\n".join(chrome_parts)
        chrome_section = f'  <div class="chrome">\n{chrome_inner}\n  </div>\n'

    body = "\n".join(zone_sections)
    return (
        "<!DOCTYPE html>\n"
        '<html lang="en">\n'
        "<head>\n"
        '<meta charset="UTF-8">\n'
        '<meta name="viewport" content="width=1920">\n'
        f'<link rel="stylesheet" href="{css_href}">\n'
        "</head>\n"
        "<body>\n"
        f'<section class="slide" data-layout="{escape(layout_name)}">\n'
        f"{body}\n"
        f"{chrome_section}"
        "</section>\n"
        "</body>\n"
        "</html>\n"
    )


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def pptx_to_html(
    pptx_path: Path,
    output_dir: Path,
) -> list[Path]:
    """Convert a PPTX presentation to a set of HTML slide files.

    Args:
        pptx_path: Path to the source .pptx file.
        output_dir: Directory where HTML files and images are written.

    Returns:
        Sorted list of generated HTML file paths.
    """
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    images_dir = output_dir / "images"

    prs = Presentation(str(pptx_path))
    theme_colors = _extract_theme_colors(str(pptx_path))
    default_fonts = _extract_default_fonts(str(pptx_path))

    try:
        css_href = os.path.relpath(CSS_PATH, output_dir).replace("\\", "/")
    except ValueError:
        css_href = "slide.css"

    html_files: list[Path] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        layout_name = slide.slide_layout.name

        # Build name->text maps for chrome detection
        layout_shapes: dict[str, str] = {}
        for s in slide.slide_layout.shapes:
            if not s.is_placeholder:
                txt = s.text_frame.text.strip() if s.has_text_frame else ""
                layout_shapes[s.name] = txt
        master_shapes: dict[str, str] = {}
        for s in slide.slide_layout.slide_master.shapes:
            if not s.is_placeholder:
                txt = s.text_frame.text.strip() if s.has_text_frame else ""
                master_shapes[s.name] = txt

        zone_shapes: dict[str, list[str]] = {
            "title": [], "subtitle": [], "content": [],
            "footnote": [], "source": [],
        }
        chrome_parts: list[str] = []

        for shape in slide.shapes:
            if _is_chrome(shape, layout_shapes, master_shapes):
                left = emu_to_px_x(shape.left)
                top = emu_to_px_y(shape.top)
                w = emu_to_px_x(shape.width)
                h = emu_to_px_y(shape.height)
                style = (
                    f"position:absolute; left:{left:.1f}px; top:{top:.1f}px; "
                    f"width:{w:.1f}px; height:{h:.1f}px"
                )
                inner = ""
                if shape.has_text_frame:
                    inner = _text_to_html(
                        shape.text_frame, theme_colors,
                        default_font=default_fonts.get("other", ("Arial", 18.0, False)),
                        shape=shape,
                    )
                chrome_parts.append(
                    f'  <div data-pptx="chrome" style="{style}">'
                    f"\n{inner}\n  </div>"
                )
                continue

            tag = "placeholder" if shape.is_placeholder else _get_pptx_tag(shape)

            top_px = emu_to_px_y(shape.top)
            h_px = emu_to_px_y(shape.height)
            zname, ztop, _zcls = _classify_zone(top_px, h_px)

            html = _shape_html(
                shape, tag, slide_idx, images_dir, ztop, _ZONE_LEFT,
                theme_colors, default_fonts,
            )
            zone_shapes[zname].append(html)

        page = _build_page(layout_name, css_href, zone_shapes, chrome_parts)

        filename = f"slide_{slide_idx:03d}.html"
        filepath = output_dir / filename
        filepath.write_text(page, encoding="utf-8")
        html_files.append(filepath)

    return sorted(html_files)


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Import a PPTX deck to editable HTML slides",
    )
    parser.add_argument(
        "input",
        type=Path,
        help="Path to the .pptx file to import",
    )
    parser.add_argument(
        "-o", "--output",
        type=Path,
        default=None,
        help="Output directory (default: <name>_html/ next to input)",
    )
    args = parser.parse_args()

    if not args.input.is_file():
        print(f"Error: {args.input} not found", file=sys.stderr)
        sys.exit(1)

    if args.output is None:
        base = args.input.stem
        args.output = args.input.parent / f"{base}_html"

    args.output.mkdir(parents=True, exist_ok=True)
    print(f"Importing {args.input} -> {args.output}/")

    html_files = pptx_to_html(args.input, args.output)
    print(f"Done. {len(html_files)} slide(s) written to {args.output}/")


if __name__ == "__main__":
    main()
