"""Gold-standard builder: corpus paths, extraction, catalog, drafting helpers, coach, exports."""

from __future__ import annotations

import tempfile
from pathlib import Path

from pptx import Presentation

from lop_workflow.agents.coach import evaluate_lop_coach
from lop_workflow.agents.jasper_prompts import jasper_prompt_filename, load_jasper_chapter_prompt
from lop_workflow.agents.llm_client import corpus_excerpt_for_section
from lop_workflow.corpus.background_material import classify_background_source_type, load_background_sources
from lop_workflow.corpus.document_extract import ensure_extracted_markdown, extract_text_from_pptx
from lop_workflow.export.html_export import _markdown_table_to_html
from lop_workflow.export.ppt_export import write_pptx
from lop_workflow.gold_patterns.loader import get_section_pattern, load_gold_catalog
from lop_workflow.models import LOPCategory
from lop_workflow.models.brief import Brief
from lop_workflow.models.facts import FactsRegistry
from lop_workflow.models.section import LOPDocument, SectionContent, SectionId
from lop_workflow.models.source import SourceRef
from lop_workflow.orchestrator.graph import LopWorkflow, NoOpAgents
from lop_workflow.orchestrator.state import HumanCheckpoint, OrchestratorState, Phase


def test_background_material_dir_env_override(monkeypatch, tmp_path):
    monkeypatch.setenv("LOP_BACKGROUND_MATERIAL_DIR", str(tmp_path))
    from lop_workflow.workspace_paths import background_material_dir

    assert background_material_dir() == tmp_path


def test_background_material_prefers_numbered_folder(monkeypatch, tmp_path):
    monkeypatch.setattr(
        "lop_workflow.workspace_paths.documents_workspace_root",
        lambda: tmp_path,
    )
    (tmp_path / "01. Background material").mkdir(parents=True)
    from lop_workflow.workspace_paths import background_material_dir

    assert background_material_dir() == tmp_path / "01. Background material"


def test_background_material_legacy_fallback(monkeypatch, tmp_path):
    monkeypatch.setattr(
        "lop_workflow.workspace_paths.documents_workspace_root",
        lambda: tmp_path,
    )
    (tmp_path / "Background Material").mkdir(parents=True)
    from lop_workflow.workspace_paths import background_material_dir

    assert background_material_dir() == tmp_path / "Background Material"


def test_classify_source_types():
    p = Path("928635 Gold Standard GEM LOP Guide.pptx")
    assert classify_background_source_type(p) == "gold_lop"
    assert classify_background_source_type(Path("22-05-2021 - Proposal X signed.pdf")) == "prior_proposal"


def test_extract_pptx_roundtrip(tmp_path):
    ppt = tmp_path / "t.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    assert slide.shapes.title is not None
    slide.shapes.title.text = "TitleAlpha"
    prs.save(str(ppt))
    text = extract_text_from_pptx(ppt)
    assert "TitleAlpha" in text


def test_ensure_extracted_markdown_cache(tmp_path):
    ppt = tmp_path / "doc.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    assert slide.shapes.title is not None
    slide.shapes.title.text = "CachedSlide"
    prs.save(str(ppt))
    ext_root = tmp_path / "extracted"
    out1 = ensure_extracted_markdown(ppt, ext_root)
    assert out1 and out1.is_file()
    mtime = out1.stat().st_mtime
    out2 = ensure_extracted_markdown(ppt, ext_root)
    assert out2 == out1
    assert out2.stat().st_mtime == mtime
    assert "CachedSlide" in out2.read_text(encoding="utf-8")


def test_load_background_sources_indexes_pptx(tmp_path):
    ppt = tmp_path / "Gold LOF Guide.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    assert slide.shapes.title is not None
    slide.shapes.title.text = "SnippetTestXYZ"
    prs.save(str(ppt))
    ext = tmp_path / "extracted"
    refs = load_background_sources(tmp_path, extracted_dir=ext)
    assert len(refs) == 1
    assert "SnippetTestXYZ" in refs[0].snippet


def test_load_gold_catalog_package():
    cat = load_gold_catalog()
    assert cat.version >= 1
    assert "context_objectives" in cat.sections


def test_get_section_pattern_competitive_overlay():
    p = get_section_pattern(SectionId.WHY_MCKINSEY, LOPCategory.COMPETITIVE)
    hints = " ".join(p.structure_hints + p.win_logic).lower()
    assert "differentiat" in hints or "differentiation" in hints


def test_corpus_excerpt_orders_by_priority():
    corpus = [
        SourceRef(source_id="a", source_type="capability_deck", title="Cap", snippet="cap"),
        SourceRef(source_id="b", source_type="gold_lop", title="Gold", snippet="gold"),
    ]
    ex = corpus_excerpt_for_section(SectionId.WHY_MCKINSEY, corpus)
    assert ex.index("gold") < ex.index("cap")


def test_jasper_prompt_filename_maps_sections():
    assert jasper_prompt_filename(SectionId.FEES) == "chapter-fees.md"
    assert jasper_prompt_filename(SectionId.APPENDIX) == "chapter-appendix-references.md"


def test_load_jasper_prompt_from_tmp(tmp_path):
    (tmp_path / "chapter-fees.md").write_text("# Fees\nTestPromptBody", encoding="utf-8")
    text = load_jasper_chapter_prompt(SectionId.FEES, prompts_root=str(tmp_path))
    assert "TestPromptBody" in text


def test_coach_rubric_and_issue_table():
    lop = LOPDocument(
        run_id="r1",
        project_name="P",
        sections=[
            SectionContent(
                section_id=SectionId.CONTEXT,
                outline_bullets=[],
                body_markdown="Short",
                sources=[],
                confidence=0.8,
            ),
            SectionContent(
                section_id=SectionId.FEES,
                outline_bullets=[],
                body_markdown="TBD pending approval",
                sources=[SourceRef(source_id="s1", source_type="gold_lop", title="g", snippet="x")],
                confidence=0.8,
                extra={"fees_table": {"headers": ["A"], "rows": [["b"]]}},
            ),
        ],
    )
    brief = Brief(run_id="r1", rfp_tender_references=[SourceRef(source_id="rfp1", source_type="rfp_tender", title="RFP", snippet="ask")])
    rep = evaluate_lop_coach(lop, FactsRegistry(), background=[], brief=brief, lop_category=LOPCategory.COMPETITIVE)
    assert rep.issue_table_markdown
    assert "| Chapter |" in rep.issue_table_markdown
    assert "rubric_rfp_fit" in rep.rubric_dimensions
    assert "rubric_gold_pattern_fit" in rep.rubric_dimensions


def test_markdown_table_to_html():
    md = "| A | B |\n| --- | --- |\n| x | y |\n"
    html = _markdown_table_to_html(md)
    assert "<table" in html
    assert "x" in html


def test_ppt_export_section_specific_layouts():
    out = tempfile.mkdtemp()
    w = LopWorkflow(agents=NoOpAgents(), audit_path=None, out_dir=out)
    st = w.new_run(project_name="Proj", client_name="Cli", voice_transcript="")
    st = w.resolve_intake(st, lop_category=LOPCategory.NON_COMPETITIVE)
    for _ in range(120):
        st = w.advance(st)
        if st.pending_checkpoint == HumanCheckpoint.POST_TOC:
            st = w.resolve_post_toc(st)
        elif st.pending_checkpoint == HumanCheckpoint.PRE_FINAL:
            st = w.resolve_pre_final(st)
        if st.phase == Phase.DONE:
            break
    assert st.lop and len(st.lop.sections) >= 11
    path = write_pptx(st, out)
    prs2 = Presentation(path)
    # Title + one slide per section (11+)
    assert len(prs2.slides) >= 12
