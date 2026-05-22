"""Build native PPTX charts from extracted chart descriptors.

Each chart descriptor contains a ``chartData`` dict (the JSON definition
from the HTML) and a ``rect`` with the bounding box in viewport pixels.

Extensible via CHART_TYPE_MAP -- add new chart types by mapping their
JSON type string to the corresponding ``XL_CHART_TYPE`` enum.  Chart
types that require a different data structure (e.g. XY/Bubble) can be
dispatched to a dedicated builder via CHART_BUILDER_MAP.

Legends are authored as explicit HTML elements (colored divs + text
labels) and auto-detected by the compiler as regular shapes.  The native
chart legend and chart title are always suppressed.
"""

from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_LABEL_POSITION,
    XL_TICK_LABEL_POSITION,
    XL_TICK_MARK,
)
from pptx.util import Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))

from shared.emu import px_to_emu_x, px_to_emu_y, CSS_PX_TO_PT

# ---------------------------------------------------------------------------
# Chart type registry
# ---------------------------------------------------------------------------

CHART_TYPE_MAP: dict[str, int] = {
    "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
    "bar_stacked_100": XL_CHART_TYPE.BAR_STACKED_100,
    "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
    "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "column_stacked_100": XL_CHART_TYPE.COLUMN_STACKED_100,
    "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
}

_STACKED_TYPES = {
    "bar_stacked", "bar_stacked_100",
    "column_stacked", "column_stacked_100",
}


# ---------------------------------------------------------------------------
# Color parser
# ---------------------------------------------------------------------------

_RGB_RE = re.compile(
    r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)(?:\s*,\s*[\d.]+)?\s*\)"
)
_HEX_RE = re.compile(r"^#?([0-9a-fA-F]{3,8})$")

_BLACK = RGBColor(0x00, 0x00, 0x00)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _srgb_hex(rgb: RGBColor) -> str:
    """OOXML a:srgbClr @val (six hex digits, no #)."""
    return "".join(f"{c:02X}" for c in tuple(rgb))


def _parse_color(css_color: str) -> RGBColor | None:
    if not css_color:
        return None
    css_color = css_color.strip()
    if css_color in ("transparent", "rgba(0, 0, 0, 0)"):
        return None
    m = _RGB_RE.match(css_color)
    if m:
        return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    m = _HEX_RE.match(css_color)
    if m:
        h = m.group(1)
        if len(h) == 3:
            h = h[0] * 2 + h[1] * 2 + h[2] * 2
        elif len(h) in (4, 8):
            h = h[:6]
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    return None


# ---------------------------------------------------------------------------
# Total labels via hidden helper series + c15 CELLRANGE
# ---------------------------------------------------------------------------

# Collected during chart building; consumed by patch_total_labels_xlsx()
# after prs.save().  Each entry: (chart_part_relpath, num_series, totals)
_pending_total_labels: list[tuple[str, int, list[float]]] = []

NS_C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_C15 = "http://schemas.microsoft.com/office/drawing/2012/chart"
NS_C16 = "http://schemas.microsoft.com/office/drawing/2014/chart"


def _set_plot_area_full_frame(chart) -> None:
    """Set a manual layout on the plot area so it fills the entire chart frame.

    Removes PowerPoint's automatic padding which otherwise misaligns bars
    with externally-positioned HTML labels.
    """
    from lxml import etree

    plot_area = chart._chartSpace.chart.plotArea
    layout = plot_area.find(f"{{{NS_C}}}layout")
    if layout is None:
        layout = etree.SubElement(plot_area, f"{{{NS_C}}}layout")
        plot_area.insert(0, layout)

    ml = layout.find(f"{{{NS_C}}}manualLayout")
    if ml is not None:
        layout.remove(ml)
    ml = etree.SubElement(layout, f"{{{NS_C}}}manualLayout")

    for tag, val in [
        ("layoutTarget", "inner"),
        ("xMode", "edge"),
        ("yMode", "edge"),
        ("x", "0"),
        ("y", "0"),
        ("w", "1"),
        ("h", "1"),
    ]:
        el = etree.SubElement(ml, f"{{{NS_C}}}{tag}")
        el.set("val", val)


def _add_total_labels_series(
    cd: CategoryChartData,
    series_list: list[dict],
    categories: list[str],
) -> list[float]:
    """Append a zero-valued helper series and return computed totals."""
    totals = []
    for ci in range(len(categories)):
        total = sum(s.get("values", [0] * (ci + 1))[ci] for s in series_list)
        totals.append(total)
    cd.add_series("Total bar", tuple([0] * len(categories)))
    return totals


def _insert_before_cat_or_val(ser_el, new_el):
    """Insert *new_el* into *ser_el* before <c:cat> or <c:val>."""
    from lxml import etree
    cat_el = ser_el.find(f"{{{NS_C}}}cat")
    if cat_el is not None:
        ser_el.insert(list(ser_el).index(cat_el), new_el)
        return
    val_el = ser_el.find(f"{{{NS_C}}}val")
    if val_el is not None:
        ser_el.insert(list(ser_el).index(val_el), new_el)
        return
    ser_el.append(new_el)


def _format_total_labels_series(
    chart, series_idx: int, tl_cfg: dict, totals: list[float],
    num_real_series: int,
) -> None:
    """Build CELLRANGE-based total labels on the zero-valued helper series.

    Mirrors the PowerPoint pattern: per-point dLbl with a:fld
    type="CELLRANGE" referencing a data-labels-range stored in the
    embedded Excel column F.
    """
    import uuid
    from lxml import etree

    helper = chart.series[series_idx]
    helper.format.fill.background()
    helper.format.line.fill.background()

    ser_elements = chart.plots[0]._element.findall(f"{{{NS_C}}}ser")
    if series_idx >= len(ser_elements):
        return
    ser_el = ser_elements[series_idx]

    cat_count = len(totals)
    font_size_px = tl_cfg.get("fontSize", 28)
    font_size_pt = float(font_size_px) * CSS_PX_TO_PT
    font_color = _parse_color(tl_cfg.get("fontColor", "#000000")) or _BLACK
    sz_hundredths_pt = str(int(font_size_pt * 100))
    srgb = _srgb_hex(font_color)

    # -- Per-point dLbl elements (must precede dLbls in OOXML order) --------
    dlbl_elements = []
    for ci in range(cat_count):
        dLbl = etree.Element(f"{{{NS_C}}}dLbl")

        idx_el = etree.SubElement(dLbl, f"{{{NS_C}}}idx")
        idx_el.set("val", str(ci))

        tx = etree.SubElement(dLbl, f"{{{NS_C}}}tx")
        rich = etree.SubElement(tx, f"{{{NS_C}}}rich")
        etree.SubElement(rich, f"{{{NS_A}}}bodyPr")
        etree.SubElement(rich, f"{{{NS_A}}}lstStyle")
        p = etree.SubElement(rich, f"{{{NS_A}}}p")
        fld = etree.SubElement(p, f"{{{NS_A}}}fld")
        fld.set("id", "{" + str(uuid.uuid4()).upper() + "}")
        fld.set("type", "CELLRANGE")
        rPr = etree.SubElement(fld, f"{{{NS_A}}}rPr")
        rPr.set("lang", "en-US")
        rPr.set("sz", sz_hundredths_pt)
        rPr.set("b", "1")
        solid_fill = etree.SubElement(rPr, f"{{{NS_A}}}solidFill")
        srgb_el = etree.SubElement(solid_fill, f"{{{NS_A}}}srgbClr")
        srgb_el.set("val", srgb)
        latin = etree.SubElement(rPr, f"{{{NS_A}}}latin")
        latin.set("typeface", "Arial")
        etree.SubElement(fld, f"{{{NS_A}}}pPr")
        t = etree.SubElement(fld, f"{{{NS_A}}}t")
        t.text = "[CELLRANGE]"
        end_rpr = etree.SubElement(p, f"{{{NS_A}}}endParaRPr")
        end_rpr.set("lang", "en-US")
        end_rpr.set("sz", sz_hundredths_pt)

        dLblPos = etree.SubElement(dLbl, f"{{{NS_C}}}dLblPos")
        dLblPos.set("val", "inBase")

        for tag, val in [
            ("showLegendKey", "0"), ("showVal", "0"), ("showCatName", "0"),
            ("showSerName", "0"), ("showPercent", "0"), ("showBubbleSize", "1"),
        ]:
            el = etree.SubElement(dLbl, f"{{{NS_C}}}{tag}")
            el.set("val", val)

        extLst = etree.SubElement(dLbl, f"{{{NS_C}}}extLst")
        ext = etree.SubElement(extLst, f"{{{NS_C}}}ext")
        ext.set("uri", "{CE6537A1-D6FC-4f65-9D91-7224C49458BB}")
        ext.set(f"{{{NS_C15}}}__dummy__", "")
        del ext.attrib[f"{{{NS_C15}}}__dummy__"]
        etree.SubElement(ext, f"{{{NS_C15}}}dlblFieldTable")
        xfs = etree.SubElement(ext, f"{{{NS_C15}}}xForSave")
        xfs.set("val", "1")
        sdlr = etree.SubElement(ext, f"{{{NS_C15}}}showDataLabelsRange")
        sdlr.set("val", "1")

        dlbl_elements.append(dLbl)

    # -- Build dLbls: per-point dLbl children + series-level defaults ---------
    # In OOXML, dLbl elements go INSIDE dLbls as the first children,
    # followed by series-level default formatting properties.
    dLbls = etree.Element(f"{{{NS_C}}}dLbls")

    for dlbl in dlbl_elements:
        dLbls.append(dlbl)

    numFmt = etree.SubElement(dLbls, f"{{{NS_C}}}numFmt")
    numFmt.set("formatCode", tl_cfg.get("numberFormat", "0"))
    numFmt.set("sourceLinked", "0")

    spPr = etree.SubElement(dLbls, f"{{{NS_C}}}spPr")
    etree.SubElement(spPr, f"{{{NS_A}}}noFill")
    ln = etree.SubElement(spPr, f"{{{NS_A}}}ln")
    etree.SubElement(ln, f"{{{NS_A}}}noFill")
    etree.SubElement(spPr, f"{{{NS_A}}}effectLst")

    txPr = etree.SubElement(dLbls, f"{{{NS_C}}}txPr")
    etree.SubElement(txPr, f"{{{NS_A}}}bodyPr")
    etree.SubElement(txPr, f"{{{NS_A}}}lstStyle")
    p = etree.SubElement(txPr, f"{{{NS_A}}}p")
    pPr = etree.SubElement(p, f"{{{NS_A}}}pPr")
    defRPr = etree.SubElement(pPr, f"{{{NS_A}}}defRPr")
    defRPr.set("sz", sz_hundredths_pt)
    defRPr.set("b", "1")
    solidFill = etree.SubElement(defRPr, f"{{{NS_A}}}solidFill")
    srgbClr = etree.SubElement(solidFill, f"{{{NS_A}}}srgbClr")
    srgbClr.set("val", srgb)
    latin = etree.SubElement(defRPr, f"{{{NS_A}}}latin")
    latin.set("typeface", "Arial")
    etree.SubElement(p, f"{{{NS_A}}}endParaRPr").set("lang", "LID4096")

    dLblPos_def = etree.SubElement(dLbls, f"{{{NS_C}}}dLblPos")
    dLblPos_def.set("val", "inBase")

    for tag, val in [
        ("showLegendKey", "0"), ("showVal", "0"), ("showCatName", "0"),
        ("showSerName", "0"), ("showPercent", "0"), ("showBubbleSize", "1"),
        ("showLeaderLines", "0"),
    ]:
        el = etree.SubElement(dLbls, f"{{{NS_C}}}{tag}")
        el.set("val", val)

    dLbls_extLst = etree.SubElement(dLbls, f"{{{NS_C}}}extLst")
    dLbls_ext = etree.SubElement(dLbls_extLst, f"{{{NS_C}}}ext")
    dLbls_ext.set("uri", "{CE6537A1-D6FC-4f65-9D91-7224C49458BB}")
    sdlr2 = etree.SubElement(dLbls_ext, f"{{{NS_C15}}}showDataLabelsRange")
    sdlr2.set("val", "1")
    sll = etree.SubElement(dLbls_ext, f"{{{NS_C15}}}showLeaderLines")
    sll.set("val", "0")

    # Insert dLbls before cat/val in the ser element
    _insert_before_cat_or_val(ser_el, dLbls)

    # -- c15:datalabelsRange in the series extLst ---------------------------
    last_row = cat_count + 1
    cell_ref = f"Sheet1!$F$2:$F${last_row}"

    ext_list_el = ser_el.find(f"{{{NS_C}}}extLst")
    if ext_list_el is None:
        ext_list_el = etree.SubElement(ser_el, f"{{{NS_C}}}extLst")

    dlr_ext = etree.SubElement(ext_list_el, f"{{{NS_C}}}ext")
    dlr_ext.set("uri", "{02D57815-91ED-43cb-92C2-25804820EDAC}")
    dlr = etree.SubElement(dlr_ext, f"{{{NS_C15}}}datalabelsRange")
    f_el = etree.SubElement(dlr, f"{{{NS_C15}}}f")
    f_el.text = cell_ref
    cache = etree.SubElement(dlr, f"{{{NS_C15}}}dlblRangeCache")
    ptCount = etree.SubElement(cache, f"{{{NS_C}}}ptCount")
    ptCount.set("val", str(cat_count))
    for ci, total in enumerate(totals):
        pt = etree.SubElement(cache, f"{{{NS_C}}}pt")
        pt.set("idx", str(ci))
        v = etree.SubElement(pt, f"{{{NS_C}}}v")
        v.text = str(int(total)) if total == int(total) else str(total)

    # Record for xlsx patching after save
    _pending_total_labels.append((num_real_series, totals))


# ---------------------------------------------------------------------------
# Category chart builder (bar, column, line, pie)
# ---------------------------------------------------------------------------

def _build_category_chart(slide, rect: dict, chart_data_json: dict) -> None:
    """Create a category-based PPTX chart from the extracted JSON."""

    chart_type_str = chart_data_json.get("type", "bar_stacked")
    xl_chart_type = CHART_TYPE_MAP.get(chart_type_str)
    if xl_chart_type is None:
        print(f"  Warning: unknown chart type '{chart_type_str}'", file=sys.stderr)
        return

    categories = chart_data_json.get("categories", [])
    series_list = chart_data_json.get("series", [])
    if not categories or not series_list:
        return

    is_stacked = chart_type_str in _STACKED_TYPES
    multi_series = len(series_list) > 1

    tl_cfg = chart_data_json.get("totalLabels", {})
    show_totals = is_stacked and multi_series and tl_cfg.get("show", True)

    cd = CategoryChartData()
    cd.categories = categories
    for s in series_list:
        cd.add_series(s.get("name", ""), tuple(s.get("values", [])))

    totals = []
    if show_totals:
        totals = _add_total_labels_series(cd, series_list, categories)

    left = px_to_emu_x(rect["x"])
    top = px_to_emu_y(rect["y"])
    width = px_to_emu_x(rect["width"])
    height = px_to_emu_y(rect["height"])

    graphic_frame = slide.shapes.add_chart(
        xl_chart_type, left, top, width, height, cd
    )
    chart = graphic_frame.chart

    # --- Series colors ---
    for i, s_json in enumerate(series_list):
        color = _parse_color(s_json.get("color", ""))
        if color and i < len(chart.series):
            fill = chart.series[i].format.fill
            fill.solid()
            fill.fore_color.rgb = color

    # --- Total labels helper series formatting ---
    if show_totals:
        _format_total_labels_series(
            chart, len(series_list), tl_cfg, totals, len(series_list),
        )

    # --- Plot-level formatting ---
    plot = chart.plots[0]

    gap_width = chart_data_json.get("gapWidth")
    if gap_width is not None:
        try:
            plot.gap_width = int(gap_width)
        except (AttributeError, TypeError):
            pass

    overlap = chart_data_json.get("overlap")
    if overlap is not None:
        try:
            plot.overlap = int(overlap)
        except (AttributeError, TypeError):
            pass

    # --- Data labels (on real series) ---
    dl_cfg = chart_data_json.get("dataLabels", {})
    if dl_cfg.get("show"):
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.position = XL_LABEL_POSITION.CENTER

        font_size = dl_cfg.get("fontSize")
        if font_size:
            data_labels.font.size = Pt(font_size * CSS_PX_TO_PT)

        font_color = _parse_color(dl_cfg.get("fontColor", ""))
        if font_color:
            data_labels.font.color.rgb = font_color

        number_format = dl_cfg.get("numberFormat")
        if number_format:
            data_labels.number_format = number_format
            data_labels.number_format_is_linked = False

    # --- Title / legend: always off (authored as HTML on the slide) ---
    chart.has_title = False
    chart.has_legend = False

    # --- Value axis ---
    val_cfg = chart_data_json.get("valueAxis", {})
    try:
        value_axis = chart.value_axis
    except Exception:
        value_axis = None

    if value_axis is not None:
        if val_cfg.get("visible") is not True:
            value_axis.visible = False
            value_axis.has_major_gridlines = False
        else:
            value_axis.visible = True
            if val_cfg.get("gridlines") is not True:
                value_axis.has_major_gridlines = False
            else:
                value_axis.has_major_gridlines = True

            axis_max = val_cfg.get("max")
            if axis_max is not None:
                value_axis.maximum_scale = float(axis_max)

            axis_min = val_cfg.get("min")
            if axis_min is not None:
                value_axis.minimum_scale = float(axis_min)

            number_format = val_cfg.get("numberFormat")
            if number_format:
                value_axis.tick_labels.number_format = number_format
                value_axis.tick_labels.number_format_is_linked = False

        value_axis.major_tick_mark = XL_TICK_MARK.NONE
        value_axis.minor_tick_mark = XL_TICK_MARK.NONE

    # --- Category axis (always visible; labels are authored in HTML) ---
    try:
        category_axis = chart.category_axis
    except Exception:
        category_axis = None

    if category_axis is not None:
        category_axis.visible = True

        category_axis.major_tick_mark = XL_TICK_MARK.NONE
        category_axis.minor_tick_mark = XL_TICK_MARK.NONE
        category_axis.tick_label_position = XL_TICK_LABEL_POSITION.NONE

        # Horizontal bar charts: first category at top to match HTML reading
        # order and external label columns (PowerPoint default is bottom-first).
        if chart_type_str.startswith("bar_"):
            category_axis.reverse_order = True

    _set_plot_area_full_frame(chart)


# ---------------------------------------------------------------------------
# Builder dispatch (extensible for XY, Bubble, etc.)
# ---------------------------------------------------------------------------

CHART_BUILDER_MAP: dict[str, callable] = {}


def add_chart(slide, shape_desc: dict) -> None:
    """Main entry point -- called by html_to_pptx for tag == 'chart'."""
    chart_data_json = shape_desc.get("chartData")
    if not chart_data_json:
        return

    rect = shape_desc.get("rect")
    if not rect:
        return

    chart_type = chart_data_json.get("type", "")

    builder = CHART_BUILDER_MAP.get(chart_type)
    if builder:
        builder(slide, rect, chart_data_json)
    elif chart_type in CHART_TYPE_MAP:
        _build_category_chart(slide, rect, chart_data_json)
    else:
        print(f"  Warning: no chart builder for type '{chart_type}'",
              file=sys.stderr)


# ---------------------------------------------------------------------------
# Post-save: patch embedded Excel with total-label column
# ---------------------------------------------------------------------------

def patch_total_labels_xlsx(pptx_path: str | Path) -> None:
    """Add column F (totals) to every chart xlsx that has CELLRANGE labels.

    Must be called after ``prs.save()`` because python-pptx builds the
    embedded xlsx at save time.  We reopen the pptx zip, locate each
    Excel workbook, and inject the extra column.
    """
    import io
    import zipfile
    import openpyxl

    if not _pending_total_labels:
        return

    pptx_path = Path(pptx_path)
    tmp_path = pptx_path.with_suffix(".pptx.tmp")

    with zipfile.ZipFile(pptx_path, "r") as zin, \
         zipfile.ZipFile(tmp_path, "w") as zout:

        xlsx_names: list[str] = []
        for item in zin.infolist():
            if (item.filename.startswith("ppt/embeddings/")
                    and item.filename.endswith(".xlsx")):
                xlsx_names.append(item.filename)

        chart_idx = 0
        patched_xlsx = {}

        for xlsx_name in xlsx_names:
            if chart_idx >= len(_pending_total_labels):
                break
            num_real_series, totals = _pending_total_labels[chart_idx]
            chart_idx += 1

            xlsx_bytes = zin.read(xlsx_name)
            wb = openpyxl.load_workbook(io.BytesIO(xlsx_bytes))
            ws = wb.active

            label_col = num_real_series + 2  # 1-indexed: A=cat, B..=series, next=helper, next=labels
            for ci, total in enumerate(totals):
                ws.cell(row=ci + 2, column=label_col, value=total)

            buf = io.BytesIO()
            wb.save(buf)
            patched_xlsx[xlsx_name] = buf.getvalue()

        for item in zin.infolist():
            if item.filename in patched_xlsx:
                zout.writestr(item, patched_xlsx[item.filename])
            else:
                zout.writestr(item, zin.read(item.filename))

    tmp_path.replace(pptx_path)
    _pending_total_labels.clear()
