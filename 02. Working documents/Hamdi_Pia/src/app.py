"""
LoP Builder — Phase 1 Streamlit Application

Run from the src/ directory:
    pip install -r requirements.txt
    cp .env.example .env          # then add your OPENAI_API_KEY
    streamlit run app.py
"""

import json
import os
import sys
import traceback
from datetime import datetime
from io import BytesIO
from pathlib import Path

import streamlit as st
import streamlit.components.v1 as st_components
from docx import Document as DocxDocument
from dotenv import load_dotenv

# ─── PATH + ENV ───────────────────────────────────────────────────────────────
_HAMDI_PIA_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(Path(__file__).parent))
sys.path.insert(0, str(_HAMDI_PIA_ROOT))  # lop_eval package
load_dotenv(Path(__file__).parent / ".env")

from datetime import date

from lop_chapters import render_chapter_brief
from mock_answers import generate_mock_answers
from orchestrator import load_agent_spec, run_agent
from slide_style_guide import render_slide_style_guide
from run_logger import (
    get_run_logger,
    log_agent_result,
    log_agent_start,
    log_event,
    log_evals_bundle,
    log_gate_event,
    log_session_start,
)
from schemas import (
    Answer,
    AnswerList,
    BASupportPack,
    ChapterSlideAuthorOutput,
    Citation,
    ClientEvaluationReport,
    EvalsBundle,
    ContextDoc,
    DotDashDoc,
    DotDashSlide,
    EmailDraft,
    IntakePackage,
    LossAnalysisReport,
    Question,
    QuestionList,
    SlideDeck,
    SlideHTML,
    SourceItem,
    SynthesisDoc,
    Todo,
    ValidationReport,
)
from deep_research import DeepResearchUnavailable, call_deep_research
from slide_renderer import (
    build_client_css_text,
    parse_client_pptx_style,
    write_single_deck,
)
from voice_answers import transcribe_and_map
from voice_transcription import transcribe_audio
from web_search import WebSearchUnavailable, call_with_web_search

# ─── PAGE CONFIG ─────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="LoP Builder — Phase 1",
    layout="wide",
    initial_sidebar_state="expanded",
)

# Apply the McKinsey-style theme. Must run after `set_page_config` and
# before any other Streamlit elements are rendered so the CSS applies
# globally on first paint. See src/_mckinsey_theme.py and
# src/.streamlit/config.toml.
from _mckinsey_theme import inject_theme_css

inject_theme_css()

# ─── PER-STEP MODEL SELECTION ────────────────────────────────────────────────
# Three model identifiers are exposed at every step's per-step selector:
#   _MODEL_FAST   — OpenAI smaller / faster / cheaper default.
#   _MODEL_FULL   — OpenAI higher-capability alternative.
#   _MODEL_CLAUDE — Anthropic premium-quality track (Opus). Slower than
#                   _MODEL_FULL but strongest on partner-grade prose,
#                   declarative action titles, and faithful long-document
#                   compression. Routed to the Anthropic gateway by
#                   `orchestrator.run_agent` based on the `claude-` prefix.
# The UI pre-seeds each step with the recommended choice (see
# `_MODEL_RECOMMENDED`) and shows a step-specific tradeoff caption per
# track so the user understands the effect before changing it.
_MODEL_FAST:   str = os.environ.get("OPENAI_MODEL", "gpt-5.4-mini-2026-03-17")
_MODEL_FULL:   str = os.environ.get("OPENAI_MODEL_FULL", "gpt-5.5-2026-04-23")
_MODEL_CLAUDE: str = os.environ.get("ANTHROPIC_MODEL", "claude-opus-4-7")

# Recommended pre-seed per step. "fast" | "full" | "claude".
# Claude is reserved for steps where the prose / faithfulness lift justifies
# the latency cost (synthesis, dot-dash, BA support, context structuring,
# slides). For slides specifically: the deck dispatches into 10 chapter-
# specialised authoring agents and Claude is the only track that
# consistently produces partner-grade McKinsey prose without AI-language
# tells; the full deck renders in ~12-20 min on Claude vs ~3 min on Faster.
# Faster / Better remain available as overrides for re-rendering ONE
# drifting slide.
_MODEL_RECOMMENDED: dict[str, str] = {
    "intake":            "fast",
    "context":           "claude",
    "synthesis":         "claude",
    "validation":        "fast",
    "dotdash":           "claude",
    "ba_support":        "claude",
    "slides":            "claude",
    "voice_structurer":  "claude",
    "client_eval":       "claude",
    "loss_eval":         "claude",
}

# Per-step tradeoff explanations (caption_fast, caption_full, caption_claude).
# Visible inside the model selector expander above each Run button.
_MODEL_TRADEOFFS: dict[str, tuple[str, str, str]] = {
    "intake": (
        "Pure extraction — RFP requirements, chapter buckets, gap "
        "list. Mini handles this reliably and is roughly 3x faster. "
        "Typical: 25–40 s.",
        "Marginally better at picking up subtle pursuit signals (RFI "
        "vs. RFP nuance, partner / competitor name disambiguation). "
        "Rarely worth the latency. Typical: 70–110 s.",
        "Overkill for pure extraction — Opus's prose strengths add no "
        "value here. Use only when an unusually large or messy multi-PDF "
        "RFP needs careful long-context reading. Typical: 90–150 s.",
    ),
    "context": (
        "Quick mode: fine — the chat model just structures `web_search` "
        "results into JSON. Deep mode: mini sometimes compresses the "
        "research report, dropping named individuals, dated events, or "
        "per-dimension specifics. Typical: 15–30 s.",
        "Recommended for Deep mode. Preserves more specifics from the "
        "4–10k-word research report and respects the FAITHFULNESS "
        "rules better. Typical: 40–70 s.",
        "Recommended for Deep mode when the research report is long and "
        "names / dates / per-dimension specifics matter. Strongest "
        "long-context faithfulness — drops fewest details from the "
        "research report. Typical: 60–110 s.",
    ),
    "synthesis": (
        "Usable, but problem statements and chapter-construction "
        "questions tend to be safer and less differentiated. "
        "Typical: 25–40 s.",
        "Sharper problem statements and more on-spec chapter-"
        "construction questions. The synthesis brief is the first "
        "artefact a partner reads — worth the latency. "
        "Typical: 60–100 s.",
        "Recommended. Strongest concise, rhetorical, partner-grade "
        "prose for the brief and problem statement. The first artefact "
        "a partner reads — biggest visible quality win once per "
        "pursuit. Typical: 90–150 s.",
    ),
    "validation": (
        "Rubric-based work — 'did this answer fill the gap?'. Mini is "
        "already very good at this. Typical: 25–40 s.",
        "Slightly better follow-up question wording. Rarely changes "
        "the verdict. Typical: 60–100 s.",
        "Overkill — rubric work does not benefit from premium prose. "
        "Use only if follow-up wording feels off in repeated runs. "
        "Typical: 80–140 s.",
    ),
    "dotdash": (
        "Usable, but action titles tend to be safer and less "
        "differentiated. Typical: 30–50 s.",
        "Sharper, more partner-grade action titles and "
        "confidence-honest dashes. Typical: 90–150 s.",
        "Recommended — the flagship slot. Sharpest declarative "
        "action titles, best 'so what' framing, most disciplined "
        "confidence-honest dashes. The McKinsey signature artefact. "
        "Typical: 120–200 s.",
    ),
    "ba_support": (
        "Todos and source items are well structured; email drafts "
        "read more boilerplate. Fine if the BA edits emails before "
        "sending. Typical: 25–40 s.",
        "Noticeably more natural, partner-grade prose. Worth it if "
        "you send the emails as-is. Typical: 60–100 s.",
        "Recommended if you send the emails as-is. Strongest "
        "natural, professional prose — fewest hedges, cleanest "
        "rhythm. Typical: 90–150 s.",
    ),
    "slides": (
        "Useful only for re-rendering ONE drifting slide. Mini follows "
        "the authoring contract but its prose reads as templated AI "
        "language and rarely clears partner review. Typical (full "
        "deck): 2–3 minutes.",
        "Useful only for re-rendering ONE drifting slide. Better layout "
        "decisions than Faster, but still not partner-grade prose at the "
        "level the chapter specialists need. Typical (full deck): 6–9 "
        "minutes.",
        "Recommended. Partner-grade prose for all 10 chapter "
        "specialists — declarative action titles, no AI hedging, "
        "real grounding in intake / context / BA pack. The full deck "
        "renders in ~12–20 minutes. Switch to Faster only when "
        "re-running a single slide that drifted.",
    ),
    "voice_structurer": (
        "Usable, but the structured perspective and per-chapter "
        "signals tend to be flatter and occasionally drop nuance the "
        "partner stated in passing. Typical: 10–20 s.",
        "Sharper structure and better chapter-mapping than Faster, "
        "with fewer dropped nuances. Reasonable middle ground when "
        "the Claude track is unavailable. Typical: 20–40 s.",
        "Recommended. Strongest faithful summarisation — preserves "
        "the partner's hedges, names, and numbers, and maps chapter "
        "signals most reliably. Runs once per memo, so the latency "
        "cost is a one-off. Typical: 40–90 s.",
    ),
    "client_eval": (
        "Usable for a rough buyer-perspective verdict, but tends to "
        "soften critiques and miss the nuanced 'this reads as "
        "boilerplate' calls that make this evaluator useful. "
        "Typical: 30–50 s.",
        "Sharper critiques than Faster, more willing to call out "
        "weak chapters specifically. Reasonable middle ground when "
        "Claude is rate-limited. Typical: 70–120 s.",
        "Recommended. Strongest at owner-voice critique — direct, "
        "specific, willing to flag generic prose as generic. The "
        "buyer-perspective slot benefits most from premium prose. "
        "Typical: 90–160 s.",
    ),
    "loss_eval": (
        "Usable for a directional loss-risk verdict, but tends to "
        "hedge severity and produce generic competitor angles. "
        "Misses the specific 'where exactly will the rival land the "
        "punch' calls that make red-teaming useful. Typical: 30–50 s.",
        "Sharper severity calls and more specific competitor angles "
        "than Faster. Reasonable when Claude is rate-limited, though "
        "model knowledge of named rival firms is still thinner. "
        "Typical: 70–120 s.",
        "Recommended. Strongest at partner-grade red-team prose and "
        "richest model knowledge of named competitor firms — willing "
        "to call critical risks critical, specific about which "
        "chapter the rival will exploit. Typical: 90–160 s.",
    ),
}


def _model_state_key(step_id: str) -> str:
    return f"model_{step_id}"


def _track_for_model(model_id: str) -> str:
    """Return the short track label ('Faster' | 'Better' | 'Claude') for a model id."""
    if model_id == _MODEL_CLAUDE:
        return "Claude"
    if model_id == _MODEL_FULL:
        return "Better"
    return "Faster"


def _format_model_option(model_id: str) -> str:
    """Format the radio-option label for a given model id."""
    return f"{_track_for_model(model_id)} — {model_id}"


def _model_for_track(track: str) -> str:
    """Resolve a track label ('fast' | 'full' | 'claude') to its model id."""
    if track == "claude":
        return _MODEL_CLAUDE
    if track == "full":
        return _MODEL_FULL
    return _MODEL_FAST


def _render_model_selector(step_id: str) -> str:
    """
    Render an expandable model selector for one step. Pre-seeds the
    recommended choice the first time it's shown. Returns the currently
    selected model identifier so the call site can pass it as
    `run_agent(..., model=...)` or `call_with_web_search(..., model=...)`.
    """
    state_key = _model_state_key(step_id)
    if state_key not in st.session_state:
        rec = _MODEL_RECOMMENDED.get(step_id, "fast")
        st.session_state[state_key] = _model_for_track(rec)

    current = st.session_state[state_key]
    # Migrate stale state from the older 2-option UI: if a previous session
    # stored a value that is no longer one of the three offered ids, snap it
    # back to the step's recommendation rather than crashing the radio.
    options = [_MODEL_FAST, _MODEL_FULL, _MODEL_CLAUDE]
    if current not in options:
        rec = _MODEL_RECOMMENDED.get(step_id, "fast")
        current = _model_for_track(rec)
        st.session_state[state_key] = current

    fast_cap, full_cap, claude_cap = _MODEL_TRADEOFFS.get(
        step_id,
        (
            "Faster, smaller model.",
            "Slower, more capable model.",
            "Premium prose track — slowest but strongest at partner-grade "
            "prose. Use sparingly.",
        ),
    )
    rec_key = _MODEL_RECOMMENDED.get(step_id, "fast")
    rec_label = {"fast": "Faster", "full": "Better", "claude": "Claude"}.get(
        rec_key, "Faster"
    )

    with st.expander(
        f"Model for this step: **{current}** "
        f"({_track_for_model(current)}) — click to change",
        expanded=False,
    ):
        st.caption(
            f"Pick the LLM for this step. Recommended for this step: "
            f"**{rec_label}**. Defaults are pre-set; re-run the step "
            f"after switching."
        )
        choice = st.radio(
            "Model",
            options=options,
            index=options.index(current),
            format_func=_format_model_option,
            key=f"radio_{state_key}",
            label_visibility="collapsed",
        )
        st.session_state[state_key] = choice
        st.markdown(
            f"- **Faster — `{_MODEL_FAST}`**: {fast_cap}\n"
            f"- **Better — `{_MODEL_FULL}`**: {full_cap}\n"
            f"- **Claude — `{_MODEL_CLAUDE}`**: {claude_cap}"
        )
    return st.session_state[state_key]


# ─── SESSION STATE INIT ───────────────────────────────────────────────────────
_DEFAULTS: dict = {
    "processed_docs": [],
    "upload_done": False,
    # Step 0 — Partner Context voice memos awaiting Confirm Uploads
    "partner_voice_memos": [],            # [{filename, text}, ...]
    "partner_voice_pending_transcript": "",  # last transcribed memo, awaiting partner edit
    "partner_voice_pending_filename": "",
    # Structured-memo state — populated when the partner clicks "Structure
    # transcript" on the pending raw transcript. The Markdown text shown
    # in the second editable textarea lives in the widget key
    # `step0_structured_memo_edit`; `step0_structured_memo` holds the
    # initial value to seed it with after the agent runs.
    "step0_structured_memo": "",
    "step0_structured_memo_done": False,
    "intake_package": None,
    "intake_done": False,
    "context_doc": None,
    "context_done": False,
    "synthesis_doc": None,
    "synthesis_done": False,
    "gate_a_done": False,
    "answer_list": None,
    "answers_generated": False,
    "answers_done": False,
    "answers_source": None,      # "voice" | "mock" — which path produced answer_list
    "voice_transcript": "",      # raw whisper-1 output, kept for partner review
    "validation_report": None,
    "validation_done": False,
    "gate_b_done": False,
    "dotdash_doc": None,
    "dotdash_done": False,
    "gate_c_done": False,
    "ba_support_pack": None,
    "ba_support_done": False,
    "slide_format_mode": "mckinsey",
    "client_pptx_bytes": None,
    "client_style_summary": "",
    "slide_deck": None,
    "slides_done": False,
    # Step 9 — Consistency evaluation (optional buyer/loss reviews in expander)
    "evals_bundle": None,
    "evals_done": False,
    "client_eval_report": None,
    "client_eval_done": False,
    "client_eval_source": "in_app_deck",   # "in_app_deck" | "uploaded_file"
    "client_eval_upload_name": "",
    "loss_eval_report": None,
    "loss_eval_done": False,
    "loss_eval_competitor_override": "",
    "consistency_eval_result": None,
    "consistency_eval_done": False,
    "consistency_eval_use_llm_judge": False,
    "consistency_eval_pass_threshold": 70.0,
    "consistency_checker_weights": None,  # dict[str, float] — None → defaults
    "_last_evals_qp_signature": None,
    "evals_scroll_pending": False,
}
for _k, _v in _DEFAULTS.items():
    if _k not in st.session_state:
        st.session_state[_k] = _v

ss = st.session_state

# Assign a unique run ID to this session (persists across Streamlit reruns)
if "run_id" not in ss:
    ss.run_id = datetime.now().strftime("%Y%m%d_%H%M%S")
    log_session_start(ss.run_id)

_, _log_path = get_run_logger(ss.run_id)


def _evals_qp_signature() -> str:
    if "evals" not in st.query_params:
        return "__no_evals_param__"
    raw = st.query_params.get("evals")
    if isinstance(raw, list):
        raw = raw[0] if raw else ""
    return f"evals={raw!r}"


def _evals_qp_means_scroll() -> bool:
    if "evals" not in st.query_params:
        return False
    raw = st.query_params.get("evals")
    if isinstance(raw, list):
        raw = raw[0] if raw else ""
    s = str(raw).lower().strip()
    return s in ("", "1", "true", "yes")


_sig = _evals_qp_signature()
if _sig != ss.get("_last_evals_qp_signature"):
    ss._last_evals_qp_signature = _sig
    if _evals_qp_means_scroll():
        ss.evals_scroll_pending = True


# ─── STEP DEFINITIONS ────────────────────────────────────────────────────────
_STEPS = [
    ("upload_done",     "Upload Documents"),
    ("intake_done",     "Intake Agent"),
    ("context_done",    "Context Agent"),
    ("synthesis_done",  "Synthesis Agent"),
    ("gate_a_done",     "Gate A — Approve Questions"),
    ("answers_done",    "Partner Answers"),
    ("validation_done", "Validation Agent"),
    ("gate_b_done",     "Gate B — Approve to Dot-Dash"),
    ("dotdash_done",    "Dot-Dash Agent"),
    ("gate_c_done",     "Gate C — Approve Dot-Dash"),
    ("ba_support_done", "BA Support Pack"),
    ("slides_done",     "Render Slides"),
    ("evals_done",      "Consistency eval"),
]

# ─── HELPERS ─────────────────────────────────────────────────────────────────

def _quality_badge(quality: str) -> None:
    if quality == "complete":
        st.success("complete")
    elif quality == "partial":
        st.warning("partial")
    else:
        st.error("missing")


def _completeness_badge(completeness: str) -> None:
    if completeness == "complete":
        st.success("complete")
    elif completeness == "partial":
        st.warning("partial")
    else:
        st.error("missing")


def _refresh_evals_done() -> None:
    ss.evals_done = bool(ss.consistency_eval_done)


def _sync_evals_bundle() -> None:
    consistency_dump = None
    if ss.consistency_eval_result is not None:
        consistency_dump = (
            ss.consistency_eval_result.model_dump()
            if hasattr(ss.consistency_eval_result, "model_dump")
            else ss.consistency_eval_result
        )
    ss.evals_bundle = EvalsBundle(
        run_id=ss.run_id,
        proposal_source=ss.client_eval_source,
        proposal_upload_name=ss.client_eval_upload_name or "",
        client_eval=ss.client_eval_report,
        loss_eval=ss.loss_eval_report,
        consistency_eval=consistency_dump,
        updated_at=datetime.now().isoformat(timespec="seconds"),
    )
    _refresh_evals_done()
    try:
        log_evals_bundle(ss.run_id, ss.evals_bundle)
    except Exception:
        pass


def _default_consistency_weights() -> dict[str, float]:
    from lop_eval.scoring import DEFAULT_CHECKER_WEIGHTS

    return dict(DEFAULT_CHECKER_WEIGHTS)


def _build_eval_config_from_session() -> object:
    from lop_eval.models import EvalConfig

    weights = ss.consistency_checker_weights
    if weights is None:
        weights = _default_consistency_weights()
    return EvalConfig(
        pass_threshold=float(ss.consistency_eval_pass_threshold),
        checker_weights=weights,
    )


def _render_consistency_weights_ui() -> None:
    from lop_eval.scoring import CHECKER_LABELS, DEFAULT_CHECKER_WEIGHTS

    if ss.consistency_checker_weights is None:
        ss.consistency_checker_weights = dict(DEFAULT_CHECKER_WEIGHTS)

    with st.expander("Scoring weights (per check)", expanded=False):
        st.caption(
            "Multiply severity deductions by each weight (0 = ignore check, "
            "1 = default, 2 = double penalty)."
        )
        ss.consistency_eval_pass_threshold = st.slider(
            "Pass threshold (score)",
            min_value=50.0,
            max_value=100.0,
            value=float(ss.consistency_eval_pass_threshold),
            step=5.0,
            key="consistency_pass_threshold_slider",
        )
        cols = st.columns(2)
        keys = list(DEFAULT_CHECKER_WEIGHTS.keys())
        for i, key in enumerate(keys):
            label = CHECKER_LABELS.get(key, key)
            with cols[i % 2]:
                ss.consistency_checker_weights[key] = st.slider(
                    label,
                    min_value=0.0,
                    max_value=2.0,
                    value=float(ss.consistency_checker_weights.get(key, 1.0)),
                    step=0.1,
                    key=f"consistency_weight_{key}",
                )
        if st.button("Reset weights to defaults", key="reset_consistency_weights"):
            ss.consistency_checker_weights = _default_consistency_weights()
            st.rerun()


def _render_evals_sidebar_status() -> None:
    if not ss.slides_done:
        return
    mark = "✓" if ss.consistency_eval_done else "○"
    st.caption(f"  {mark} Consistency evaluation")


def _parse_firms_input(raw: str) -> list[str]:
    """
    Parse the partner's free-text competitor-firms input.

    Splits on newlines and commas, strips whitespace, drops empties, and
    deduplicates while preserving the partner's ordering.
    """
    if not raw:
        return []
    pieces: list[str] = []
    for line in raw.splitlines():
        for piece in line.split(","):
            cleaned = piece.strip()
            if cleaned:
                pieces.append(cleaned)
    seen: set[str] = set()
    deduped: list[str] = []
    for p in pieces:
        key = p.casefold()
        if key in seen:
            continue
        seen.add(key)
        deduped.append(p)
    return deduped


def _reset_from(step: str) -> None:
    """Clear state for `step` and all downstream steps.

    Note: consistency eval can re-run independently; resetting from any upstream
    step clears the Step 9 bundle (including optional buyer/loss reviews).
    """
    _order = [
        "intake", "context", "synthesis", "gate_a",
        "answers", "validation", "gate_b", "dotdash", "gate_c",
        "ba_support", "slides", "evals",
    ]
    if step not in _order:
        return
    idx = _order.index(step)

    # Flags
    for s in _order[idx:]:
        ss[f"{s}_done"] = False

    # Data blobs per step
    _data: dict[str, list] = {
        "intake":       ["intake_package"],
        "context":      ["context_doc"],
        "synthesis":    ["synthesis_doc"],
        "gate_a":       [],
        "answers":      ["answer_list"],
        "validation":   ["validation_report"],
        "gate_b":       [],
        "dotdash":      ["dotdash_doc"],
        "gate_c":       [],
        "ba_support":   ["ba_support_pack"],
        "slides":       ["slide_deck"],
        "evals":        [
            "evals_bundle",
            "client_eval_report",
            "loss_eval_report",
            "consistency_eval_result",
        ],
    }
    for s in _order[idx:]:
        for field in _data.get(s, []):
            ss[field] = None
    if "evals" in _order[idx:]:
        ss.client_eval_done = False
        ss.loss_eval_done = False
        ss.consistency_eval_done = False
        ss.evals_done = False

    # Clear Step 1 partner-input widgets when intake is reset (e.g. Change uploads).
    if "intake" in _order[idx:]:
        for key in ("intake_competitive_status", "intake_competitor_firms"):
            ss.pop(key, None)
    if "answers" in _order[idx:]:
        ss.answers_generated = False
        ss.answers_source = None
        ss.voice_transcript = ""
        for key in [
            k for k in ss.keys()
            if k.startswith("voice_recording_")
            or k.startswith("voice_upload_")
            or k.startswith("answer_edit_")
        ]:
            del ss[key]

    # Clear Gate A question edits when synthesis is reset
    if idx <= _order.index("synthesis"):
        for key in [k for k in ss.keys() if k.startswith("q_edit_")]:
            del ss[key]

    # Clear Gate C dot-dash edits when dot-dash is reset
    if idx <= _order.index("dotdash"):
        for key in [
            k for k in ss.keys()
            if k.startswith("dd_headline_")
            or k.startswith("dd_dashes_")
            or k.startswith("dd_notes_")
        ]:
            del ss[key]


def _format_structured_memo(structured: dict) -> str:
    """
    Render the `voice-structurer-agent` JSON output as a Markdown block
    the partner can read and edit in a `st.text_area`.

    Skips empty fields so the partner sees only the populated sections —
    the agent is allowed to leave fields empty when the transcript
    carries no signal, and we do not want to surface placeholders like
    "Win themes / differentiators: (none)" that would invite the partner
    to invent content.
    """
    if not isinstance(structured, dict):
        return ""

    lines: list[str] = []

    perspective = structured.get("lop_perspective") or {}
    if isinstance(perspective, dict):
        section_lines: list[str] = []

        client_situation = (perspective.get("client_situation") or "").strip()
        if client_situation:
            section_lines.append("**Client situation**")
            section_lines.append(client_situation)
            section_lines.append("")

        why_now = (perspective.get("why_now") or "").strip()
        if why_now:
            section_lines.append("**Why now**")
            section_lines.append(why_now)
            section_lines.append("")

        wants = (perspective.get("what_partner_wants_in_lop") or "").strip()
        if wants:
            section_lines.append("**What the partner wants in the LoP**")
            section_lines.append(wants)
            section_lines.append("")

        win_themes = perspective.get("win_themes") or []
        win_themes = [str(t).strip() for t in win_themes if str(t).strip()]
        if win_themes:
            section_lines.append("**Win themes / differentiators**")
            for t in win_themes:
                section_lines.append(f"- {t}")
            section_lines.append("")

        open_questions = perspective.get("open_questions") or []
        open_questions = [str(q).strip() for q in open_questions if str(q).strip()]
        if open_questions:
            section_lines.append("**Open questions**")
            for q in open_questions:
                section_lines.append(f"- {q}")
            section_lines.append("")

        if section_lines:
            lines.append("## LoP perspective")
            lines.append("")
            lines.extend(section_lines)

    chapter_signals = structured.get("chapter_signals") or []
    chapter_lines: list[str] = []
    for entry in chapter_signals:
        if not isinstance(entry, dict):
            continue
        chapter = (entry.get("chapter") or "").strip()
        signal = (entry.get("signal") or "").strip()
        if not chapter or not signal:
            continue
        chapter_lines.append(f"**{chapter}**")
        chapter_lines.append(signal)
        chapter_lines.append("")

    if chapter_lines:
        if lines and lines[-1] != "":
            lines.append("")
        lines.append("## Signal per LoP chapter")
        lines.append("")
        lines.extend(chapter_lines)

    return "\n".join(lines).rstrip() + ("\n" if lines else "")


def _reset_partner_voice_memos() -> None:
    """Clear all Step 0 partner voice-memo state (pending + committed)."""
    ss.partner_voice_memos = []
    ss.partner_voice_pending_transcript = ""
    ss.partner_voice_pending_filename = ""
    ss.step0_structured_memo = ""
    ss.step0_structured_memo_done = False
    for key in [
        k for k in list(ss.keys())
        if k.startswith("step0_voice_recording")
        or k.startswith("step0_voice_upload")
        or k.startswith("step0_voice_transcript_edit")
        or k.startswith("step0_structured_memo_edit")
    ]:
        del ss[key]


def _render_partner_voice_memo_block() -> None:
    """
    Step 0 — partner voice-memo recorder/uploader.

    Records (or accepts an uploaded) audio file, runs whisper-1 to
    transcribe, lets the partner clean the transcript, and on commit
    appends a `{filename, text}` entry to `ss.partner_voice_memos`. The
    confirmed memos are merged into `ss.processed_docs` as
    `doc_type="Partner Context"` when the user clicks **Confirm Uploads**.
    """
    with st.expander(
        "Or: record / upload a partner voice memo (transcribed via whisper-1)",
        expanded=False,
    ):
        st.caption(
            "Record in-browser or upload an audio file (mp3, m4a, wav, "
            "webm, mp4, mpga — up to 25 MB). The transcript is added as "
            "a **Partner Context** document the Intake Agent reads as "
            "authoritative partner input."
        )

        # Already-committed memos (rendered with Remove buttons).
        if ss.partner_voice_memos:
            st.markdown("**Voice memos staged for upload:**")
            for i, memo in enumerate(ss.partner_voice_memos):
                cols = st.columns([5, 2, 1])
                cols[0].write(memo["filename"])
                cols[1].caption(f"{len(memo['text'])} chars")
                if cols[2].button("Remove", key=f"step0_voice_remove_{i}"):
                    ss.partner_voice_memos.pop(i)
                    log_event(ss.run_id, f"[STEP0 VOICE] memo removed: {memo['filename']}")
                    st.rerun()
            st.divider()

        recorded = st.audio_input(
            "Record memo", key="step0_voice_recording"
        )
        uploaded_audio = st.file_uploader(
            "Or upload an audio file",
            type=["mp3", "m4a", "wav", "webm", "mp4", "mpga"],
            key="step0_voice_upload",
            accept_multiple_files=False,
        )

        picked = uploaded_audio if uploaded_audio is not None else recorded
        disabled = picked is None

        if picked is not None:
            size_kb = len(picked.getvalue()) / 1024
            src_label = (
                "uploaded file"
                if uploaded_audio is not None
                else "in-browser recording"
            )
            st.caption(
                f"Ready to transcribe — {src_label}: "
                f"{getattr(picked, 'name', 'memo')} ({size_kb:.1f} KB)"
            )

        if st.button(
            "Transcribe with whisper-1",
            type="secondary",
            key="step0_voice_transcribe",
            disabled=disabled,
        ):
            audio_bytes = picked.getvalue()
            raw_name = getattr(picked, "name", "") or ""
            log_event(
                ss.run_id,
                f"[STEP0 VOICE] audio received: "
                f"{raw_name or '(in-browser memo)'}, "
                f"{len(audio_bytes) / 1024:.1f} KB",
            )

            try:
                with st.spinner(
                    "Transcribing with whisper-1 (typically 10–40 seconds)..."
                ):
                    transcript = transcribe_audio(
                        audio_bytes, raw_name or "memo.wav"
                    )
                log_event(
                    ss.run_id,
                    f"[STEP0 VOICE] transcribed {len(transcript)} chars",
                )

                next_idx = len(ss.partner_voice_memos) + 1
                pending_filename = f"partner_context_voice_{next_idx}.txt"

                ss.partner_voice_pending_transcript = transcript
                ss.partner_voice_pending_filename = pending_filename
                if "step0_voice_transcript_edit" in ss:
                    del ss["step0_voice_transcript_edit"]
                st.rerun()
            except ValueError as exc:
                log_event(ss.run_id, f"[STEP0 VOICE] ERROR: {exc}")
                st.error(str(exc))
            except Exception as exc:
                traceback.print_exc()
                log_event(ss.run_id, f"[STEP0 VOICE] ERROR: {exc}")
                st.error(f"Transcription failed: {exc}")

        if ss.partner_voice_pending_transcript:
            st.markdown("**Transcript — edit before adding:**")
            st.caption(
                "Clean up obvious whisper-1 mishears (proper nouns, numbers, "
                "acronyms). You can also click **Structure transcript** "
                "below to turn the rambling memo into a structured LoP "
                "perspective the Intake Agent reads more cleanly."
            )
            st.text_area(
                "transcript",
                value=ss.partner_voice_pending_transcript,
                key="step0_voice_transcript_edit",
                height=200,
                label_visibility="collapsed",
            )

            # ── Structure transcript: turn rambling speech into a
            # structured LoP perspective + per-chapter signal map. The
            # partner can edit the structured output before adding the
            # memo; we save the structured version when present, fall
            # back to the raw transcript otherwise.
            voice_structurer_model = _render_model_selector("voice_structurer")
            if st.button(
                "Structure transcript",
                key="step0_voice_structure",
                help=(
                    "Run the voice-structurer-agent on the transcript above. "
                    "Output appears in a second editable textarea so you can "
                    "refine before adding the memo."
                ),
            ):
                source_text = ss.get(
                    "step0_voice_transcript_edit",
                    ss.partner_voice_pending_transcript,
                ).strip()
                if not source_text:
                    st.error("Transcript is empty — nothing to structure.")
                else:
                    try:
                        log_agent_start(
                            ss.run_id,
                            "voice-structurer-agent",
                            f"{ss.partner_voice_pending_filename} | "
                            f"{len(source_text)} chars | "
                            f"model={voice_structurer_model}",
                        )
                        with st.spinner(
                            "Structuring transcript into LoP perspective "
                            "(15–90 seconds)..."
                        ):
                            result = run_agent(
                                "voice-structurer-agent",
                                user_message=(
                                    "Structure the following partner voice "
                                    "transcript into the LoP perspective + "
                                    "chapter signal map per the Output "
                                    "Schema. Use ONLY what the partner "
                                    "actually said.\n\n"
                                    "Transcript:\n\"\"\"\n"
                                    f"{source_text}\n"
                                    "\"\"\""
                                ),
                                files=None,
                                use_extended_thinking=False,
                                model=voice_structurer_model,
                            )
                        log_agent_result(
                            ss.run_id, "voice-structurer-agent", result
                        )
                        formatted = _format_structured_memo(result)
                        if not formatted.strip():
                            st.warning(
                                "Structurer ran but returned no content — "
                                "transcript may be too short or unclear. "
                                "Edit the raw transcript and try again, or "
                                "add the raw version directly."
                            )
                        else:
                            ss.step0_structured_memo = formatted
                            ss.step0_structured_memo_done = True
                            if "step0_structured_memo_edit" in ss:
                                del ss["step0_structured_memo_edit"]
                            st.rerun()
                    except Exception as exc:
                        traceback.print_exc()
                        log_event(
                            ss.run_id,
                            f"[VOICE STRUCTURER AGENT] ERROR: {exc}",
                        )
                        st.error(f"Structuring failed: {exc}")

            if ss.step0_structured_memo_done:
                st.markdown("**Structured LoP perspective — edit before adding:**")
                st.caption(
                    "This is what gets added as Partner Context when you "
                    "click **Add as Partner Context** below. Edits here "
                    "override the structurer's output."
                )
                st.text_area(
                    "structured memo",
                    value=ss.step0_structured_memo,
                    key="step0_structured_memo_edit",
                    height=320,
                    label_visibility="collapsed",
                )

            commit_cols = st.columns([2, 2, 4])
            with commit_cols[0]:
                if st.button(
                    "Add as Partner Context",
                    type="primary",
                    key="step0_voice_commit",
                ):
                    use_structured = bool(ss.step0_structured_memo_done)
                    if use_structured:
                        edited = ss.get(
                            "step0_structured_memo_edit",
                            ss.step0_structured_memo,
                        ).strip()
                    else:
                        edited = ss.get(
                            "step0_voice_transcript_edit",
                            ss.partner_voice_pending_transcript,
                        ).strip()
                    if not edited:
                        st.error("Memo is empty — nothing to add.")
                    else:
                        # Tag structured-version memos in the filename so
                        # the partner sees at a glance which memos went
                        # through the structurer when reviewing the staged
                        # list above.
                        original_name = ss.partner_voice_pending_filename
                        if use_structured:
                            stem, _, ext = original_name.rpartition(".")
                            if stem and ext:
                                final_name = f"{stem}_structured.{ext}"
                            else:
                                final_name = f"{original_name} (structured)"
                        else:
                            final_name = original_name

                        ss.partner_voice_memos.append(
                            {
                                "filename": final_name,
                                "text": edited,
                            }
                        )
                        log_event(
                            ss.run_id,
                            f"[STEP0 VOICE] memo committed: "
                            f"{final_name} "
                            f"({len(edited)} chars, "
                            f"{'structured' if use_structured else 'raw'})",
                        )
                        ss.partner_voice_pending_transcript = ""
                        ss.partner_voice_pending_filename = ""
                        ss.step0_structured_memo = ""
                        ss.step0_structured_memo_done = False
                        if "step0_voice_transcript_edit" in ss:
                            del ss["step0_voice_transcript_edit"]
                        if "step0_structured_memo_edit" in ss:
                            del ss["step0_structured_memo_edit"]
                        st.rerun()
            with commit_cols[1]:
                if st.button(
                    "Discard transcript",
                    key="step0_voice_discard",
                ):
                    ss.partner_voice_pending_transcript = ""
                    ss.partner_voice_pending_filename = ""
                    ss.step0_structured_memo = ""
                    ss.step0_structured_memo_done = False
                    if "step0_voice_transcript_edit" in ss:
                        del ss["step0_voice_transcript_edit"]
                    if "step0_structured_memo_edit" in ss:
                        del ss["step0_structured_memo_edit"]
                    log_event(ss.run_id, "[STEP0 VOICE] pending transcript discarded")
                    st.rerun()


# ─── SIDEBAR ─────────────────────────────────────────────────────────────────
with st.sidebar:
    st.title("LoP Builder")
    st.caption("Phase 1 — Intake → Synthesis → Gate A → Mocks → Validation → Gate B → Dot-Dash → Gate C")
    st.divider()

    for _flag, _label in _STEPS:
        if _flag == "evals_done" and ss.slides_done:
            if ss.get(_flag):
                st.success(_label)
            else:
                st.text(_label)
            _render_evals_sidebar_status()
        elif ss.get(_flag):
            st.success(_label)
        else:
            st.text(_label)

    st.divider()
    st.caption(f"Faster model: {_MODEL_FAST}")
    st.caption(f"Better model: {_MODEL_FULL}")
    st.caption(f"Claude model: {_MODEL_CLAUDE}")
    st.caption(
        "Per-step model choice is configurable above each Run button."
    )
    st.caption(f"Run: {ss.run_id}")
    if _log_path.exists():
        st.caption(f"Log: runs/{_log_path.name}")
        try:
            log_bytes = _log_path.read_bytes()
            st.download_button(
                "Download run log",
                data=log_bytes,
                file_name=_log_path.name,
                mime="text/plain",
                key="dl_run_log",
                use_container_width=True,
            )
        except OSError:
            pass

    st.divider()
    st.caption("Consistency eval deep link")
    if ss.slides_done:
        if st.button(
            "Jump to Step 9 — Consistency eval (?evals=1)",
            use_container_width=True,
            key="sidebar_jump_evals",
        ):
            st.query_params["evals"] = "1"
            ss.evals_scroll_pending = True
            st.rerun()
    else:
        st.caption("Complete Step 8 to unlock the consistency eval jump link.")

    if st.button("Start Over", use_container_width=True):
        log_event(ss.run_id, "SESSION ended by user (Start Over)")
        for _k in list(ss.keys()):
            del ss[_k]
        st.rerun()

# ─── HEADER ──────────────────────────────────────────────────────────────────
st.title("LoP Builder — Phase 1")
st.caption(
    "Upload  →  Intake  →  Context  →  Synthesis  →  Gate A  →  "
    "Partner Answers (voice memo)  →  Validation  →  Gate B  →  Dot-Dash  →  Gate C  →  Consistency eval"
)
st.caption(
    "After slides are rendered, open **Step 9 — Consistency evaluation** via the sidebar "
    "or add **`?evals=1`** to this page URL (bookmark friendly)."
)

if not os.environ.get("OPENAI_API_KEY"):
    st.error(
        "**OPENAI_API_KEY not set.**  \n"
        "Copy `src/.env.example` to `src/.env`, add your API key, then restart: "
        "`streamlit run app.py`"
    )
    st.stop()

_done_count = sum(ss.get(f, False) for f, _ in _STEPS)
st.progress(_done_count / len(_STEPS))
st.divider()


# ─────────────────────────────────────────────────────────────────────────────
# STEP 0 — UPLOAD DOCUMENTS
# ─────────────────────────────────────────────────────────────────────────────
st.subheader("Step 0 — Upload Documents")

if ss.upload_done:
    with st.expander(
        f"Uploaded — {len(ss.processed_docs)} document(s) confirmed", expanded=False
    ):
        for d in ss.processed_docs:
            st.caption(f"{d['filename']}  [{d['doc_type']}]")
        if st.button("Change uploads (resets all downstream steps)"):
            ss.upload_done = False
            ss.processed_docs = []
            _reset_partner_voice_memos()
            _reset_from("intake")
            st.rerun()
else:
    st.markdown(
        "**What this does:** ingests the source documents the rest of the "
        "pipeline reads from. Files stay local to this session.  \n"
        "**What you do:** drag in the RFP, RFI, Gold Standard LoP, or any "
        "partner context files (Word / text / PDF); or record a partner "
        "voice memo below and have it transcribed; tag each item, then "
        "click **Confirm Uploads**.  \n"
        "_RFPs and RFIs are uncommon — most pursuits are driven by a "
        "partner voice memo or written brief, which goes in as **Partner "
        "Context**._"
    )

    uploaded = st.file_uploader(
        "Select files (PDF, DOCX, or TXT — multiple files allowed)",
        accept_multiple_files=True,
        type=["pdf", "docx", "txt"],
    )

    # ── Partner voice memo (whisper-1) ──
    _render_partner_voice_memo_block()

    pending_voice_memos: list[dict] = list(ss.partner_voice_memos)

    if uploaded or pending_voice_memos:
        st.markdown("**Tag each document:**")
        st.caption(
            "**RFP** — formal tender from the client.  \n"
            "**RFI** — pre-RFP exploratory document; signals what the client is "
            "looking for before the formal tender.  \n"
            "**Gold Standard LoP — Guidance** — a document describing what a "
            "strong LoP should look like (rules, template, checklist). Used as "
            "the binding reference for chapter shape.  \n"
            "**Gold Standard LoP — Examples** — one or more past LoPs to learn "
            "from. The intake agent synthesizes patterns from them into a guide.  \n"
            "**Partner Context** — partner brief / context dump (Word, text, "
            "PDF, or transcribed voice memo). Authoritative partner input but "
            "not a formal tender — use this when there is no RFP/RFI and the "
            "partner is dictating context directly."
        )
        hdr = st.columns([5, 3])
        hdr[0].markdown("Filename")
        hdr[1].markdown("Type")
        for f in uploaded or []:
            row = st.columns([5, 3])
            row[0].write(f.name)
            row[1].selectbox(
                "type",
                [
                    "RFP",
                    "RFI",
                    "Gold Standard LoP — Guidance",
                    "Gold Standard LoP — Examples",
                    "Partner Context",
                ],
                key=f"tag_{f.name}",
                label_visibility="collapsed",
            )
        for memo in pending_voice_memos:
            row = st.columns([5, 3])
            row[0].write(f"{memo['filename']}  _(voice memo, transcribed)_")
            row[1].text("Partner Context")

        if st.button("Confirm Uploads", type="primary"):
            docs: list[dict] = []
            for f in uploaded or []:
                doc_type = ss.get(f"tag_{f.name}", "RFP")
                raw = f.getvalue()

                if f.name.lower().endswith(".pdf"):
                    docs.append(
                        {
                            "filename": f.name,
                            "doc_type": doc_type,
                            "is_pdf": True,
                            "bytes": raw,
                            "text": "",
                        }
                    )
                elif f.name.lower().endswith(".docx"):
                    docx_doc = DocxDocument(BytesIO(raw))
                    text = "\n\n".join(
                        p.text for p in docx_doc.paragraphs if p.text.strip()
                    )
                    docs.append(
                        {
                            "filename": f.name,
                            "doc_type": doc_type,
                            "is_pdf": False,
                            "bytes": b"",
                            "text": text,
                        }
                    )
                else:  # .txt
                    docs.append(
                        {
                            "filename": f.name,
                            "doc_type": doc_type,
                            "is_pdf": False,
                            "bytes": b"",
                            "text": raw.decode("utf-8", errors="replace"),
                        }
                    )

            for memo in pending_voice_memos:
                docs.append(
                    {
                        "filename": memo["filename"],
                        "doc_type": "Partner Context",
                        "is_pdf": False,
                        "bytes": b"",
                        "text": memo["text"],
                    }
                )

            ss.processed_docs = docs
            ss.upload_done = True
            ss.partner_voice_memos = []
            ss.partner_voice_pending_transcript = ""
            ss.partner_voice_pending_filename = ""
            ss.step0_structured_memo = ""
            ss.step0_structured_memo_done = False
            if "step0_structured_memo_edit" in ss:
                del ss["step0_structured_memo_edit"]
            log_event(
                ss.run_id,
                f"[UPLOAD] {len(docs)} document(s) confirmed: "
                + ", ".join(f"{d['filename']} [{d['doc_type']}]" for d in docs),
            )
            st.rerun()

st.divider()


# ─────────────────────────────────────────────────────────────────────────────
# STEP 1 — INTAKE AGENT
# ─────────────────────────────────────────────────────────────────────────────
if not ss.upload_done:
    st.text("Step 1 — Intake Agent  [complete upload first]")
    st.divider()
else:
    st.subheader("Step 1 — Intake Agent")

    if ss.intake_done and ss.intake_package:
        pkg: IntakePackage = ss.intake_package
        with st.expander(
            f"Done — {pkg.client_name}  |  {pkg.industry}  |  {len(pkg.gap_list)} gaps found",
            expanded=False,
        ):
            c1, c2 = st.columns(2)
            c1.metric("Client", pkg.client_name)
            c1.metric("Industry", pkg.industry)
            c2.metric("Geography", pkg.geography)
            c2.metric("Gaps", len(pkg.gap_list))

            # Pursuit profile — type + competitive status
            _PURSUIT_LABELS = {
                "rfp": "RFP",
                "rfi_only": "RFI only",
                "rfp_with_rfi": "RFP + RFI",
                "partner_brief": "Partner brief (no RFP/RFI)",
                "unclear": "Unclear",
            }
            p1, p2 = st.columns(2)
            with p1:
                st.markdown("**Pursuit type**")
                st.info(_PURSUIT_LABELS.get(pkg.pursuit_type, pkg.pursuit_type))
            with p2:
                st.markdown("**Competitive status**")
                if pkg.competitive_status == "non_competitive":
                    st.success("Non-competitive")
                elif pkg.competitive_status == "competitive":
                    st.warning("Competitive")
                else:
                    st.info("Unclear")

            if pkg.competitor_firms:
                st.caption(
                    "Competing firms named: "
                    + ", ".join(pkg.competitor_firms)
                )
            elif pkg.competitive_status == "competitive":
                st.caption(
                    "Competitive pursuit — competing firms not named in inputs "
                    "(see gap list)."
                )

            if pkg.gold_standard_mode != "none":
                _GOLD_LABELS = {
                    "guidance": "Gold Standard — Guidance (authoritative)",
                    "examples_synthesis": "Gold Standard — Examples synthesis",
                }
                with st.expander(
                    _GOLD_LABELS.get(pkg.gold_standard_mode, pkg.gold_standard_mode),
                    expanded=False,
                ):
                    if pkg.gold_standard_mode == "guidance" and pkg.gold_standard_guidance:
                        st.markdown(pkg.gold_standard_guidance)
                    elif (
                        pkg.gold_standard_mode == "examples_synthesis"
                        and pkg.gold_standard_synthesis
                    ):
                        st.markdown(pkg.gold_standard_synthesis)
                    else:
                        st.caption("(no content extracted)")

            if pkg.rfi_signals:
                st.markdown("**Pre-RFP RFI signals:**")
                for s in pkg.rfi_signals:
                    st.markdown(f"- {s}")

            st.markdown("**Chapter coverage:**")
            for b in pkg.chapter_buckets:
                row = st.columns([5, 1])
                row[0].markdown(f"**{b.chapter}**  — {b.notes}" if b.notes else f"**{b.chapter}**")
                with row[1]:
                    _quality_badge(b.quality)

            if pkg.gap_list:
                st.markdown("**Gap list:**")
                for g in pkg.gap_list:
                    st.markdown(f"- {g}")

            if pkg.key_facts:
                st.markdown("**Key facts:**")
                for kf in pkg.key_facts:
                    st.markdown(f"- {kf}")

        if st.button("Re-run Intake Agent", key="rerun_intake"):
            _reset_from("intake")
            st.rerun()
    else:
        st.markdown(
            "**What this does:** the Intake Agent reads your documents, "
            "classifies content by LoP chapter, extracts key facts and RFP "
            "requirements, and produces a gap list.  \n"
            "**What you do:** confirm whether this pursuit is competitive or "
            "non-competitive, list any firms you know are competing, then "
            "click Run."
        )

        st.radio(
            "Is this a competitive or non-competitive pursuit?",
            options=["competitive", "non_competitive"],
            format_func=lambda v: (
                "Competitive — other firms tendering"
                if v == "competitive"
                else "Non-competitive — sole-source / relationship-led"
            ),
            index=None,
            horizontal=True,
            key="intake_competitive_status",
        )

        if ss.get("intake_competitive_status") == "competitive":
            st.text_area(
                "Competing firms (optional, one per line or comma-separated)",
                placeholder="e.g. BCG, Bain, Roland Berger",
                height=90,
                key="intake_competitor_firms",
            )

        run_disabled = ss.get("intake_competitive_status") not in {
            "competitive",
            "non_competitive",
        }
        intake_model = _render_model_selector("intake")
        if st.button(
            "Run Intake Agent",
            type="primary",
            key="run_intake",
            disabled=run_disabled,
        ):
            comp_status = ss.get("intake_competitive_status")
            firms = _parse_firms_input(
                ss.get("intake_competitor_firms", "")
                if comp_status == "competitive"
                else ""
            )

            with st.spinner(
                "Intake Agent running — classifying documents and identifying gaps "
                "(may take 30–90 seconds)..."
            ):
                try:
                    doc_summary = ", ".join(
                        f"{d['filename']} [{d['doc_type']}]"
                        for d in ss.processed_docs
                    )
                    start_summary = (
                        f"{doc_summary} | partner: {comp_status}"
                        + (f" | firms={len(firms)}" if firms else "")
                    )
                    log_agent_start(ss.run_id, "intake-agent", start_summary)

                    partner_block = (
                        "PARTNER-PROVIDED INPUTS — these are AUTHORITATIVE, "
                        "treat as ground truth. Echo them verbatim into the "
                        "output. Do NOT silently change them; if the documents "
                        "contradict them, surface the contradiction in the "
                        "relevant chapter notes AND raise a gap-list item.\n"
                        f"- competitive_status: {comp_status}\n"
                        + (
                            f"- competitor_firms (partner-listed): "
                            f"{json.dumps(firms)}\n"
                            if firms
                            else "- competitor_firms (partner-listed): []\n"
                        )
                    )

                    result = run_agent(
                        "intake-agent",
                        user_message=(
                            f"{partner_block}\n"
                            "Analyse the provided documents and produce a structured "
                            "intake package. Each document is labelled with its type "
                            "(RFP, RFI, Partner Context, Gold Standard LoP — "
                            "Guidance, or Gold Standard LoP — Examples). Use the "
                            "partner-provided competitive_status and "
                            "competitor_firms verbatim, determine pursuit_type "
                            "(use 'partner_brief' when only Partner Context is "
                            "present with no RFP/RFI), route Gold Standard "
                            "content per its tag (guidance is authoritative; "
                            "examples drive a pattern synthesis), populate "
                            "rfi_signals when an RFI is present, treat Partner "
                            "Context as authoritative partner input but never "
                            "extract rfp_requirements from it, and classify all "
                            "RFP/RFI/Partner Context content across the nine LoP "
                            "chapters using the canonical definitions below.\n\n"
                            f"{render_chapter_brief()}"
                        ),
                        files=ss.processed_docs,
                        use_extended_thinking=True,
                        model=intake_model,
                    )
                    log_agent_result(ss.run_id, "intake-agent", result)
                    ss.intake_package = IntakePackage.model_validate(result)
                    ss.intake_done = True
                    st.rerun()
                except Exception as exc:
                    traceback.print_exc()
                    log_event(ss.run_id, f"[INTAKE AGENT] ERROR: {exc}")
                    st.error(f"Intake Agent failed: {exc}")

    st.divider()


# ─────────────────────────────────────────────────────────────────────────────
# STEP 2 — CONTEXT AGENT (web search: Quick / Deep + model-knowledge fallback)
# ─────────────────────────────────────────────────────────────────────────────


def _stamp_web_citations(result: dict, retrieved_at: str) -> dict:
    """
    Stamp `retrieved_at` on every Citation where `kind == "web"` so the UI
    can show the date clearly. Idempotent and safe on missing fields.
    """
    citations = result.get("citations", []) or []
    for c in citations:
        if isinstance(c, dict) and c.get("kind") == "web" and not c.get("retrieved_at"):
            c["retrieved_at"] = retrieved_at
    return result


def _ensure_field(result: dict, key: str, default) -> None:
    """Set `result[key]` to `default` if missing or empty. In-place."""
    if not result.get(key):
        result[key] = default


# Deep Research model identifiers. The user picks between these in the
# Step 2 UI; the choice is passed into call_deep_research and overrides
# the DEEP_RESEARCH_MODEL env var for that run. The env var stays as the
# fallback default when the user has not made a UI choice yet.
_DR_MODEL_MINI = "o4-mini-deep-research-2025-06-26"
_DR_MODEL_FULL = "o3-deep-research-2025-06-26"


# Dedicated research-phase prompt for the Deep Research model. Deliberately
# different from the context-agent system prompt — Deep Research produces
# a long-form prose briefing, not JSON, and the framing avoids LoP-pipeline
# jargon AND positioning/competitive-advice language that has triggered
# content-policy refusals in past runs. The structuring step (second pass)
# uses the context-agent system prompt unchanged, which is correct.
_DEEP_RESEARCH_SYSTEM_PROMPT = (
    "You are an autonomous public-record research analyst. The user "
    "describes a company and a research scope. Plan and run web searches "
    "across public sources (company filings, press releases, regulatory "
    "bodies, reputable news, official websites, trade press), then write "
    "a comprehensive long-form briefing.\n\n"
    "Length and depth — important: this is not a short summary. The "
    "downstream consumer needs the specifics, not a digest. Aim for at "
    "least 4,000 words and do not artificially compress; let the depth "
    "match what the research material justifies. Each of the six "
    "dimensions below should receive a full treatment of at least "
    "500–800 words with multiple cited sources, named entities, and "
    "concrete figures. If a dimension surfaced more material than that "
    "(common after 100+ searches), use it. Do not truncate, paraphrase "
    "away, or omit specific facts to stay short — preserve every "
    "specific number, named individual, dated event, and named "
    "regulation. Density and faithfulness to the sources are the goal.\n\n"
    "Cover these dimensions in prose, with inline URL citations on every "
    "specific factual claim, named individual, named figure, and named "
    "regulation:\n\n"
    "1. CLIENT PROFILE — business model, scale, primary markets and "
    "geographies, ownership structure, recent strategic priorities "
    "relevant to the research scope, leadership where material.\n"
    "2. MARKET TRENDS — three to five trends shaping the client's "
    "industry that bear on the research scope, with sources for any "
    "quantified claim.\n"
    "3. PEERS AND PUBLIC POSITIONING — describe relevant peers and any "
    "other firms the user names. For each named firm, summarise its "
    "publicly disclosed positioning, recent public moves, and any "
    "publicly stated approach to the same problem area. Stay strictly "
    "factual and source-grounded; this is public-record research, not "
    "tactical advice on how to compete with any firm.\n"
    "4. SECTOR CHALLENGES — recurring strategic and operational "
    "challenges in this sector that bear on the research scope.\n"
    "5. RECENT SIGNALS — specific recent material events about the "
    "client itself (last ~24 months): leadership changes (CEO, CFO, "
    "board), M&A and divestments, financial milestones (earnings, "
    "ratings, refinancings), operational events (restructurings, "
    "technology rollouts, named partnerships), risk events (lawsuits, "
    "regulatory enforcement, cyber incidents, recalls). Include the "
    "date and the source for each. Do not artificially cap how many you "
    "include — every material event surfaced in the research belongs "
    "here.\n"
    "6. REGULATORY ENVIRONMENT — laws, regulations, directives, and "
    "enforcement actions that materially affect the client's industry "
    "and the research scope, with effective dates and named sources. "
    "Include every material item; do not artificially cap.\n\n"
    "Style and constraints:\n"
    "- Prose only. Do not return JSON, YAML, or any other structured "
    "format.\n"
    "- Cite a public URL for every specific claim. If a claim cannot be "
    "sourced, omit it or hedge.\n"
    "- Do not invent figures, dates, or quotes.\n"
    "- Do not provide medical, legal, financial, or political advice.\n"
    "- Preserve every specific fact, named individual, named figure, "
    "dated event, and named regulation surfaced in your research. "
    "Treat all subjects as public-record research."
)


def _build_deep_research_user_message(pkg: IntakePackage, extra_ctx: str) -> str:
    """
    Build a research-phase user message scoped to public-record research.
    Deliberately avoids LoP-pipeline jargon (no "build a context document",
    no "Mode: deep", no JSON examples) — that is the structuring step's job.
    """
    competitor_block = ""
    firms = pkg.competitor_firms or []
    if firms:
        firms_str = ", ".join(firms)
        competitor_block = (
            f"\n\nOther firms named by the user — research each firm's "
            f"public positioning, recent public moves, and any publicly "
            f"stated approach to the problem area. This is factual "
            f"public-record research; do not produce tactical advice on "
            f"how to compete with any of these firms: {firms_str}"
        )

    extra_block = ""
    extra_clean = (extra_ctx or "").strip()
    if extra_clean:
        extra_block = (
            f"\n\nAdditional research context from the user (may name "
            f"recent events or specific angles to investigate):\n"
            f"\"\"\"\n{extra_clean}\n\"\"\""
        )

    return (
        f"Research scope:\n"
        f"  Subject: {pkg.client_name}\n"
        f"  Industry: {pkg.industry}\n"
        f"  Geography: {pkg.geography}\n"
        f"  Topic: {pkg.problem_area}"
        f"{competitor_block}"
        f"{extra_block}\n\n"
        f"Produce the briefing as instructed in the system prompt. Run "
        f"as many searches as you need, then write the final report."
    )


def _run_deep_research_path(
    pkg: IntakePackage,
    extra_ctx: str,
    progress_cb,
    model: str | None = None,
    structuring_model: str | None = None,
) -> dict:
    """
    Run the Deep Research engine then structure the resulting report into a
    ContextDoc dict via the existing context-agent (no tools, JSON mode).

    Stages:
      1. call_deep_research(...) -> (report_text, url_citations)
         Uses a dedicated research-phase system prompt + user message,
         NOT the context-agent's structuring system prompt.
      2. run_agent("context-agent", ...) on a structuring user message
         that embeds the report and the URL citation list. This is where
         the context-agent system prompt is applied.
      3. Stamp `search_mode = "deep"` and `additional_context_used`.

    `model` overrides DEEP_RESEARCH_MODEL for this run when supplied
    (typically driven by the UI toggle between the mini and full
    Deep Research variants).

    Raises DeepResearchUnavailable on stage 1 failure so the call-site can
    fall back to call_with_web_search.
    """
    research_user_msg = _build_deep_research_user_message(pkg, extra_ctx)

    report_text, url_citations = call_deep_research(
        system_prompt=_DEEP_RESEARCH_SYSTEM_PROMPT,
        user_message=research_user_msg,
        model=model,
        progress_cb=progress_cb,
    )

    log_event(
        ss.run_id,
        f"[CONTEXT AGENT] Deep Research returned report "
        f"({len(report_text)} chars, {len(url_citations)} URL citation(s))",
    )

    if url_citations:
        url_lines: list[str] = []
        for i, c in enumerate(url_citations, start=1):
            title = (c.get("title") or "").strip() or "(no title)"
            url = (c.get("url") or "").strip()
            url_lines.append(f"{i}. {title} — {url}")
        urls_block = "\n".join(url_lines)
    else:
        urls_block = "(no URL citations were attached to the research report)"

    intake_dict = pkg.model_dump()
    structuring_msg = (
        "STRUCTURING TASK — convert the research report below into a "
        "ContextDoc JSON. Mode: **deep** (research path). Preserve every "
        "URL citation. Populate `recent_signals`, `regulatory_environment`, "
        "and `chapter_takeaways`. Do NOT add facts the report does not "
        "contain. Do NOT run any tools — the research is already complete.\n\n"
        "FAITHFULNESS RULES — read carefully. Compression is NOT the "
        "goal of this step. Density and faithfulness are. Specifically:\n"
        "- Preserve every specific number, named individual, dated event, "
        "and named regulation that appears in the research report.\n"
        "- Do not paraphrase specifics into generic phrasing to stay "
        "short; honour the per-section Deep word ranges as floors not "
        "ceilings.\n"
        "- Include every material recent event in `recent_signals` (10–25 "
        "entries is normal when the research is rich); every material "
        "regulation in `regulatory_environment` (5–10 entries is normal); "
        "do not drop sourced items to stay short.\n"
        "- Every Citation MUST have a non-empty `claim` stating the "
        "actual fact, even on `model_knowledge` citations.\n\n"
        f"Intake Package:\n{json.dumps(intake_dict, indent=2)}\n\n"
        f"Additional context from the user (may be empty):\n"
        f"\"\"\"\n{extra_ctx}\n\"\"\"\n\n"
        "RESEARCH REPORT (prose with inline citation markers — treat as "
        "the authoritative factual foundation):\n"
        f"\"\"\"\n{report_text}\n\"\"\"\n\n"
        "URL CITATIONS FROM THE RESEARCH (use these URLs verbatim in "
        "ContextDoc.citations where kind == \"web\"; do not invent new "
        "URLs):\n"
        f"{urls_block}"
    )

    parsed = run_agent(
        "context-agent",
        user_message=structuring_msg,
        files=None,
        use_extended_thinking=True,
        model=structuring_model,
    )
    parsed["search_mode"] = "deep"
    return parsed


if not ss.intake_done:
    st.text("Step 2 — Context Agent  [complete intake first]")
    st.divider()
else:
    st.subheader("Step 2 — Context Agent")

    if ss.context_done and ss.context_doc:
        ctx: ContextDoc = ss.context_doc

        web_citations = [c for c in ctx.citations if c.kind == "web"]
        model_citations = [c for c in ctx.citations if c.kind != "web"]

        with st.expander("Done — context document ready", expanded=False):
            # Four-state search-mode banner.
            if ctx.search_mode == "deep":
                st.success(
                    f"Deep Research — autonomous multi-source research, "
                    f"{len(web_citations)} web source(s) cited"
                )
            elif ctx.search_mode == "deep_fallback":
                st.warning(
                    "Deep Research was unavailable — fell back to enriched "
                    f"web search. {len(web_citations)} web source(s) cited."
                )
            elif ctx.search_mode == "quick":
                st.success(
                    f"Quick web search — {len(web_citations)} web source(s) cited"
                )
            else:
                st.warning(
                    "Web search was unavailable on this run — fell back to model "
                    "knowledge. Treat all claims as directional and validate "
                    "before client-facing use."
                )

            if ctx.additional_context_used:
                with st.expander("Your additional context", expanded=False):
                    st.markdown(ctx.additional_context_used)

            st.markdown("**Client profile**")
            st.markdown(ctx.client_profile)
            st.markdown("**Market trends**")
            st.markdown(ctx.market_trends)
            st.markdown("**Competitive landscape**")
            st.markdown(ctx.competitive_landscape)
            st.markdown("**Relevant challenges**")
            st.markdown(ctx.relevant_challenges)

            # ── Recent signals (grouped by category) ──
            if ctx.recent_signals:
                st.markdown("---")
                st.markdown(
                    f"**Recent signals — {len(ctx.recent_signals)} event(s)**"
                )

                _CATEGORY_ORDER = [
                    "leadership_change",
                    "m&a",
                    "financial",
                    "operational",
                    "risk_event",
                    "news",
                    "",
                ]
                _CATEGORY_LABELS = {
                    "leadership_change": "Leadership change",
                    "m&a": "M&A",
                    "financial": "Financial",
                    "operational": "Operational",
                    "risk_event": "Risk event",
                    "news": "News",
                    "": "Other",
                }

                grouped: dict[str, list] = {}
                for sig in ctx.recent_signals:
                    key = (sig.category or "").strip().lower()
                    if key not in _CATEGORY_LABELS:
                        key = ""
                    grouped.setdefault(key, []).append(sig)

                for cat in _CATEGORY_ORDER:
                    bucket = grouped.get(cat, [])
                    if not bucket:
                        continue
                    st.markdown(f"_{_CATEGORY_LABELS[cat]}_")
                    for sig in bucket:
                        date_str = f" · {sig.date}" if sig.date else ""
                        st.markdown(
                            f"- **{sig.headline}**{date_str}"
                        )
                        if sig.detail:
                            st.caption(sig.detail)
                        if sig.citation_urls:
                            link_md = ", ".join(
                                f"[source {i + 1}]({u})"
                                for i, u in enumerate(sig.citation_urls)
                            )
                            st.caption(f"Sources: {link_md}")

            # ── Regulatory environment ──
            if ctx.regulatory_environment:
                st.markdown("---")
                st.markdown(
                    f"**Regulatory environment — "
                    f"{len(ctx.regulatory_environment)} item(s)**"
                )
                for item in ctx.regulatory_environment:
                    eff = (
                        f" ({item.effective_date})"
                        if item.effective_date
                        else ""
                    )
                    st.markdown(f"**{item.topic}**{eff}")
                    if item.summary:
                        st.caption(f"Summary: {item.summary}")
                    if item.client_impact:
                        st.caption(f"Client impact: {item.client_impact}")
                    if item.citation_urls:
                        link_md = ", ".join(
                            f"[source {i + 1}]({u})"
                            for i, u in enumerate(item.citation_urls)
                        )
                        st.caption(f"Sources: {link_md}")

            # ── Chapter takeaways ──
            ct = ctx.chapter_takeaways
            ct_pairs = [
                ("Context and objectives", ct.context_and_objectives),
                ("Why McKinsey", ct.why_mckinsey),
                ("Approach", ct.approach),
                ("Market trends", ct.market_trends),
                ("Credentials", ct.credentials),
            ]
            populated_pairs = [(k, v) for k, v in ct_pairs if v.strip()]
            if populated_pairs:
                st.markdown("---")
                st.markdown(
                    f"**Chapter takeaways — {len(populated_pairs)}/5 chapter(s) "
                    f"with research-driven guidance**"
                )
                for label, text in populated_pairs:
                    with st.expander(label, expanded=False):
                        st.markdown(text)

            if web_citations or model_citations:
                st.markdown("---")
                st.markdown("**Citations**")
                for c in web_citations:
                    label = c.title or c.url
                    line = f"- [{label}]({c.url}) — {c.claim}"
                    if c.source_note:
                        line += f"  \n  _{c.source_note}_"
                    if c.retrieved_at:
                        line += f"  \n  _retrieved {c.retrieved_at}_"
                    st.markdown(line)
                for c in model_citations:
                    st.markdown(f"- {c.claim} — _{c.source_note}_")

            if ctx.searches_performed:
                with st.expander(
                    f"Searches performed ({len(ctx.searches_performed)})",
                    expanded=False,
                ):
                    for q in ctx.searches_performed:
                        st.markdown(f"- {q}")

            if ctx.evidence_gaps:
                st.warning(
                    "Evidence gaps — validate before using in LoP:  \n"
                    + "\n".join(f"- {g}" for g in ctx.evidence_gaps)
                )
            st.info(f"Knowledge note: {ctx.knowledge_cutoff_note}")

        if st.button("Re-run Context Agent", key="rerun_context"):
            _reset_from("context")
            st.rerun()
    else:
        st.markdown(
            "**What this does:** the Context Agent is your strongest "
            "research tool — Quick gives a tailored company/market snapshot "
            "in ~30 seconds; Deep runs OpenAI's Deep Research model "
            "autonomously across the public web for ~5–15 minutes to "
            "produce a sourced, LoP-aligned briefing covering company, "
            "market, named competitors, recent leadership / M&A / financial "
            "signals, regulatory environment, and chapter-by-chapter "
            "takeaways for the LoP itself.  \n"
            "**What you do:** add any extra context the agent should know "
            "(optional), pick Quick or Deep, then click Run."
        )

        st.text_area(
            "Additional context (optional)",
            placeholder=(
                "Anything else the agent should know — relationship history, "
                "internal hypotheses, recent news the partner mentioned, "
                "specific angles to investigate..."
            ),
            height=110,
            key="context_extra_input",
        )

        st.radio(
            "Search depth",
            options=["quick", "deep"],
            format_func=lambda m: (
                "Quick — targeted web search, ~30 seconds"
                if m == "quick"
                else "Deep — Deep Research model"
            ),
            horizontal=True,
            key="context_search_mode",
        )

        if ss.get("context_search_mode") == "deep":
            # Pre-seed the model toggle from the env var (if any) so users
            # who configured DEEP_RESEARCH_MODEL in .env see their preference
            # honored on first render.
            if "deep_research_model_choice" not in ss:
                env_default = os.environ.get("DEEP_RESEARCH_MODEL", "")
                ss["deep_research_model_choice"] = (
                    "full" if env_default == _DR_MODEL_FULL else "mini"
                )

            st.radio(
                "Deep Research model — pick depth vs. time",
                options=["mini", "full"],
                format_func=lambda v: (
                    "Faster — o4-mini-deep-research (~5–10 min, lighter "
                    "synthesis, lower cost)"
                    if v == "mini"
                    else "Deeper — o3-deep-research (~10–25 min, richer "
                    "synthesis, higher cost)"
                ),
                key="deep_research_model_choice",
            )

            chosen_label = (
                "o3-deep-research (deeper)"
                if ss.get("deep_research_model_choice") == "full"
                else "o4-mini-deep-research (faster)"
            )
            st.info(
                f"Deep mode launches OpenAI's Deep Research model "
                f"(**{chosen_label}**), which autonomously plans and runs "
                f"many web searches across company filings, news, "
                f"competitor moves, regulations, and sector trends. The "
                f"result is a sourced briefing with recent signals "
                f"(leadership / M&A / financial / operational / risk), "
                f"regulatory environment with client-specific impact, and "
                f"chapter-by-chapter takeaways for Context, Why McKinsey, "
                f"Approach, Market Trends, and Credentials. Pick **Faster** "
                f"for a first pass on tight timelines; pick **Deeper** when "
                f"you want maximum density and have ~20 min to wait. If "
                f"Deep Research is not exposed by the gateway, the agent "
                f"falls back automatically to enriched web search "
                f"(~1–3 minutes)."
            )

        context_model = _render_model_selector("context")
        if st.button("Run Context Agent", type="primary", key="run_context"):
            mode = ss.get("context_search_mode", "quick")
            extra_ctx = (ss.get("context_extra_input", "") or "").strip()

            # Resolve the chosen Deep Research model (only used in Deep mode).
            dr_choice = ss.get("deep_research_model_choice", "mini")
            dr_model = _DR_MODEL_FULL if dr_choice == "full" else _DR_MODEL_MINI

            if mode == "quick":
                spinner_msg = (
                    "Quick web search running — targeted queries, typically "
                    "~30 seconds..."
                )
            elif dr_choice == "full":
                spinner_msg = (
                    "Deep Research running on o3-deep-research — deeper, "
                    "richer synthesis, typically 10–25 minutes. You can "
                    "leave this tab open and come back."
                )
            else:
                spinner_msg = (
                    "Deep Research running on o4-mini-deep-research — "
                    "lighter synthesis, typically 5–10 minutes. You can "
                    "leave this tab open and come back."
                )

            with st.spinner(spinner_msg):
                progress_slot = st.empty()
                try:
                    pkg = ss.intake_package
                    log_agent_start(
                        ss.run_id, "context-agent",
                        f"{pkg.client_name} | {pkg.industry} | "
                        f"{pkg.geography} | mode={mode}"
                        + (f" | dr_model={dr_model}" if mode == "deep" else "")
                        + (f" | extra_ctx={len(extra_ctx)} chars"
                           if extra_ctx else ""),
                    )
                    intake_dict = pkg.model_dump()

                    user_msg = (
                        f"Build a context document for this LoP pursuit. "
                        f"Mode: **{mode}**.\n\n"
                        f"Intake Package:\n{json.dumps(intake_dict, indent=2)}\n\n"
                        f"Additional context from the user (may be empty):\n"
                        f"\"\"\"\n{extra_ctx}\n\"\"\""
                    )

                    used_path = ""
                    fallback_reason = ""
                    parsed: dict | None = None

                    def _progress(msg: str) -> None:
                        log_event(ss.run_id, f"[CONTEXT AGENT] {msg}")
                        try:
                            progress_slot.info(msg)
                        except Exception:
                            pass

                    if mode == "deep":
                        # Stage 1: Deep Research engine.
                        # Note: Deep Research uses its own research-phase
                        # system prompt + user message (built inside the
                        # helper), NOT `user_msg`. `user_msg` is reserved
                        # for the web_search/model-knowledge fallback paths
                        # which DO use the context-agent's structuring
                        # system prompt. The model identifier comes from
                        # the UI toggle (Faster vs Deeper).
                        try:
                            parsed = _run_deep_research_path(
                                pkg,
                                extra_ctx,
                                _progress,
                                model=dr_model,
                                structuring_model=context_model,
                            )
                            used_path = "deep"
                        except DeepResearchUnavailable as exc:
                            fallback_reason = str(exc)
                            log_event(
                                ss.run_id,
                                f"[CONTEXT AGENT] Deep Research unavailable, "
                                f"falling back to enriched web_search: {exc}",
                            )
                            try:
                                progress_slot.warning(
                                    "Deep Research not available; running "
                                    "enriched web search instead — typically "
                                    "1–3 minutes."
                                )
                            except Exception:
                                pass

                            # Stage 2: enriched web_search fallback.
                            try:
                                spec = load_agent_spec("context-agent")
                                parsed, _ann = call_with_web_search(
                                    system_prompt=spec["system_prompt"],
                                    user_message=user_msg,
                                    output_schema=spec["output_schema"],
                                    mode="deep",
                                    model=context_model,
                                )
                                parsed["search_mode"] = "deep_fallback"
                                used_path = "deep_fallback"
                            except WebSearchUnavailable as exc2:
                                fallback_reason = (
                                    f"deep research: {exc} | "
                                    f"web_search: {exc2}"
                                )
                                log_event(
                                    ss.run_id,
                                    f"[CONTEXT AGENT] web_search also "
                                    f"unavailable, falling back to model "
                                    f"knowledge: {exc2}",
                                )
                                parsed = run_agent(
                                    "context-agent",
                                    user_message=user_msg,
                                    files=None,
                                    use_extended_thinking=True,
                                    model=context_model,
                                )
                                parsed["search_mode"] = (
                                    "model_knowledge_fallback"
                                )
                                used_path = "model_knowledge_fallback"
                    else:
                        # Quick mode — go straight to web_search.
                        try:
                            spec = load_agent_spec("context-agent")
                            parsed, _ann = call_with_web_search(
                                system_prompt=spec["system_prompt"],
                                user_message=user_msg,
                                output_schema=spec["output_schema"],
                                mode="quick",
                                model=context_model,
                            )
                            parsed["search_mode"] = "quick"
                            used_path = "quick"
                        except WebSearchUnavailable as exc:
                            fallback_reason = str(exc)
                            log_event(
                                ss.run_id,
                                f"[CONTEXT AGENT] web_search unavailable, "
                                f"falling back to model knowledge: {exc}",
                            )
                            parsed = run_agent(
                                "context-agent",
                                user_message=user_msg,
                                files=None,
                                use_extended_thinking=True,
                                model=context_model,
                            )
                            parsed["search_mode"] = (
                                "model_knowledge_fallback"
                            )
                            used_path = "model_knowledge_fallback"

                    assert parsed is not None

                    # Ensure audit-trail fields are populated regardless of path.
                    _ensure_field(parsed, "additional_context_used", extra_ctx)
                    _stamp_web_citations(parsed, date.today().isoformat())

                    log_event(
                        ss.run_id,
                        f"[CONTEXT AGENT] used_path={used_path}",
                    )
                    log_agent_result(ss.run_id, "context-agent", parsed)
                    ss.context_doc = ContextDoc.model_validate(parsed)
                    ss.context_done = True

                    if used_path == "deep_fallback":
                        st.warning(
                            "Deep Research was not available on this run — "
                            "fell back to enriched web search. "
                            f"Reason: {fallback_reason}"
                        )
                    elif used_path == "model_knowledge_fallback":
                        st.warning(
                            "Live web search was not available on this run — "
                            "fell back to model knowledge. Treat all claims as "
                            "directional and validate before client-facing use. "
                            f"Reason: {fallback_reason}"
                        )

                    progress_slot.empty()
                    st.rerun()
                except Exception as exc:
                    traceback.print_exc()
                    log_event(ss.run_id, f"[CONTEXT AGENT] ERROR: {exc}")
                    progress_slot.empty()
                    st.error(f"Context Agent failed: {exc}")

    st.divider()


# ─────────────────────────────────────────────────────────────────────────────
# STEP 3 — SYNTHESIS AGENT
# ─────────────────────────────────────────────────────────────────────────────
if not ss.context_done:
    st.text("Step 3 — Synthesis Agent  [complete context first]")
    st.divider()
else:
    st.subheader("Step 3 — Synthesis Agent")

    if ss.synthesis_done and ss.synthesis_doc:
        syn: SynthesisDoc = ss.synthesis_doc
        with st.expander("Done — synthesis ready for Gate A review", expanded=False):
            st.markdown("**Problem statement**")
            st.info(syn.problem_statement)
            st.markdown("**Win themes**")
            for wt in syn.win_themes:
                st.markdown(f"- {wt}")
            st.caption(
                f"Question list: {len(syn.question_list.questions)} questions generated"
            )
    else:
        st.markdown(
            "**What this does:** the Synthesis Agent merges the intake "
            "package and context document into a synthesis brief, problem "
            "statement, win themes, and a partner question list (the "
            "questions you and the partner will answer at Step 4).  \n"
            "**What you do:** no inputs needed here — click Run, then "
            "review the question list at Gate A."
        )
        synthesis_model = _render_model_selector("synthesis")
        if st.button("Run Synthesis Agent", type="primary", key="run_synthesis"):
            with st.spinner(
                "Synthesis Agent running — producing brief and question list "
                "(30–90 seconds)..."
            ):
                try:
                    log_agent_start(ss.run_id, "synthesis-agent")
                    intake_dict = ss.intake_package.model_dump()
                    context_dict = ss.context_doc.model_dump()
                    result = run_agent(
                        "synthesis-agent",
                        user_message=(
                            "Produce a synthesis brief, problem statement, win themes, "
                            "and partner question list.\n\n"
                            f"Intake Package:\n{json.dumps(intake_dict, indent=2)}\n\n"
                            f"Context Document:\n{json.dumps(context_dict, indent=2)}"
                        ),
                        files=None,
                        use_extended_thinking=True,
                        model=synthesis_model,
                    )
                    log_agent_result(ss.run_id, "synthesis-agent", result)
                    ss.synthesis_doc = SynthesisDoc.model_validate(result)
                    ss.synthesis_done = True
                    st.rerun()
                except Exception as exc:
                    traceback.print_exc()
                    log_event(ss.run_id, f"[SYNTHESIS AGENT] ERROR: {exc}")
                    st.error(f"Synthesis Agent failed: {exc}")

    st.divider()


# ─────────────────────────────────────────────────────────────────────────────
# GATE A — REVIEW AND APPROVE QUESTION LIST
# ─────────────────────────────────────────────────────────────────────────────
if not ss.synthesis_done:
    st.text("Gate A — Review Question List  [complete synthesis first]")
    st.divider()
else:
    st.subheader("Gate A — Review and Approve Question List")

    if ss.gate_a_done:
        q_count = len(ss.synthesis_doc.question_list.questions)
        with st.expander(f"Approved — {q_count} questions confirmed", expanded=False):
            for q in ss.synthesis_doc.question_list.questions:
                st.markdown(f"**{q.id}** ({q.chapter})  \n{q.question}")
        if st.button("Re-open Gate A", key="reopen_gate_a"):
            ss.gate_a_done = False
            _reset_from("answers")
            st.rerun()
    else:
        syn = ss.synthesis_doc

        st.markdown(
            "**What this is:** the human checkpoint before any partner work. "
            "The Synthesis Agent has produced the question list the partner "
            "will answer at Step 4.  \n"
            "**What you do:** read the synthesis brief and question list "
            "below, edit any question wording inline, then click **Approve** "
            "to move on. Use **Iterate** if questions need to be added, "
            "removed, or reframed."
        )

        # Synthesis overview
        with st.expander("Synthesis brief", expanded=False):
            st.markdown(syn.brief_summary)

        st.markdown("**Problem statement**")
        st.info(syn.problem_statement)

        st.markdown("**Win themes**")
        for wt in syn.win_themes:
            st.markdown(f"- {wt}")

        st.divider()
        st.markdown(
            "**Question list** — edit question text inline, then Approve or use "
            "Iterate to re-run synthesis with notes."
        )
        st.caption(
            "Tip: edit any question text directly in the box below. "
            "Changes are saved when you click Approve."
        )

        for i, q in enumerate(syn.question_list.questions):
            c_label, c_field = st.columns([2, 5])
            c_label.markdown(f"**{q.id}**  \n{q.chapter}")
            c_label.caption(f"Type: {q.expected_answer_type}")
            c_field.text_area(
                f"q_{i}",
                value=ss.get(f"q_edit_{i}", q.question),
                key=f"q_edit_{i}",
                height=80,
                label_visibility="collapsed",
            )
            c_field.caption(f"Why asked: {q.why_asked}")

        st.markdown("---")
        col_approve, col_iter = st.columns([1, 2])

        with col_approve:
            if st.button("Approve Question List", type="primary", key="approve_gate_a"):
                updated_qs = [
                    Question(
                        id=q.id,
                        chapter=q.chapter,
                        question=ss.get(f"q_edit_{i}", q.question),
                        why_asked=q.why_asked,
                        expected_answer_type=q.expected_answer_type,
                    )
                    for i, q in enumerate(syn.question_list.questions)
                ]
                syn.question_list = QuestionList(questions=updated_qs)
                ss.synthesis_doc = syn
                ss.gate_a_done = True
                log_gate_event(
                    ss.run_id, "gate-a", "APPROVED",
                    f"{len(updated_qs)} question(s) confirmed",
                )
                _reset_from("answers")
                st.rerun()

        with col_iter:
            with st.expander("Iterate synthesis with notes"):
                notes = st.text_area(
                    "Describe what to change (e.g. 'Add a question on budget; "
                    "remove Q9 as it overlaps with Q3')",
                    key="gate_a_notes",
                    height=90,
                )
                if st.button("Re-run Synthesis with These Notes", key="rerun_synthesis"):
                    if notes.strip():
                        with st.spinner("Re-running Synthesis Agent..."):
                            try:
                                log_gate_event(
                                    ss.run_id, "gate-a", "ITERATE",
                                    f"notes: {notes[:120]}",
                                )
                                log_agent_start(ss.run_id, "synthesis-agent", "re-run with user notes")
                                intake_dict = ss.intake_package.model_dump()
                                context_dict = ss.context_doc.model_dump()

                                # Capture any inline edits the user made to question text
                                # so the agent revises against the user's current draft,
                                # not the original model output.
                                edited_questions = [
                                    Question(
                                        id=q.id,
                                        chapter=q.chapter,
                                        question=ss.get(f"q_edit_{i}", q.question),
                                        why_asked=q.why_asked,
                                        expected_answer_type=q.expected_answer_type,
                                    )
                                    for i, q in enumerate(
                                        ss.synthesis_doc.question_list.questions
                                    )
                                ]
                                current_draft = SynthesisDoc(
                                    brief_summary=ss.synthesis_doc.brief_summary,
                                    problem_statement=ss.synthesis_doc.problem_statement,
                                    win_themes=ss.synthesis_doc.win_themes,
                                    question_list=QuestionList(questions=edited_questions),
                                )
                                draft_dict = current_draft.model_dump()

                                result = run_agent(
                                    "synthesis-agent",
                                    user_message=(
                                        "**REVISION MODE.** You are revising an existing "
                                        "synthesis based on explicit user feedback from the "
                                        "Gate A reviewer. The user has reviewed the previous "
                                        "draft and given you specific changes to apply. Your "
                                        "output MUST visibly reflect their feedback — if they "
                                        "asked you to remove a question, that question must "
                                        "not appear; if they asked you to add a question, a "
                                        "new question must appear; if they asked you to "
                                        "reword something, the rewording must be applied. "
                                        "Do NOT return a near-identical question list.\n\n"
                                        f"Intake Package:\n"
                                        f"{json.dumps(intake_dict, indent=2)}\n\n"
                                        f"Context Document:\n"
                                        f"{json.dumps(context_dict, indent=2)}\n\n"
                                        f"Current synthesis draft (with any inline question "
                                        f"edits already applied by the user):\n"
                                        f"{json.dumps(draft_dict, indent=2)}\n\n"
                                        f"User feedback to apply:\n{notes.strip()}\n\n"
                                        "Return the FULL revised synthesis (brief, problem "
                                        "statement, win themes, and complete question list). "
                                        "Preserve any of the user's inline question edits "
                                        "unless their feedback explicitly contradicts them."
                                    ),
                                    files=None,
                                    use_extended_thinking=True,
                                    model=ss.get(
                                        _model_state_key("synthesis"),
                                        _MODEL_FULL,
                                    ),
                                )
                                log_agent_result(ss.run_id, "synthesis-agent", result)
                                ss.synthesis_doc = SynthesisDoc.model_validate(result)
                                for _k in [k for k in ss.keys() if k.startswith("q_edit_")]:
                                    del ss[_k]
                                _reset_from("answers")
                                st.rerun()
                            except Exception as exc:
                                traceback.print_exc()
                                log_event(ss.run_id, f"[SYNTHESIS AGENT] RE-RUN ERROR: {exc}")
                                st.error(f"Re-synthesis failed: {exc}")
                    else:
                        st.warning("Add notes before re-running.")

    st.divider()


# ─────────────────────────────────────────────────────────────────────────────
# STEP 4 — PARTNER ANSWERS (voice memo primary, mock answers as test fallback)
# ─────────────────────────────────────────────────────────────────────────────
if not ss.gate_a_done:
    st.text("Step 4 — Partner Answers  [approve Gate A first]")
    st.divider()
else:
    st.subheader("Step 4 — Partner Answers")
    st.markdown(
        "**What this does:** captures the partner's answers to every question "
        "approved at Gate A. Voice path: the partner records one consolidated "
        "memo answering each question in order; whisper-1 transcribes it and "
        "the Voice Splitter Agent maps the transcript onto the question list. "
        "Mock path: an LLM fills plausible answers for testing without a "
        "partner.  \n"
        "**What you do:** record (or upload) the partner's memo and click "
        "**Transcribe & map**, then edit any answer inline before continuing. "
        "Use **Generate Mock Answers** only when no partner is in the loop."
    )

    if not ss.answers_generated:
        col_voice, col_mock = st.columns([3, 2])

        # ── Voice memo (primary) ──
        with col_voice:
            st.markdown("**Voice memo (primary)**")
            st.caption(
                "Record in-browser, or upload a memo recorded on your phone "
                "(mp3, m4a, wav, webm — up to 25 MB)."
            )

            recorded = st.audio_input(
                "Record memo", key="voice_recording_input"
            )
            uploaded_audio = st.file_uploader(
                "Or upload an audio file",
                type=["mp3", "m4a", "wav", "webm", "mp4", "mpga"],
                key="voice_upload_input",
                accept_multiple_files=False,
            )

            picked = uploaded_audio if uploaded_audio is not None else recorded
            if picked is not None:
                size_kb = len(picked.getvalue()) / 1024
                src_label = (
                    "uploaded file" if uploaded_audio is not None else "in-browser recording"
                )
                st.caption(
                    f"Ready to transcribe — {src_label}: "
                    f"{getattr(picked, 'name', 'memo')} ({size_kb:.1f} KB)"
                )

            disabled = picked is None
            if st.button(
                "Transcribe & map to answers",
                type="primary",
                key="run_voice_pipeline",
                disabled=disabled,
            ):
                audio_bytes = picked.getvalue()
                filename = getattr(picked, "name", "memo.wav") or "memo.wav"
                size_kb = len(audio_bytes) / 1024
                log_event(
                    ss.run_id,
                    f"[VOICE] audio received: {filename}, {size_kb:.1f} KB",
                )

                try:
                    n_q = len(ss.synthesis_doc.question_list.questions)
                    log_agent_start(
                        ss.run_id, "voice-splitter-agent",
                        f"transcript → {n_q} question(s)",
                    )
                    with st.spinner(
                        "Transcribing with whisper-1 and routing answers "
                        "(may take 20–60 seconds)..."
                    ):
                        transcript, answer_list = transcribe_and_map(
                            audio_bytes,
                            filename,
                            ss.synthesis_doc.question_list,
                        )

                    log_event(
                        ss.run_id,
                        f"[WHISPER] transcribed {len(transcript)} chars",
                    )
                    log_agent_result(
                        ss.run_id, "voice-splitter-agent",
                        answer_list.model_dump(),
                    )

                    ss.voice_transcript = transcript
                    ss.answer_list = answer_list
                    ss.answers_source = "voice"
                    ss.answers_generated = True
                    st.rerun()
                except ValueError as exc:
                    log_event(ss.run_id, f"[VOICE] ERROR: {exc}")
                    st.error(str(exc))
                except Exception as exc:
                    traceback.print_exc()
                    log_event(ss.run_id, f"[VOICE PIPELINE] ERROR: {exc}")
                    st.error(f"Voice pipeline failed: {exc}")

        # ── Mock answers (test mode) ──
        with col_mock:
            st.markdown("**Mock answers (test mode)**")
            st.caption(
                "For testing only — generates plausible, intentionally uneven "
                "partner answers so the Validation Agent has something to "
                "audit. Will be removed once a real partner is in the loop."
            )
            if st.button("Generate Mock Answers", key="gen_answers"):
                with st.spinner("Mock Partner Agent running..."):
                    try:
                        n_q = len(ss.synthesis_doc.question_list.questions)
                        log_agent_start(
                            ss.run_id, "mock-partner-agent",
                            f"{n_q} question(s) to answer",
                        )
                        ss.answer_list = generate_mock_answers(
                            ss.synthesis_doc.question_list
                        )
                        ss.answers_source = "mock"
                        ss.answers_generated = True
                        log_agent_result(
                            ss.run_id, "mock-partner-agent",
                            ss.answer_list.model_dump(),
                        )
                        st.rerun()
                    except Exception as exc:
                        traceback.print_exc()
                        log_event(ss.run_id, f"[MOCK PARTNER AGENT] ERROR: {exc}")
                        st.error(f"Mock answer generation failed: {exc}")
    else:
        # Banner reflecting source
        if ss.answers_source == "voice":
            st.success(
                "Voice memo transcribed and mapped. Edit any answer below — "
                "questions the memo did not cover are blank for you to fill in."
            )
            if ss.voice_transcript:
                with st.expander("Raw transcript (whisper-1 output)", expanded=False):
                    st.write(ss.voice_transcript)
        elif ss.answers_source == "mock":
            st.info(
                "Mock answers generated. Edit any answer before running validation."
            )

        q_map = {q.id: q for q in ss.synthesis_doc.question_list.questions}

        for answer in ss.answer_list.answers:
            q = q_map.get(answer.question_id)
            hdr = (
                f"**{answer.question_id}** — {q.chapter}"
                if q
                else f"**{answer.question_id}**"
            )
            st.markdown(hdr)
            if q:
                st.caption(q.question)
            st.text_area(
                f"ans_{answer.question_id}",
                value=ss.get(f"answer_edit_{answer.question_id}", answer.answer_text),
                key=f"answer_edit_{answer.question_id}",
                height=90,
                label_visibility="collapsed",
            )

        col_confirm, col_regen = st.columns([2, 2])

        with col_confirm:
            if not ss.answers_done:
                if st.button("Confirm Answers", type="primary", key="confirm_answers"):
                    updated = [
                        Answer(
                            question_id=a.question_id,
                            answer_text=ss.get(
                                f"answer_edit_{a.question_id}", a.answer_text
                            ),
                        )
                        for a in ss.answer_list.answers
                    ]
                    ss.answer_list = AnswerList(answers=updated)
                    ss.answers_done = True
                    ss.validation_done = False
                    ss.validation_report = None
                    st.rerun()
            else:
                st.success("Answers confirmed.")
                if st.button("Edit answers again", key="re_edit_answers"):
                    ss.answers_done = False
                    ss.validation_done = False
                    ss.validation_report = None
                    st.rerun()

        with col_regen:
            if ss.answers_source == "voice":
                if st.button(
                    "Re-record / re-upload memo",
                    key="regen_voice",
                ):
                    ss.answers_generated = False
                    ss.answers_done = False
                    ss.answers_source = None
                    ss.voice_transcript = ""
                    ss.answer_list = None
                    ss.validation_done = False
                    ss.validation_report = None
                    for _k in [
                        k for k in ss.keys()
                        if k.startswith("answer_edit_")
                        or k.startswith("voice_recording_")
                        or k.startswith("voice_upload_")
                    ]:
                        del ss[_k]
                    log_event(ss.run_id, "[VOICE] partner reset memo path")
                    st.rerun()
            else:
                if st.button("Regenerate mock answers", key="regen_answers"):
                    with st.spinner("Regenerating..."):
                        try:
                            ss.answer_list = generate_mock_answers(
                                ss.synthesis_doc.question_list
                            )
                            ss.answers_done = False
                            ss.validation_done = False
                            ss.validation_report = None
                            for _k in [
                                k for k in ss.keys()
                                if k.startswith("answer_edit_")
                            ]:
                                del ss[_k]
                            st.rerun()
                        except Exception as exc:
                            st.error(f"Regeneration failed: {exc}")

    st.divider()


# ─────────────────────────────────────────────────────────────────────────────
# STEP 5 — VALIDATION AGENT
# ─────────────────────────────────────────────────────────────────────────────
if not ss.answers_done:
    st.text("Step 5 — Validation Agent  [confirm answers first]")
else:
    st.subheader("Step 5 — Validation Agent")

    if ss.validation_done and ss.validation_report:
        report: ValidationReport = ss.validation_report

        # Overall metrics row
        m1, m2, m3 = st.columns([1, 1, 4])
        m1.metric("Readiness score", f"{report.readiness_score}/100")
        with m2:
            st.markdown("**Status**")
            if report.overall_readiness == "ready":
                st.success(report.overall_readiness)
            elif report.overall_readiness == "conditional":
                st.warning(report.overall_readiness)
            else:
                st.error(report.overall_readiness)
        m3.info(report.recommendation)

        st.progress(report.readiness_score / 100)
        st.markdown("---")

        # Per-question verdicts
        st.markdown("**Per-question verdicts:**")
        for v in report.verdicts:
            v_cols = st.columns([5, 1])
            with v_cols[0]:
                st.markdown(f"**{v.question_id}** — {v.question_text}")
                st.caption(
                    f"Answer: {v.answer_text[:200]}"
                    + ("..." if len(v.answer_text) > 200 else "")
                )
                st.caption(v.assessment)
                if v.follow_up:
                    st.caption(f"Follow-up needed: {v.follow_up}")
            with v_cols[1]:
                _completeness_badge(v.completeness)

        st.markdown("---")
        # Dot-dash readiness verdict
        rc1, rc2 = st.columns([1, 4])
        with rc1:
            st.markdown("**Dot-dash readiness**")
            if report.can_proceed_to_dot_dash:
                st.success("can proceed")
            else:
                st.error("not yet")
        with rc2:
            if report.dot_dash_blockers:
                label = (
                    "Blockers — must close before dot-dash:"
                    if not report.can_proceed_to_dot_dash
                    else "Material risks — dot-dash will mark these as placeholder/partial:"
                )
                st.markdown(f"**{label}**")
                for b in report.dot_dash_blockers:
                    st.markdown(f"- {b}")
            else:
                st.caption("No material blockers or risks flagged.")

        # Follow-up questions
        if report.follow_up_questions:
            st.markdown("---")
            st.markdown(
                f"**Follow-up questions for the partner — "
                f"{len(report.follow_up_questions)} proposed:**"
            )
            for fq in report.follow_up_questions:
                st.markdown(f"**{fq.id}** ({fq.chapter})  \n{fq.question}")
                st.caption(f"Why asked: {fq.why_asked}")

        # Residual gaps
        if report.residual_gaps:
            st.markdown("---")
            st.markdown("**Residual gaps:**")
            for g in report.residual_gaps:
                st.markdown(f"- {g}")

        if st.button("Re-run Validation Agent", key="rerun_validation"):
            _reset_from("validation")
            st.rerun()

    else:
        st.markdown(
            "**What this does:** audits each Step 4 answer for completeness, "
            "proposes follow-up questions for residual gaps, and gives an "
            "explicit go / no-go verdict on whether the inputs are strong "
            "enough to draft the dot-dash storyline.  \n"
            "**What you do:** click Run, then review the verdict at Gate B "
            "(any 'no-go' sends you back to Step 4 to tighten answers)."
        )
        validation_model = _render_model_selector("validation")
        if st.button("Run Validation Agent", type="primary", key="run_validation"):
            with st.spinner(
                "Validation Agent running — auditing answer completeness "
                "(30–90 seconds)..."
            ):
                try:
                    log_agent_start(
                        ss.run_id, "validation-agent",
                        f"{len(ss.answer_list.answers)} answer(s) to audit",
                    )
                    qs_dict = ss.synthesis_doc.question_list.model_dump()
                    ans_dict = ss.answer_list.model_dump()
                    result = run_agent(
                        "validation-agent",
                        user_message=(
                            "Audit the partner answers against the question list "
                            "and produce a completeness report.\n\n"
                            f"Question List:\n{json.dumps(qs_dict, indent=2)}\n\n"
                            f"Answer List:\n{json.dumps(ans_dict, indent=2)}"
                        ),
                        files=None,
                        use_extended_thinking=True,
                        model=validation_model,
                    )
                    log_agent_result(ss.run_id, "validation-agent", result)
                    ss.validation_report = ValidationReport.model_validate(result)
                    ss.validation_done = True
                    st.rerun()
                except Exception as exc:
                    traceback.print_exc()
                    log_event(ss.run_id, f"[VALIDATION AGENT] ERROR: {exc}")
                    st.error(f"Validation Agent failed: {exc}")

    st.divider()


# ─────────────────────────────────────────────────────────────────────────────
# GATE B — APPROVE TO PROCEED TO DOT-DASH (or loop back with follow-ups)
# ─────────────────────────────────────────────────────────────────────────────
if not ss.validation_done:
    st.text("Gate B — Approve to Dot-Dash  [complete validation first]")
    st.divider()
else:
    st.subheader("Gate B — Approve to Dot-Dash")

    report: ValidationReport = ss.validation_report

    if ss.gate_b_done:
        with st.expander(
            f"Approved — proceeding to dot-dash with "
            f"readiness_score={report.readiness_score}/100",
            expanded=False,
        ):
            if report.dot_dash_blockers:
                st.markdown("**Risks carried into dot-dash:**")
                for b in report.dot_dash_blockers:
                    st.markdown(f"- {b}")
        if st.button("Re-open Gate B", key="reopen_gate_b"):
            ss.gate_b_done = False
            _reset_from("dotdash")
            st.rerun()
    else:
        st.markdown(
            "**What this is:** the human checkpoint that approves the partner "
            "answers as a strong-enough basis for drafting the LoP storyline.  \n"
            "**What you do:** read the validation verdict and per-question "
            "blockers, then pick a path: loop back to Step 4 with follow-up "
            "questions, or proceed to the Dot-Dash Agent."
        )

        # Surface the verdict explicitly at the gate
        if report.can_proceed_to_dot_dash:
            st.success(
                "Validation says: proceed to dot-dash. "
                "Any flagged risks will be carried as `partial` or `placeholder` slides."
            )
        else:
            st.error(
                "Validation says: do NOT proceed yet — at least one critical chapter "
                "is blocked. Recommended: append the proposed follow-up questions to "
                "the question list and re-run mock answers + validation."
            )

        st.markdown("**Recommendation from validation agent:**")
        st.info(report.recommendation)

        st.markdown("---")
        st.markdown(
            "**Two paths from here.** Pick one. (You can override the validation "
            "verdict if you want to push on, but the affected slides will be "
            "marked as placeholders.)"
        )

        col_loop, col_proceed = st.columns([1, 1])

        # ── Path 1: Loop back with follow-ups ──
        with col_loop:
            st.markdown("**Loop back — add follow-ups to question list**")
            n_fu = len(report.follow_up_questions or [])
            if n_fu == 0:
                st.caption("No follow-up questions proposed — nothing to loop back with.")
            else:
                st.caption(
                    f"This will append {n_fu} follow-up question(s) to the question "
                    f"list, clear the current answers and validation, and send you "
                    f"back to Step 4 (Partner Answers) so the partner can answer the "
                    f"follow-ups before re-running validation."
                )
                if st.button(
                    f"Append {n_fu} follow-up(s) and re-run answers",
                    key="loopback_followups",
                ):
                    syn = ss.synthesis_doc
                    existing_qs = list(syn.question_list.questions)
                    # Re-id follow-ups so they continue Q-numbering after existing
                    next_idx = len(existing_qs) + 1
                    appended = []
                    for fq in report.follow_up_questions:
                        appended.append(
                            Question(
                                id=f"Q{next_idx}",
                                chapter=fq.chapter,
                                question=fq.question,
                                why_asked=fq.why_asked
                                + " (follow-up from validation)",
                                expected_answer_type=fq.expected_answer_type,
                            )
                        )
                        next_idx += 1
                    syn.question_list = QuestionList(
                        questions=existing_qs + appended
                    )
                    ss.synthesis_doc = syn
                    log_gate_event(
                        ss.run_id, "gate-b", "LOOPBACK",
                        f"appended {n_fu} follow-up question(s) — "
                        f"total now {len(syn.question_list.questions)}",
                    )
                    # Reset everything from answers onward; keep gate_a approved
                    _reset_from("answers")
                    ss.gate_a_done = True
                    st.rerun()

        # ── Path 2: Proceed to dot-dash ──
        with col_proceed:
            st.markdown("**Proceed — send current state to Dot-Dash Agent**")
            st.caption(
                "Use this when validation says proceed, OR when you want to draft "
                "a placeholder dot-dash now and close the gaps later."
            )
            override = False
            if not report.can_proceed_to_dot_dash:
                override = st.checkbox(
                    "I understand validation says NOT to proceed, and I want to "
                    "draft a placeholder dot-dash anyway.",
                    key="gate_b_override",
                )
            can_press = report.can_proceed_to_dot_dash or override
            if st.button(
                "Approve & proceed to Dot-Dash",
                type="primary",
                disabled=not can_press,
                key="approve_gate_b",
            ):
                ss.gate_b_done = True
                log_gate_event(
                    ss.run_id, "gate-b", "APPROVED",
                    "proceed to dot-dash"
                    + (" (with override)" if override else ""),
                )
                _reset_from("dotdash")
                st.rerun()

    st.divider()


# ─────────────────────────────────────────────────────────────────────────────
# STEP 6 — DOT-DASH AGENT
# ─────────────────────────────────────────────────────────────────────────────
if not ss.gate_b_done:
    st.text("Step 6 — Dot-Dash Agent  [approve Gate B first]")
    st.divider()
else:
    st.subheader("Step 6 — Dot-Dash Agent")

    if ss.dotdash_done and ss.dotdash_doc:
        dd: DotDashDoc = ss.dotdash_doc
        with st.expander(
            f"Done — {len(dd.slides)} slide(s) drafted", expanded=False
        ):
            st.markdown("**Storyline summary**")
            st.info(dd.storyline_summary)
            for s in dd.slides:
                st.markdown(f"**{s.chapter}**  \n{s.headline}")
                for p in s.supporting_points:
                    st.markdown(f"  - {p}")
                if s.notes:
                    st.caption(f"Notes: {s.notes}")
        if st.button("Re-run Dot-Dash Agent", key="rerun_dotdash"):
            _reset_from("dotdash")
            st.rerun()
    else:
        st.markdown(
            "**What this does:** drafts the LoP storyline — one slide per "
            "chapter with a declarative McKinsey-style action title (the "
            "'dot') and 3–5 supporting bullet points (the 'dashes').  \n"
            "**What you do:** click Run, then review and edit the slides at "
            "Gate C before export."
        )
        dotdash_model = _render_model_selector("dotdash")
        if st.button("Run Dot-Dash Agent", type="primary", key="run_dotdash"):
            with st.spinner(
                "Dot-Dash Agent running — drafting LoP storyline (30–90 seconds)..."
            ):
                try:
                    log_agent_start(ss.run_id, "dot-dash-agent")
                    intake_dict = ss.intake_package.model_dump()
                    context_dict = ss.context_doc.model_dump()
                    syn_dict = ss.synthesis_doc.model_dump()
                    ans_dict = ss.answer_list.model_dump()
                    val_dict = ss.validation_report.model_dump()
                    result = run_agent(
                        "dot-dash-agent",
                        user_message=(
                            "Produce the LoP storyline as a chapter-by-chapter "
                            "dot-dash. Use the validation report to set slide "
                            "confidence honestly — flag placeholder content "
                            "rather than fabricating.\n\n"
                            f"Intake Package:\n{json.dumps(intake_dict, indent=2)}\n\n"
                            f"Context Document:\n{json.dumps(context_dict, indent=2)}\n\n"
                            f"Synthesis Document:\n{json.dumps(syn_dict, indent=2)}\n\n"
                            f"Answer List:\n{json.dumps(ans_dict, indent=2)}\n\n"
                            f"Validation Report:\n{json.dumps(val_dict, indent=2)}"
                        ),
                        files=None,
                        use_extended_thinking=True,
                        model=dotdash_model,
                    )
                    log_agent_result(ss.run_id, "dot-dash-agent", result)
                    ss.dotdash_doc = DotDashDoc.model_validate(result)
                    ss.dotdash_done = True
                    st.rerun()
                except Exception as exc:
                    traceback.print_exc()
                    log_event(ss.run_id, f"[DOT-DASH AGENT] ERROR: {exc}")
                    st.error(f"Dot-Dash Agent failed: {exc}")

    st.divider()


# ─────────────────────────────────────────────────────────────────────────────
# GATE C — REVIEW AND ITERATE DOT-DASH
# ─────────────────────────────────────────────────────────────────────────────
if not ss.dotdash_done:
    st.text("Gate C — Review Dot-Dash  [complete dot-dash agent first]")
    st.divider()
else:
    st.subheader("Gate C — Review and Iterate Dot-Dash")

    if ss.gate_c_done:
        dd = ss.dotdash_doc
        with st.expander(
            f"Approved — final dot-dash locked ({len(dd.slides)} slide(s))",
            expanded=False,
        ):
            st.markdown("**Storyline summary**")
            st.info(dd.storyline_summary)
            for s in dd.slides:
                st.markdown(f"**{s.chapter}** [{s.confidence}]  \n{s.headline}")
                for p in s.supporting_points:
                    st.markdown(f"  - {p}")
                if s.notes:
                    st.caption(f"Notes: {s.notes}")
        if st.button("Re-open Gate C", key="reopen_gate_c"):
            ss.gate_c_done = False
            st.rerun()
    else:
        dd = ss.dotdash_doc

        st.markdown(
            "**What this is:** the final human checkpoint before exporting "
            "the LoP. The Dot-Dash Agent has drafted the slide-level "
            "storyline.  \n"
            "**What you do:** read the summary, edit headlines / supporting "
            "points / notes inline, then **Approve** to lock the dot-dash, "
            "or **Iterate** to re-draft with feedback notes."
        )

        st.markdown("**Storyline summary**")
        st.info(dd.storyline_summary)

        if dd.open_risks:
            st.markdown("**Open risks flagged by the agent:**")
            for r in dd.open_risks:
                st.markdown(f"- {r}")

        st.markdown("---")
        st.markdown(
            "**Slides** — edit headline, supporting points (one per line), "
            "and notes inline. Edits are saved when you click Approve."
        )

        for i, s in enumerate(dd.slides):
            st.markdown(f"**Slide {i + 1} — {s.chapter}**")
            st.caption(f"Confidence: {s.confidence}")

            st.text_input(
                f"Headline ({s.chapter})",
                value=ss.get(f"dd_headline_{i}", s.headline),
                key=f"dd_headline_{i}",
            )
            st.text_area(
                f"Supporting points — one per line ({s.chapter})",
                value=ss.get(
                    f"dd_dashes_{i}",
                    "\n".join(s.supporting_points),
                ),
                key=f"dd_dashes_{i}",
                height=120,
            )
            st.text_input(
                f"Notes ({s.chapter})",
                value=ss.get(f"dd_notes_{i}", s.notes or ""),
                key=f"dd_notes_{i}",
            )
            st.markdown("---")

        col_approve_c, col_iter_c = st.columns([1, 2])

        with col_approve_c:
            if st.button("Approve Dot-Dash", type="primary", key="approve_gate_c"):
                updated_slides = []
                for i, s in enumerate(dd.slides):
                    new_dashes = [
                        line.strip().lstrip("-").strip()
                        for line in ss.get(
                            f"dd_dashes_{i}", "\n".join(s.supporting_points)
                        ).split("\n")
                        if line.strip()
                    ]
                    updated_slides.append(
                        DotDashSlide(
                            chapter=s.chapter,
                            headline=ss.get(f"dd_headline_{i}", s.headline),
                            supporting_points=new_dashes,
                            confidence=s.confidence,
                            notes=ss.get(f"dd_notes_{i}", s.notes or ""),
                        )
                    )
                ss.dotdash_doc = DotDashDoc(
                    storyline_summary=dd.storyline_summary,
                    slides=updated_slides,
                    open_risks=dd.open_risks,
                )
                ss.gate_c_done = True
                log_gate_event(
                    ss.run_id, "gate-c", "APPROVED",
                    f"{len(updated_slides)} slide(s) confirmed",
                )
                st.rerun()

        with col_iter_c:
            with st.expander("Iterate dot-dash with notes"):
                dd_notes = st.text_area(
                    "Describe what to change (e.g. 'Sharpen the Why McKinsey "
                    "headline; the Approach slide is too generic — add a dash on "
                    "the integrated bench point')",
                    key="gate_c_notes",
                    height=100,
                )
                if st.button(
                    "Re-run Dot-Dash with These Notes", key="rerun_dotdash_notes"
                ):
                    if dd_notes.strip():
                        with st.spinner("Re-running Dot-Dash Agent..."):
                            try:
                                log_gate_event(
                                    ss.run_id, "gate-c", "ITERATE",
                                    f"notes: {dd_notes[:120]}",
                                )
                                log_agent_start(
                                    ss.run_id, "dot-dash-agent",
                                    "re-run with user notes",
                                )

                                # Capture inline edits before sending to agent
                                edited_slides = []
                                for i, s in enumerate(dd.slides):
                                    new_dashes = [
                                        line.strip().lstrip("-").strip()
                                        for line in ss.get(
                                            f"dd_dashes_{i}",
                                            "\n".join(s.supporting_points),
                                        ).split("\n")
                                        if line.strip()
                                    ]
                                    edited_slides.append(
                                        DotDashSlide(
                                            chapter=s.chapter,
                                            headline=ss.get(
                                                f"dd_headline_{i}", s.headline
                                            ),
                                            supporting_points=new_dashes,
                                            confidence=s.confidence,
                                            notes=ss.get(
                                                f"dd_notes_{i}", s.notes or ""
                                            ),
                                        )
                                    )
                                current_draft = DotDashDoc(
                                    storyline_summary=dd.storyline_summary,
                                    slides=edited_slides,
                                    open_risks=dd.open_risks,
                                )
                                draft_dict = current_draft.model_dump()

                                intake_dict = ss.intake_package.model_dump()
                                context_dict = ss.context_doc.model_dump()
                                syn_dict = ss.synthesis_doc.model_dump()
                                ans_dict = ss.answer_list.model_dump()
                                val_dict = ss.validation_report.model_dump()

                                result = run_agent(
                                    "dot-dash-agent",
                                    user_message=(
                                        "**REVISION MODE.** You are revising an "
                                        "existing dot-dash based on explicit user "
                                        "feedback from the Gate C reviewer. The "
                                        "user has reviewed the current draft and "
                                        "given you specific changes to apply. Your "
                                        "output MUST visibly reflect their feedback "
                                        "— do NOT return a near-identical "
                                        "dot-dash. Preserve any inline edits the "
                                        "user already applied unless their feedback "
                                        "explicitly contradicts them.\n\n"
                                        f"Intake Package:\n"
                                        f"{json.dumps(intake_dict, indent=2)}\n\n"
                                        f"Context Document:\n"
                                        f"{json.dumps(context_dict, indent=2)}\n\n"
                                        f"Synthesis Document:\n"
                                        f"{json.dumps(syn_dict, indent=2)}\n\n"
                                        f"Answer List:\n"
                                        f"{json.dumps(ans_dict, indent=2)}\n\n"
                                        f"Validation Report:\n"
                                        f"{json.dumps(val_dict, indent=2)}\n\n"
                                        f"Current dot-dash draft (with any inline "
                                        f"edits already applied by the user):\n"
                                        f"{json.dumps(draft_dict, indent=2)}\n\n"
                                        f"User feedback to apply:\n"
                                        f"{dd_notes.strip()}\n\n"
                                        "Return the FULL revised dot-dash "
                                        "(storyline_summary, all nine slides in "
                                        "canonical order, and open_risks)."
                                    ),
                                    files=None,
                                    use_extended_thinking=True,
                                    model=ss.get(
                                        _model_state_key("dotdash"),
                                        _MODEL_FULL,
                                    ),
                                )
                                log_agent_result(
                                    ss.run_id, "dot-dash-agent", result
                                )
                                ss.dotdash_doc = DotDashDoc.model_validate(result)
                                # Clear inline edit keys so the new draft is shown
                                for _k in [
                                    k for k in ss.keys()
                                    if k.startswith("dd_headline_")
                                    or k.startswith("dd_dashes_")
                                    or k.startswith("dd_notes_")
                                ]:
                                    del ss[_k]
                                st.rerun()
                            except Exception as exc:
                                traceback.print_exc()
                                log_event(
                                    ss.run_id,
                                    f"[DOT-DASH AGENT] RE-RUN ERROR: {exc}",
                                )
                                st.error(f"Re-run failed: {exc}")
                    else:
                        st.warning("Add notes before re-running.")

    st.divider()


# ─────────────────────────────────────────────────────────────────────────────
# STEP 7 — BA SUPPORT PACK
# ─────────────────────────────────────────────────────────────────────────────
if not ss.gate_c_done:
    st.text("Step 7 — BA Support Pack  [approve Gate C first]")
    st.divider()
else:
    st.subheader("Step 7 — BA Support Pack")

    if ss.ba_support_done and ss.ba_support_pack:
        pack: BASupportPack = ss.ba_support_pack
        with st.expander(
            f"Done — {len(pack.todo_list)} todo(s), "
            f"{len(pack.email_drafts)} email draft(s), "
            f"{len(pack.source_pack)} source item(s)",
            expanded=False,
        ):
            st.markdown("**Summary**")
            st.info(pack.summary)
            tab_t, tab_e, tab_s = st.tabs(
                ["To-do list", "Email drafts", "Source pack"]
            )
            with tab_t:
                for t in pack.todo_list:
                    st.markdown(
                        f"- **[{t.priority}]** ({t.chapter}) {t.action}"
                        + (f" — _due {t.due_relative}_" if t.due_relative else "")
                    )
            with tab_e:
                for e in pack.email_drafts:
                    st.markdown(f"**{e.recipient_name}** — {e.subject}")
            with tab_s:
                for s_item in pack.source_pack:
                    st.markdown(
                        f"- ({s_item.chapter}) {s_item.description}"
                        + (f" — contact: {s_item.contact_name}" if s_item.contact_name else "")
                    )
        if st.button("Re-run BA Support Agent", key="rerun_ba_support"):
            _reset_from("ba_support")
            st.rerun()
    else:
        st.markdown(
            "**What this does:** converts every open item from validation, "
            "the dot-dash, and the partner answers into a concrete BA "
            "handoff pack — a to-do list (priorities + chapters), email "
            "drafts for every partner-named contact, and a chapter-by-"
            "chapter source-pack checklist of materials still to pull.  \n"
            "**What you do:** click Run, then review the three tabs. Edit "
            "email drafts inline before copy-pasting; export as markdown "
            "if you want the pack as a single document."
        )
        ba_support_model = _render_model_selector("ba_support")
        if st.button(
            "Run BA Support Agent", type="primary", key="run_ba_support"
        ):
            with st.spinner(
                "BA Support Agent running — drafting todos, emails, and "
                "source pack (30–60 seconds)..."
            ):
                try:
                    log_agent_start(ss.run_id, "ba-support-agent")
                    intake_dict = ss.intake_package.model_dump()
                    syn_dict = ss.synthesis_doc.model_dump()
                    ans_dict = ss.answer_list.model_dump()
                    val_dict = ss.validation_report.model_dump()
                    dd_dict = ss.dotdash_doc.model_dump()
                    result = run_agent(
                        "ba-support-agent",
                        user_message=(
                            "Produce the BA Support Pack for the approved "
                            "LoP storyline. Convert every blocker, risk, "
                            "residual gap, and partial/placeholder slide "
                            "into concrete BA actions. Draft one email per "
                            "partner-named contact. Build the chapter-by-"
                            "chapter source-pack checklist.\n\n"
                            f"Intake Package:\n{json.dumps(intake_dict, indent=2)}\n\n"
                            f"Synthesis Document:\n{json.dumps(syn_dict, indent=2)}\n\n"
                            f"Answer List:\n{json.dumps(ans_dict, indent=2)}\n\n"
                            f"Validation Report:\n{json.dumps(val_dict, indent=2)}\n\n"
                            f"Approved Dot-Dash:\n{json.dumps(dd_dict, indent=2)}"
                        ),
                        files=None,
                        use_extended_thinking=True,
                        model=ba_support_model,
                    )
                    log_agent_result(ss.run_id, "ba-support-agent", result)
                    ss.ba_support_pack = BASupportPack.model_validate(result)
                    ss.ba_support_done = True
                    st.rerun()
                except Exception as exc:
                    traceback.print_exc()
                    log_event(
                        ss.run_id, f"[BA SUPPORT AGENT] ERROR: {exc}"
                    )
                    st.error(f"BA Support Agent failed: {exc}")

    # Render the three working sections (todos / emails / source pack) in
    # full so the BA can edit, copy, and export — only when the pack is
    # already produced.
    if ss.ba_support_done and ss.ba_support_pack:
        pack = ss.ba_support_pack

        st.markdown("**Summary**")
        st.info(pack.summary)

        st.markdown("---")

        # ── To-do list ──
        st.markdown("### To-do list")
        if pack.todo_list:
            _priority_rank = {"blocker": 0, "high": 1, "normal": 2}
            sorted_todos = sorted(
                pack.todo_list,
                key=lambda t: (
                    _priority_rank.get(t.priority, 3),
                    t.chapter,
                ),
            )
            todo_rows = [
                {
                    "Priority": t.priority,
                    "Chapter": t.chapter,
                    "Owner": t.owner,
                    "Action": t.action,
                    "Dependency": t.dependency,
                    "Due (relative)": t.due_relative,
                    "ID": t.id,
                }
                for t in sorted_todos
            ]
            st.dataframe(
                todo_rows,
                hide_index=True,
                use_container_width=True,
            )
            todo_md_lines = ["# BA To-do list", ""]
            for t in sorted_todos:
                line = (
                    f"- **[{t.priority}]** ({t.chapter}) {t.action}"
                )
                if t.due_relative:
                    line += f" — _due {t.due_relative}_"
                if t.dependency:
                    line += f"  \n  _Dependency:_ {t.dependency}"
                todo_md_lines.append(line)
            st.download_button(
                "Download to-do list as markdown",
                data="\n".join(todo_md_lines).encode("utf-8"),
                file_name=f"todo-list-{ss.run_id}.md",
                mime="text/markdown",
                key="dl_todos",
            )
        else:
            st.caption("No todos produced — all items closed at Gate C.")

        st.markdown("---")

        # ── Email drafts ──
        st.markdown("### Email drafts")
        if pack.email_drafts:
            for i, e in enumerate(pack.email_drafts):
                with st.expander(
                    f"{e.recipient_name} — {e.subject}", expanded=False
                ):
                    st.caption(
                        f"Recipient: {e.recipient_name}"
                        + (f" ({e.recipient_role})" if e.recipient_role else "")
                        + f"  ·  Purpose: {e.purpose}"
                        + (
                            f"  ·  Linked chapter: {e.linked_chapter}"
                            if e.linked_chapter else ""
                        )
                    )
                    st.text_input(
                        "Subject",
                        value=ss.get(f"email_subj_{i}", e.subject),
                        key=f"email_subj_{i}",
                    )
                    st.text_area(
                        "Body",
                        value=ss.get(f"email_body_{i}", e.body),
                        key=f"email_body_{i}",
                        height=240,
                    )
            email_md_lines = ["# BA email drafts", ""]
            for i, e in enumerate(pack.email_drafts):
                subj = ss.get(f"email_subj_{i}", e.subject)
                body = ss.get(f"email_body_{i}", e.body)
                email_md_lines.append(f"## To: {e.recipient_name} — {e.recipient_role}")
                email_md_lines.append(f"**Subject:** {subj}")
                email_md_lines.append("")
                email_md_lines.append(body)
                email_md_lines.append("")
                email_md_lines.append("---")
                email_md_lines.append("")
            st.download_button(
                "Download all email drafts as markdown",
                data="\n".join(email_md_lines).encode("utf-8"),
                file_name=f"email-drafts-{ss.run_id}.md",
                mime="text/markdown",
                key="dl_emails",
            )
        else:
            st.caption(
                "No email drafts — partner did not name specific contacts "
                "in the answers."
            )

        st.markdown("---")

        # ── Source pack ──
        st.markdown("### Source-pack checklist")
        if pack.source_pack:
            source_rows = [
                {
                    "Chapter": s_item.chapter,
                    "Item type": s_item.item_type,
                    "Description": s_item.description,
                    "Contact": s_item.contact_name,
                    "Status": s_item.status,
                    "ID": s_item.id,
                }
                for s_item in pack.source_pack
            ]
            st.dataframe(
                source_rows,
                hide_index=True,
                use_container_width=True,
            )
            source_md_lines = ["# BA source-pack checklist", ""]
            for s_item in pack.source_pack:
                line = (
                    f"- **{s_item.chapter}** — {s_item.description} "
                    f"({s_item.item_type})"
                )
                if s_item.contact_name:
                    line += f" — contact: {s_item.contact_name}"
                line += f" [{s_item.status}]"
                source_md_lines.append(line)
            st.download_button(
                "Download source pack as markdown",
                data="\n".join(source_md_lines).encode("utf-8"),
                file_name=f"source-pack-{ss.run_id}.md",
                mime="text/markdown",
                key="dl_sources",
            )
        else:
            st.caption("No outstanding source items.")

    st.divider()


# ─── SLIDE AUTHOR DISPATCH ───────────────────────────────────────────────────
# Maps each canonical chapter to its specialized authoring agent. The cover
# slide is now ALWAYS prepended as the first deck slide via a separate
# `slide-cover-agent` call (it no longer hijacks slide_index == 0 of the
# dot-dash, so the Context and Objectives chapter gets its own real slide).
# Unknown chapters fall back to the legacy single-agent author so partial
# dot-dash docs still render.
_AGENT_FOR_CHAPTER: dict[str, str] = {
    "Context and Objectives":  "slide-context-agent",
    "Why McKinsey":            "slide-why-mckinsey-agent",
    "Timeline and Team":       "slide-timeline-team-agent",
    "Team":                    "slide-team-agent",
    "Credentials":             "slide-credentials-agent",
    "Market Trends":           "slide-market-trends-agent",
    "Approach":                "slide-approach-agent",
    "Fees":                    "slide-fees-agent",
    "Appendix":                "slide-appendix-agent",
}

# Hard cap on total deck size (cover + chapter slides). Chapter agents
# may emit 1-2 slides; this is a guardrail against unbounded fan-out.
_MAX_DECK_SLIDES = 25


def _agent_for_chapter(chapter: str) -> str:
    """Pick the chapter-specialized slide author agent for one DotDashSlide."""
    return _AGENT_FOR_CHAPTER.get(chapter, "slide-author-agent")


def _ba_pack_items_for_chapter(
    pack: BASupportPack | None,
    chapter: str,
    item_types: tuple[str, ...] | None = None,
) -> list[dict]:
    """
    Return the BA-pack source-pack items relevant to one chapter (excluding
    `received` items, which have already been folded in upstream). If
    `item_types` is provided, only items whose `item_type` is in that tuple
    are returned.
    """
    if not pack or not pack.source_pack:
        return []
    out: list[dict] = []
    for s_item in pack.source_pack:
        if s_item.chapter != chapter:
            continue
        if (s_item.status or "").lower() == "received":
            continue
        if item_types and s_item.item_type not in item_types:
            continue
        out.append(s_item.model_dump())
    return out


def _ba_pack_emails_for_chapter(
    pack: BASupportPack | None,
    chapter: str,
) -> list[dict]:
    """Email drafts whose `linked_chapter` matches this chapter (anchors named experts on the slide)."""
    if not pack or not pack.email_drafts:
        return []
    return [
        e.model_dump() for e in pack.email_drafts if e.linked_chapter == chapter
    ]


def _build_slide_payload(
    *,
    agent_name: str,
    slide: DotDashSlide,
    idx: int,
    intake: IntakePackage | None,
    context: ContextDoc | None,
    synthesis: SynthesisDoc | None,
    dotdash: DotDashDoc,
    ba_pack: BASupportPack | None,
    format_mode: str,
    client_style_summary: str,
) -> str:
    """
    Build the chapter-specific user_message for one slide. Prepends the
    shared style guide and canonical chapter brief, then serialises the
    minimum-viable input slice for the selected agent so each specialist
    sees only what its chapter needs (intake, context dimensions, BA-pack
    items filtered by chapter and item_type).
    """
    intake_essentials = (
        {
            "client_name":        intake.client_name,
            "industry":           intake.industry,
            "geography":          intake.geography,
            "problem_area":       intake.problem_area,
            "pursuit_type":       intake.pursuit_type,
            "competitive_status": intake.competitive_status,
            "competitor_firms":   list(intake.competitor_firms or []),
            "rfp_requirements":   list(intake.rfp_requirements or []),
            "key_facts":          list(intake.key_facts or []),
        }
        if intake is not None
        else {}
    )

    problem_statement = synthesis.problem_statement if synthesis else ""
    win_themes = list(synthesis.win_themes or []) if synthesis else []
    storyline_summary = dotdash.storyline_summary

    # Per-chapter slice of ContextDoc and BA pack — keep payloads tight.
    chapter_slice: dict = {}

    if context is not None:
        chapter_takeaways = context.chapter_takeaways

        if agent_name == "slide-context-agent":
            chapter_slice["relevant_challenges"] = context.relevant_challenges
            chapter_slice["chapter_takeaway"] = (
                chapter_takeaways.context_and_objectives
            )
        elif agent_name == "slide-why-mckinsey-agent":
            chapter_slice["competitive_landscape"] = context.competitive_landscape
            chapter_slice["chapter_takeaway"] = chapter_takeaways.why_mckinsey
        elif agent_name == "slide-market-trends-agent":
            chapter_slice["market_trends"] = context.market_trends
            chapter_slice["recent_signals"] = [
                s.model_dump() for s in context.recent_signals
            ]
            chapter_slice["regulatory_environment"] = [
                r.model_dump() for r in context.regulatory_environment
            ]
            chapter_slice["chapter_takeaway"] = chapter_takeaways.market_trends
            chapter_slice["search_mode"] = context.search_mode
        elif agent_name == "slide-approach-agent":
            chapter_slice["relevant_challenges"] = context.relevant_challenges
            chapter_slice["chapter_takeaway"] = chapter_takeaways.approach
        elif agent_name == "slide-credentials-agent":
            chapter_slice["chapter_takeaway"] = chapter_takeaways.credentials

    if ba_pack is not None:
        if agent_name == "slide-team-agent":
            chapter_slice["ba_pack_cv_items"] = _ba_pack_items_for_chapter(
                ba_pack, "Team", item_types=("cv",)
            )
            chapter_slice["ba_pack_team_emails"] = _ba_pack_emails_for_chapter(
                ba_pack, "Team"
            )
        elif agent_name == "slide-credentials-agent":
            chapter_slice["ba_pack_case_one_pagers"] = (
                _ba_pack_items_for_chapter(
                    ba_pack, "Credentials", item_types=("case_one_pager",)
                )
            )
        elif agent_name == "slide-fees-agent":
            chapter_slice["ba_pack_fee_models"] = _ba_pack_items_for_chapter(
                ba_pack, "Fees", item_types=("fee_model",)
            )
        elif agent_name == "slide-timeline-team-agent":
            chapter_slice["ba_pack_team_items"] = _ba_pack_items_for_chapter(
                ba_pack, "Timeline and Team"
            )
        elif agent_name == "slide-appendix-agent":
            # Appendix takes everything not already received, regardless of
            # chapter — the agent decides what fits the appendix narrative.
            appendix_items: list[dict] = []
            for s_item in (ba_pack.source_pack or []):
                if (s_item.status or "").lower() == "received":
                    continue
                appendix_items.append(s_item.model_dump())
            chapter_slice["ba_pack_appendix_candidates"] = appendix_items

    payload = {
        "slide_index":           idx,
        "format_mode":           format_mode,
        "client_style_summary":  client_style_summary,
        "intake_essentials":     intake_essentials,
        "problem_statement":     problem_statement,
        "win_themes":            win_themes,
        "storyline_summary":     storyline_summary,
        "slide":                 slide.model_dump(),
        "chapter_slice":         chapter_slice,
    }

    parts = [
        render_slide_style_guide(),
        "",
        render_chapter_brief(),
        "",
        "## Authoring task",
        "",
        f"Author the inner `<section class=\"slide\">` block for slide "
        f"{idx + 1} (chapter: {slide.chapter}). Use only the inputs in the "
        "JSON payload below. Return the JSON object "
        "`{ \"slides\": [{ \"html_body\": \"...\", \"notes\": \"\" }, ...], "
        "\"notes\": \"\" }` only — `slides` is an array (1+ fragments).",
        "",
        "## Inputs",
        "",
        "```json",
        json.dumps(payload, indent=2, ensure_ascii=False),
        "```",
    ]
    return "\n".join(parts)


# ─── STEP 9 HELPERS — DECK + UPLOAD EXTRACTION ────────────────────────────────
#
# `_deck_to_text` flattens the in-app rendered deck (`ss.slide_deck`) into a
# plain-text block the Client Evaluator can read like a human. Each slide is
# prefixed with `--- Slide N: <chapter> ---` so the agent can reference the
# right chapter when grading.
#
# `_extract_uploaded_proposal` accepts the optional final-deliverable upload
# (HTML / PDF / PPTX). It returns one of two payload shapes that match the
# orchestrator's existing input contract:
#   - `(text, None)`     for HTML / PPTX — passed inline as part of the
#                        user_message.
#   - `(None, [file])`   for PDF — handed off to `run_agent(..., files=...)`
#                        which forwards it as a native file block via the
#                        OpenAI / Anthropic file content type.

def _strip_html_to_text(html: str) -> str:
    """Lightweight HTML-to-text — collapses tags, decodes entities, keeps
    visible text. Good enough for evaluator input; not a full DOM parser."""
    if not html:
        return ""

    from html.parser import HTMLParser
    from html import unescape

    class _TextExtractor(HTMLParser):
        def __init__(self) -> None:
            super().__init__()
            self.parts: list[str] = []
            self._skip_depth = 0

        def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
            if tag in {"script", "style"}:
                self._skip_depth += 1
            if tag in {"br", "p", "li", "div", "section", "h1", "h2", "h3", "h4", "h5", "h6"}:
                self.parts.append("\n")

        def handle_endtag(self, tag: str) -> None:
            if tag in {"script", "style"} and self._skip_depth > 0:
                self._skip_depth -= 1
            if tag in {"p", "li", "div", "section", "h1", "h2", "h3", "h4", "h5", "h6"}:
                self.parts.append("\n")

        def handle_data(self, data: str) -> None:
            if self._skip_depth == 0:
                self.parts.append(data)

    parser = _TextExtractor()
    try:
        parser.feed(html)
    except Exception:
        return unescape(html)

    text = unescape("".join(parser.parts))
    lines = [ln.strip() for ln in text.splitlines()]
    return "\n".join(ln for ln in lines if ln)


def _deck_to_text(deck: SlideDeck | None) -> str:
    """Flatten an in-app rendered deck into a readable text block — one
    `--- Slide N: <chapter> ---` header per slide. Returns empty string
    when the deck is missing or empty."""
    if deck is None or not deck.slides:
        return ""
    blocks: list[str] = []
    for i, sl in enumerate(deck.slides, start=1):
        header = f"--- Slide {i}: {sl.chapter or 'Untitled'} ---"
        body = _strip_html_to_text(sl.html or "")
        blocks.append(f"{header}\n{body}".strip())
    return "\n\n".join(blocks)


def _extract_pptx_text(pptx_bytes: bytes) -> str:
    """Extract plain text from a PPTX upload — one slide per block, with a
    `--- Slide N ---` header so the agent can reference it. Returns empty
    string on parse failure."""
    if not pptx_bytes:
        return ""
    try:
        from pptx import Presentation
        prs = Presentation(BytesIO(pptx_bytes))
    except Exception:
        return ""

    blocks: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            text = ""
            if shape.has_text_frame:
                text = "\n".join(
                    p.text for p in shape.text_frame.paragraphs if p.text
                )
            elif getattr(shape, "has_table", False) and shape.has_table:
                rows: list[str] = []
                for row in shape.table.rows:
                    cells = [
                        (cell.text or "").strip().replace("\n", " ")
                        for cell in row.cells
                    ]
                    rows.append(" | ".join(cells))
                text = "\n".join(rows)
            if text and text.strip():
                lines.append(text.strip())
        blocks.append(f"--- Slide {i} ---\n" + "\n".join(lines))
    return "\n\n".join(blocks)


def _extract_uploaded_proposal(
    upload,
) -> tuple[str, list[dict] | None]:
    """
    Convert a Streamlit `UploadedFile` (HTML / PDF / PPTX) into the
    evaluator input contract.

    Returns
    -------
    (proposal_text, files_for_orchestrator)
        - proposal_text is non-empty for HTML / PPTX; empty for PDF
          (the PDF is passed natively via the files list instead).
        - files_for_orchestrator is `None` for HTML / PPTX; a one-element
          list of the standard `{filename, doc_type, is_pdf, bytes, text}`
          dict for PDF (matches the shape `orchestrator._build_content`
          consumes).
    """
    if upload is None:
        return "", None

    name = (upload.name or "").lower()
    raw = upload.getvalue()

    if name.endswith(".pdf"):
        return "", [
            {
                "filename": upload.name,
                "doc_type": "Final Proposal",
                "is_pdf": True,
                "bytes": raw,
                "text": "",
            }
        ]
    if name.endswith(".pptx"):
        return _extract_pptx_text(raw), None
    # Treat everything else as HTML / text.
    try:
        text = raw.decode("utf-8", errors="replace")
    except Exception:
        text = ""
    if "<" in text and ">" in text:
        text = _strip_html_to_text(text)
    return text, None


# ─────────────────────────────────────────────────────────────────────────────
# STEP 8 — RENDER SLIDES
# ─────────────────────────────────────────────────────────────────────────────
if not ss.ba_support_done:
    st.text("Step 8 — Render Slides  [complete BA Support Pack first]")
    st.divider()
else:
    st.subheader("Step 8 — Render Slides")

    if ss.slides_done and ss.slide_deck:
        deck: SlideDeck = ss.slide_deck
        with st.expander(
            f"Done — {len(deck.slides)} slide(s) rendered "
            f"({deck.format_mode} format)",
            expanded=False,
        ):
            st.caption(
                f"Format: {deck.format_mode}"
                + (
                    f"  ·  {deck.client_style_summary}"
                    if deck.client_style_summary else ""
                )
            )
            for i, sl in enumerate(deck.slides):
                st.markdown(f"- Slide {i + 1} — {sl.chapter} ({sl.filename})")
        if st.button("Re-render slides", key="rerun_slides"):
            _reset_from("slides")
            st.rerun()
    else:
        st.markdown(
            "**What this does:** renders the approved dot-dash as one "
            "1920×1080 HTML slide per chapter, using the bundled McKinsey "
            "stylesheet by default or a CSS overlay derived from a partner-"
            "supplied client PPTX style.  \n"
            "**What you do:** pick the format, optionally upload the "
            "client house-style PPTX, then click Run. The deck preview "
            "renders inline; you can also download the full deck (HTML "
            "files + CSS) as a ZIP."
        )

        format_choice = st.radio(
            "Slide format",
            options=[
                "McKinsey format (default)",
                "Client format",
            ],
            index=0 if ss.slide_format_mode == "mckinsey" else 1,
            key="slide_format_radio",
            horizontal=True,
        )
        ss.slide_format_mode = (
            "mckinsey"
            if format_choice.startswith("McKinsey")
            else "client"
        )

        client_style: dict = {}
        client_style_summary = ""
        if ss.slide_format_mode == "client":
            up = st.file_uploader(
                "Upload client house-style PPTX (optional)",
                type=["pptx"],
                key="client_style_uploader",
            )
            descr = st.text_area(
                "Or describe the client style in plain text "
                "(used as guidance when no PPTX is supplied — and as a "
                "supplementary hint when one is)",
                value=ss.client_style_summary,
                height=100,
                key="client_style_text",
            )
            if up is not None:
                pptx_bytes = up.getvalue()
                ss.client_pptx_bytes = pptx_bytes
                client_style = parse_client_pptx_style(pptx_bytes) or {}
                client_style_summary = (
                    client_style.get("summary") or ""
                )
                if not client_style_summary and descr.strip():
                    client_style_summary = descr.strip()
                if client_style_summary:
                    st.success(
                        "Client style parsed — " + client_style_summary
                    )
                else:
                    st.warning(
                        "Could not extract a usable style fingerprint "
                        "from the upload. Falling back to the description "
                        "above (or McKinsey defaults if empty)."
                    )
            elif descr.strip():
                client_style_summary = descr.strip()

            ss.client_style_summary = client_style_summary

        slides_model = _render_model_selector("slides")
        if st.button("Render Slides", type="primary", key="run_slides"):
            try:
                log_agent_start(
                    ss.run_id,
                    "slide-author-agent",
                    f"format={ss.slide_format_mode} | model={slides_model}",
                )

                # Prepare the export folder for this run (single .html file)
                export_dir = (
                    Path(__file__).resolve().parent.parent
                    / "exports"
                    / str(ss.run_id)
                )
                export_dir.mkdir(parents=True, exist_ok=True)

                # Build the optional client-style overlay as a CSS string;
                # it gets inlined into the single deck HTML at write time.
                client_css_text = ""
                if (
                    ss.slide_format_mode == "client"
                    and ss.client_pptx_bytes
                ):
                    style_dict = parse_client_pptx_style(
                        ss.client_pptx_bytes
                    ) or {}
                    if not style_dict.get("summary") and ss.client_style_summary:
                        style_dict["summary"] = ss.client_style_summary
                    client_css_text = build_client_css_text(style_dict)
                    client_summary_for_agent = (
                        style_dict.get("summary") or ss.client_style_summary
                    )
                else:
                    client_summary_for_agent = (
                        ss.client_style_summary
                        if ss.slide_format_mode == "client"
                        else ""
                    )

                # Author one slide at a time and assemble the deck.
                # Order: prepended cover (always 1) + each dot-dash chapter
                # (each chapter agent may emit 1+ slides).
                rendered_slides: list[SlideHTML] = []
                dd: DotDashDoc = ss.dotdash_doc

                # Synthetic cover slide so the cover agent's input shape is
                # the same as a chapter agent's. The cover's title text is
                # the synthesis problem_statement (per slide-cover-agent.md);
                # we pass the storyline_summary in `notes` for context only.
                cover_slide_input = DotDashSlide(
                    chapter="Cover",
                    headline=ss.synthesis_doc.problem_statement or "",
                    supporting_points=[],
                    confidence="complete",
                    notes=dd.storyline_summary or "",
                )

                run_plan: list[tuple[str, DotDashSlide]] = [
                    ("slide-cover-agent", cover_slide_input),
                ]
                for slide in dd.slides:
                    run_plan.append((_agent_for_chapter(slide.chapter), slide))

                progress = st.progress(0.0, text="Rendering slides…")
                total_calls = len(run_plan)
                for call_idx, (agent_name, slide_input) in enumerate(run_plan):
                    deck_idx = len(rendered_slides)
                    if deck_idx >= _MAX_DECK_SLIDES:
                        log_event(
                            ss.run_id,
                            f"[SLIDE BUDGET] truncated at "
                            f"{_MAX_DECK_SLIDES} slides — skipping "
                            f"remaining chapter agents",
                        )
                        break

                    progress.progress(
                        call_idx / max(total_calls, 1),
                        text=(
                            f"Rendering call {call_idx + 1}/{total_calls} — "
                            f"{slide_input.chapter} ({agent_name})"
                        ),
                    )
                    user_msg = _build_slide_payload(
                        agent_name=agent_name,
                        slide=slide_input,
                        idx=deck_idx,
                        intake=ss.intake_package,
                        context=ss.context_doc,
                        synthesis=ss.synthesis_doc,
                        dotdash=dd,
                        ba_pack=ss.ba_support_pack,
                        format_mode=ss.slide_format_mode,
                        client_style_summary=client_summary_for_agent,
                    )
                    log_event(
                        ss.run_id,
                        f"[SLIDE AUTHOR] call {call_idx + 1}/{total_calls} "
                        f"({slide_input.chapter}) -> {agent_name}",
                    )
                    result = run_agent(
                        agent_name,
                        user_message=user_msg,
                        files=None,
                        use_extended_thinking=False,
                        model=slides_model,
                    )

                    # Parse the agent output via ChapterSlideAuthorOutput,
                    # which accepts both the new multi-slide shape
                    # `{slides:[{html_body, notes}]}` and the legacy
                    # `{html_body, notes}` shape.
                    try:
                        parsed = ChapterSlideAuthorOutput.model_validate(
                            result or {}
                        )
                    except Exception as exc:
                        log_event(
                            ss.run_id,
                            f"[SLIDE AUTHOR] parse error for "
                            f"{slide_input.chapter} ({agent_name}): {exc}",
                        )
                        parsed = ChapterSlideAuthorOutput(slides=[], notes="")

                    chapter_label = slide_input.chapter or f"Slide {deck_idx + 1}"
                    for sub_idx, fragment in enumerate(parsed.slides):
                        if len(rendered_slides) >= _MAX_DECK_SLIDES:
                            break
                        # Multi-slide chapters get a "(2)", "(3)", … suffix
                        # in the on-card label so the BA can tell them apart.
                        label = (
                            chapter_label
                            if sub_idx == 0
                            else f"{chapter_label} ({sub_idx + 1})"
                        )
                        slug = (
                            chapter_label.lower()
                            .replace("&", "and")
                            .replace("/", "-")
                        )
                        slug = "".join(
                            c if c.isalnum() or c == "-" else "-"
                            for c in slug
                        ).strip("-") or f"slide-{len(rendered_slides) + 1}"
                        suffix = "" if sub_idx == 0 else f"-{sub_idx + 1}"
                        rendered_slides.append(
                            SlideHTML(
                                chapter=label,
                                filename=(
                                    f"{len(rendered_slides) + 1:02d}-"
                                    f"{slug}{suffix}.html"
                                ),
                                html=fragment.html_body or "",
                                notes=fragment.notes or "",
                            )
                        )

                progress.progress(1.0, text="Writing deck file…")

                deck = SlideDeck(
                    format_mode=ss.slide_format_mode,
                    client_style_summary=(
                        client_summary_for_agent
                        if ss.slide_format_mode == "client"
                        else ""
                    ),
                    slides=rendered_slides,
                )
                index_path = write_single_deck(
                    export_dir,
                    rendered_slides,
                    ss.slide_format_mode,
                    client_css_text,
                )

                log_agent_result(
                    ss.run_id,
                    "slide-author-agent",
                    {
                        "slides_rendered": len(rendered_slides),
                        "index_path": str(index_path),
                        "format_mode": ss.slide_format_mode,
                    },
                )

                ss.slide_deck = deck
                ss.slides_done = True
                progress.empty()
                st.rerun()
            except Exception as exc:
                traceback.print_exc()
                log_event(
                    ss.run_id, f"[SLIDE AUTHOR AGENT] ERROR: {exc}"
                )
                st.error(f"Slide rendering failed: {exc}")

    # Done view: in-app preview + single-file HTML download
    if ss.slides_done and ss.slide_deck:
        deck = ss.slide_deck

        st.caption(
            f"Format: {deck.format_mode}"
            + (
                f"  ·  {deck.client_style_summary}"
                if deck.client_style_summary else ""
            )
        )

        export_dir = (
            Path(__file__).resolve().parent.parent
            / "exports"
            / str(ss.run_id)
        )
        index_path = export_dir / "index.html"
        if index_path.exists():
            try:
                st_components.html(
                    index_path.read_text(encoding="utf-8"),
                    height=900,
                    scrolling=True,
                )
            except Exception as exc:
                st.warning(f"Inline preview failed: {exc}")
                st.caption(
                    f"You can open the deck directly at: "
                    f"{index_path}"
                )

            try:
                st.download_button(
                    "Download deck (single self-contained HTML)",
                    data=index_path.read_bytes(),
                    file_name=f"lop-deck-{ss.run_id}.html",
                    mime="text/html",
                    key="dl_deck_html",
                )
            except Exception as exc:
                st.warning(f"HTML download packaging failed: {exc}")
        else:
            st.warning(
                "Deck index.html not found on disk — re-render to "
                "regenerate the preview."
            )

    st.divider()


# ─────────────────────────────────────────────────────────────────────────────
# STEP 9 — CLIENT / OWNER PROPOSAL EVALUATOR
# ─────────────────────────────────────────────────────────────────────────────
#
# Runs after the slide deck is rendered (or against an uploaded final
# deliverable). Plays the role of the company owner / sponsor on the
# client side and produces a structured `ClientEvaluationReport`:
# RFP coverage, owner priorities (incl. inferred), per-chapter buyer
# view, reasonableness checks on timeline / fees / team / approach, top
# concerns, missing items, recommended changes, and an overall verdict.
# A second evaluator (persona TBD) will plug in alongside this one in
# the same step when defined.

def _consistency_pass_badge(passed: bool) -> None:
    if passed:
        st.success("Consistency QC: PASS")
    else:
        st.error("Consistency QC: FAIL")


def _consistency_severity_marker(severity: str) -> str:
    return {
        "critical": "🔴",
        "major": "🟠",
        "minor": "🟡",
    }.get(severity, "⚪")


def _render_evals_summary_strip(bundle: EvalsBundle) -> None:
    ce = bundle.consistency_eval or {}
    if ce:
        passed = ce.get("passed")
        score = ce.get("overall_score", "—")
        label = "Pass" if passed else "Needs fixes" if passed is False else "—"
        st.metric("Consistency score", f"{score} / 100", delta=label, delta_color="normal")
        summary = ce.get("summary") or {}
        if summary:
            st.caption(
                f"Critical: {summary.get('critical', 0)} · "
                f"Major: {summary.get('major', 0)} · "
                f"Minor: {summary.get('minor', 0)}"
            )
    else:
        st.caption("Consistency evaluation: not run yet")


def _run_consistency_eval(
    *,
    deck_text: str,
    proposal_upload_text: str,
    uploaded_filename: str,
    use_llm_judge: bool,
) -> None:
    from lop_eval.adapters.from_streamlit import (
        proposal_from_deck,
        proposal_from_text,
        source_of_truth_from_intake,
    )
    from lop_eval.evaluator import evaluate_document
    from lop_eval.llm_adapter import LlmJudgeAdapter

    doc_id = ss.run_id
    if ss.slide_deck and ss.client_eval_source == "in_app_deck":
        doc = proposal_from_deck(ss.slide_deck, doc_id)
    elif proposal_upload_text:
        doc = proposal_from_text(proposal_upload_text, doc_id)
    elif deck_text:
        doc = proposal_from_text(deck_text, doc_id)
    else:
        raise ValueError("No extractable proposal text for consistency QC.")

    sot = source_of_truth_from_intake(ss.intake_package)
    llm_adapter = None
    if use_llm_judge:
        import openai

        def _call_llm(prompt: str) -> str:
            client = openai.OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))
            resp = client.chat.completions.create(
                model=_MODEL_FAST,
                messages=[
                    {
                        "role": "system",
                        "content": (
                            "You are a strict evaluator. Reply only with "
                            "Score: 0 or 1 and Reasoning: as specified."
                        ),
                    },
                    {"role": "user", "content": prompt},
                ],
                temperature=0,
            )
            return (resp.choices[0].message.content or "").strip()

        llm_adapter = LlmJudgeAdapter(_call_llm)

    cfg = _build_eval_config_from_session()
    result = evaluate_document(
        doc, source_of_truth=sot, llm_adapter=llm_adapter, config=cfg
    )
    ss.consistency_eval_result = result
    ss.consistency_eval_done = True
    _sync_evals_bundle()


def _client_eval_verdict_badge(verdict: str) -> None:
    """Render the overall verdict as a coloured Streamlit pill."""
    label = verdict.replace("_", " ")
    if verdict == "would_buy":
        st.success(f"Verdict: {label}")
    elif verdict == "would_buy_with_revisions":
        st.warning(f"Verdict: {label}")
    else:
        st.error(f"Verdict: {label}")


def _client_eval_status_marker(status: str) -> str:
    """Single-character coloured marker for table-row status columns."""
    return {
        "covered": "🟢",
        "partial": "🟡",
        "missing": "🔴",
        "strong": "🟢",
        "acceptable": "🟡",
        "weak": "🟠",
        "reasonable": "🟢",
        "stretch": "🟡",
        "unreasonable": "🔴",
        "not_shown": "⚪",
    }.get(status, "⚪")


def _build_client_eval_payload(
    *,
    intake: IntakePackage | None,
    synthesis: SynthesisDoc | None,
    validation: ValidationReport | None,
    dotdash: DotDashDoc | None,
    deck_text: str,
    proposal_upload_text: str,
    proposal_upload_filename: str,
) -> str:
    """
    Build the user_message for `client-evaluator-agent`.

    Always feeds the upstream JSON pack (intake / synthesis / validation /
    dot-dash) so the agent knows what was meant to be in the proposal.
    Then attaches the actual proposal content — either the in-app deck
    text (`deck_text`) or the text extracted from the uploaded final
    deliverable. PDFs are NOT passed through this function; they are
    handed to the orchestrator via `files=` on `run_agent`.
    """
    intake_essentials = (
        {
            "client_name":        intake.client_name,
            "industry":           intake.industry,
            "geography":          intake.geography,
            "problem_area":       intake.problem_area,
            "pursuit_type":       intake.pursuit_type,
            "competitive_status": intake.competitive_status,
            "competitor_firms":   list(intake.competitor_firms or []),
            "rfp_requirements":   list(intake.rfp_requirements or []),
            "key_facts":          list(intake.key_facts or []),
        }
        if intake is not None
        else {}
    )

    synthesis_essentials = (
        {
            "problem_statement": synthesis.problem_statement,
            "win_themes":        list(synthesis.win_themes or []),
        }
        if synthesis is not None
        else {}
    )

    validation_essentials = (
        {
            "overall_readiness":       validation.overall_readiness,
            "readiness_score":         validation.readiness_score,
            "can_proceed_to_dot_dash": validation.can_proceed_to_dot_dash,
            "residual_gaps":           list(validation.residual_gaps or []),
            "recommendation":          validation.recommendation,
        }
        if validation is not None
        else {}
    )

    dotdash_essentials = (
        {
            "storyline_summary": dotdash.storyline_summary,
            "slides": [
                {
                    "chapter":           s.chapter,
                    "headline":          s.headline,
                    "supporting_points": list(s.supporting_points or []),
                    "confidence":        s.confidence,
                    "notes":             s.notes,
                }
                for s in dotdash.slides
            ],
            "open_risks": list(dotdash.open_risks or []),
        }
        if dotdash is not None
        else {}
    )

    payload = {
        "intake_essentials":   intake_essentials,
        "synthesis":           synthesis_essentials,
        "validation":          validation_essentials,
        "dotdash":             dotdash_essentials,
        "deck_text":           deck_text,
        "proposal_upload": {
            "filename": proposal_upload_filename,
            "text":     proposal_upload_text,
        },
    }

    parts = [
        render_chapter_brief(),
        "",
        "## Evaluation task",
        "",
        "Play the role of the company owner / business sponsor receiving "
        "this LoP. Read the proposal end-to-end (either `deck_text` or "
        "`proposal_upload` — if a PDF was uploaded, it is attached as a "
        "native file block alongside this message and should be treated "
        "as the authoritative proposal). Grade it from the buyer's "
        "perspective against the upstream intake / synthesis / "
        "validation / dot-dash pack. Return the JSON object specified "
        "in the agent spec — RFP coverage, owner priorities, chapter "
        "assessment, reasonableness checks, top concerns, missing items, "
        "recommended changes, and an overall verdict with a 0–100 score.",
        "",
        "## Inputs",
        "",
        "```json",
        json.dumps(payload, indent=2, ensure_ascii=False),
        "```",
    ]
    return "\n".join(parts)


# ─── STEP 9 — LOSS ANALYSIS HELPERS ──────────────────────────────────────────

def _loss_eval_likelihood_badge(likelihood: str) -> None:
    """Render the loss likelihood as a coloured Streamlit pill."""
    label = (likelihood or "").replace("_", " ") or "unknown"
    if likelihood == "low":
        st.success(f"Loss likelihood: {label}")
    elif likelihood in ("moderate", "high"):
        st.warning(f"Loss likelihood: {label}")
    elif likelihood == "very_high":
        st.error(f"Loss likelihood: {label}")
    else:
        st.info(f"Loss likelihood: {label}")


def _loss_eval_verdict_badge(verdict: str) -> None:
    """Render the save-or-kill verdict as a coloured Streamlit pill."""
    label = (verdict or "").replace("_", " ") or "unknown"
    if verdict == "competitive_as_is":
        st.success(f"Save-or-kill: {label}")
    elif verdict == "needs_surgical_edits":
        st.warning(f"Save-or-kill: {label}")
    elif verdict == "needs_redo":
        st.error(f"Save-or-kill: {label}")
    else:
        st.info(f"Save-or-kill: {label}")


def _loss_eval_severity_marker(severity: str) -> str:
    """Single-character coloured marker for severity columns."""
    return {
        "critical": "🔴",
        "high":     "🟠",
        "medium":   "🟡",
        "low":      "🟢",
    }.get(severity, "⚪")


def _loss_eval_priority_marker(priority: str) -> str:
    """Single-character coloured marker for improvement priority."""
    return {
        "blocker": "🔴",
        "high":    "🟠",
        "normal":  "🟡",
    }.get(priority, "⚪")


def _loss_eval_framing_preview(
    intake: IntakePackage | None,
    override: str,
) -> tuple[str, str]:
    """
    Compute the framing question the loss agent will ask and a short
    deterministic context hint for the agent payload.

    Returns
    -------
    (display_question, framing_hint)
        - `display_question` is the markdown line shown to the user above
          the Run button so the partner sees exactly what is being asked.
        - `framing_hint` is the pre-computed reading of the competitive
          context that goes into the JSON payload as informational
          context (the agent still emits `framing_question` itself).
    """
    o = (override or "").strip()
    if o:
        return (
            f"Question I will ask: **Why would we lose this proposal "
            f"to {o}?** (partner-supplied stress-test target)",
            f"competitor_override provided: {o}",
        )

    if intake is not None:
        comp_status = (intake.competitive_status or "").strip().lower()
        firms = [
            f.strip() for f in (intake.competitor_firms or []) if f.strip()
        ]
        if comp_status == "competitive" and firms:
            named = ", ".join(firms[:3])
            extra = "" if len(firms) <= 3 else f" (and {len(firms) - 3} more)"
            return (
                f"Question I will ask: **Why would we lose this proposal "
                f"to {named}?**{extra}",
                f"competitive vs named firms: {', '.join(firms)}",
            )
        if comp_status == "non_competitive":
            return (
                "Question I will ask: **Why would we lose this proposal?** "
                "(non-competitive — primary risk is loss to no-decision)",
                "non_competitive — consider loss to no-decision",
            )

    return (
        "Question I will ask: **Why would we lose this proposal?** "
        "(no named competitors — covers loss to no-decision and to "
        "unnamed rivals)",
        "competitive status unclear or no firms named — include "
        "no_decision and any inferred rivals",
    )


def _build_loss_eval_payload(
    *,
    intake: IntakePackage | None,
    synthesis: SynthesisDoc | None,
    validation: ValidationReport | None,
    dotdash: DotDashDoc | None,
    deck_text: str,
    proposal_upload_text: str,
    proposal_upload_filename: str,
    competitor_override: str,
    framing_hint: str,
) -> str:
    """
    Build the user_message for `loss-analysis-agent`. Mirrors
    `_build_client_eval_payload` but adds `competitor_override` and
    `framing_hint` to the payload so the agent can frame the single
    question correctly. PDFs are handed to `run_agent(..., files=...)`
    just like in the client tab — not through this function.
    """
    intake_essentials = (
        {
            "client_name":        intake.client_name,
            "industry":           intake.industry,
            "geography":          intake.geography,
            "problem_area":       intake.problem_area,
            "pursuit_type":       intake.pursuit_type,
            "competitive_status": intake.competitive_status,
            "competitor_firms":   list(intake.competitor_firms or []),
            "rfp_requirements":   list(intake.rfp_requirements or []),
            "key_facts":          list(intake.key_facts or []),
        }
        if intake is not None
        else {}
    )

    synthesis_essentials = (
        {
            "problem_statement": synthesis.problem_statement,
            "win_themes":        list(synthesis.win_themes or []),
        }
        if synthesis is not None
        else {}
    )

    validation_essentials = (
        {
            "overall_readiness":       validation.overall_readiness,
            "readiness_score":         validation.readiness_score,
            "can_proceed_to_dot_dash": validation.can_proceed_to_dot_dash,
            "residual_gaps":           list(validation.residual_gaps or []),
            "recommendation":          validation.recommendation,
        }
        if validation is not None
        else {}
    )

    dotdash_essentials = (
        {
            "storyline_summary": dotdash.storyline_summary,
            "slides": [
                {
                    "chapter":           s.chapter,
                    "headline":          s.headline,
                    "supporting_points": list(s.supporting_points or []),
                    "confidence":        s.confidence,
                    "notes":             s.notes,
                }
                for s in dotdash.slides
            ],
            "open_risks": list(dotdash.open_risks or []),
        }
        if dotdash is not None
        else {}
    )

    payload = {
        "intake_essentials":   intake_essentials,
        "synthesis":           synthesis_essentials,
        "validation":          validation_essentials,
        "dotdash":             dotdash_essentials,
        "deck_text":           deck_text,
        "proposal_upload": {
            "filename": proposal_upload_filename,
            "text":     proposal_upload_text,
        },
        "competitor_override": competitor_override or "",
        "framing_hint":        framing_hint or "",
    }

    parts = [
        render_chapter_brief(),
        "",
        "## Loss-analysis task",
        "",
        "Act as the McKinsey team's red team. Read the proposal "
        "end-to-end (either `deck_text` or `proposal_upload` — if a PDF "
        "was uploaded, it is attached as a native file block alongside "
        "this message and should be treated as the authoritative "
        "proposal). Answer one question, framed using "
        "`competitor_override` (if non-empty) or "
        "`intake_essentials.competitive_status` + "
        "`intake_essentials.competitor_firms`: \"Why would we lose this "
        "proposal\" or \"...to <competitor>\". Return the JSON object "
        "specified in the agent spec — top loss reasons, competitor "
        "angles, vulnerable chapters, loss likelihood + 0–100 risk "
        "score, save-or-kill verdict, punchline, and ranked key "
        "improvements.",
        "",
        "## Inputs",
        "",
        "```json",
        json.dumps(payload, indent=2, ensure_ascii=False),
        "```",
    ]
    return "\n".join(parts)


if not ss.slides_done:
    st.text("Step 9 — Consistency evaluation  [complete Render Slides first]")
    if ss.get("evals_scroll_pending") or _evals_qp_means_scroll():
        st.info(
            "Consistency eval link active. Finish **Step 8 — Render Slides** first — "
            "then reload with `?evals=1` or use **Jump to Step 9 — Consistency eval** in the sidebar."
        )
    st.divider()
else:
    st.markdown(
        '<div id="evals-section-anchor" style="scroll-margin-top:3rem"></div>',
        unsafe_allow_html=True,
    )
    st.subheader("Step 9 — Consistency evaluation")
    if ss.evals_scroll_pending:
        ss.evals_scroll_pending = False
        st_components.html(
            """
<script>
(function () {
  function scrollToEvals() {
    var doc = window.parent.document;
    var el = doc.getElementById("evals-section-anchor");
    if (el) {
      el.scrollIntoView({behavior: "smooth", block: "start"});
      return;
    }
    var hs = doc.querySelectorAll("h2, h3");
    for (var i = 0; i < hs.length; i++) {
      var t = hs[i].innerText || "";
      if (t.indexOf("Step 9") !== -1) {
        hs[i].scrollIntoView({behavior: "smooth", block: "start"});
        break;
      }
    }
  }
  setTimeout(scrollToEvals, 400);
})();
</script>
""",
            height=0,
            width=0,
        )
    st.caption(
        "Internal consistency QC on the proposal (numbers, dates, scope, "
        "terminology, units, claim–support). Pick one source below — the "
        "consistency check grades that document."
    )

    # ─── Shared source picker + uploader ────────────────────────────────
    source_choice = st.radio(
        "What should the consistency check read?",
        options=[
            "Rendered deck from Step 8",
            "Uploaded final deliverable (HTML / PDF / PPTX)",
        ],
        index=0 if ss.client_eval_source == "in_app_deck" else 1,
        key="client_eval_source_radio",
        horizontal=False,
    )
    ss.client_eval_source = (
        "in_app_deck"
        if source_choice.startswith("Rendered")
        else "uploaded_file"
    )

    uploaded_proposal = None
    if ss.client_eval_source == "uploaded_file":
        uploaded_proposal = st.file_uploader(
            "Upload the final proposal (HTML, PDF, or PPTX)",
            type=["html", "htm", "pdf", "pptx"],
            key="client_eval_upload",
        )
        if uploaded_proposal is not None:
            ss.client_eval_upload_name = uploaded_proposal.name

    # Pre-compute proposal content once per render (consistency + optional LLM reviews).
    # `_extract_uploaded_proposal` returns text for HTML/PPTX and native path for PDF.
    shared_deck_text = (
        _deck_to_text(ss.slide_deck) if ss.slide_deck else ""
    )
    shared_proposal_text = ""
    shared_proposal_filename = ""
    shared_files_for_agent: list[dict] | None = None
    if (
        ss.client_eval_source == "uploaded_file"
        and uploaded_proposal is not None
    ):
        shared_proposal_text, shared_files_for_agent = (
            _extract_uploaded_proposal(uploaded_proposal)
        )
        shared_proposal_filename = uploaded_proposal.name

    run_disabled = (
        ss.client_eval_source == "uploaded_file"
        and uploaded_proposal is None
    )
    qc_pdf_only = bool(
        ss.client_eval_source == "uploaded_file"
        and uploaded_proposal is not None
        and not shared_proposal_text
        and shared_files_for_agent
    )
    qc_disabled = run_disabled or qc_pdf_only
    if run_disabled:
        st.caption(
            "Upload a final deliverable above to enable **Run Consistency Check** "
            "(or optional buyer/loss reviews below)."
        )
    elif qc_pdf_only:
        st.caption(
            "This upload has no extractable text — use the rendered deck, "
            "HTML, or PPTX upload (not PDF-only)."
        )

    if ss.evals_bundle is not None:
        _render_evals_summary_strip(ss.evals_bundle)
        try:
            bundle_json = json.dumps(
                ss.evals_bundle.model_dump(),
                indent=2,
                ensure_ascii=False,
            )
            st.download_button(
                "Download consistency eval bundle (JSON)",
                data=bundle_json.encode("utf-8"),
                file_name=f"evals-bundle-{ss.run_id}.json",
                mime="application/json",
                key="dl_evals_bundle_json",
            )
        except Exception as exc:
            st.warning(f"Evals bundle download failed: {exc}")

    # ─── Consistency evaluation (primary) ────────────────────────────
    st.markdown(
        "Deterministic internal consistency checks (numbers, dates, "
        "scope, units, terminology, abbreviations, claim–support). "
        "Optional LLM judge resolves ambiguous terminology flags."
    )
    ss.consistency_eval_use_llm_judge = st.checkbox(
        "Run LLM judge on ambiguous flags",
        value=bool(ss.consistency_eval_use_llm_judge),
        key="consistency_eval_use_llm_judge_cb",
    )
    _render_consistency_weights_ui()

    if st.button(
        "Run Consistency Check",
        type="primary",
        key="run_consistency_eval",
        disabled=qc_disabled,
    ):
        try:
            with st.spinner("Running consistency QC…"):
                _run_consistency_eval(
                    deck_text=shared_deck_text,
                    proposal_upload_text=shared_proposal_text,
                    uploaded_filename=shared_proposal_filename,
                    use_llm_judge=ss.consistency_eval_use_llm_judge,
                )
            st.rerun()
        except Exception as exc:
            traceback.print_exc()
            log_event(ss.run_id, f"[CONSISTENCY QC] ERROR: {exc}")
            st.error(f"Consistency QC failed: {exc}")

    if ss.consistency_eval_done and ss.consistency_eval_result:
        cresult = ss.consistency_eval_result
        _consistency_pass_badge(cresult.passed)
        st.metric("Consistency score", f"{cresult.overall_score} / 100")
        st.caption(
            f"Critical: {cresult.summary.critical} · "
            f"Major: {cresult.summary.major} · "
            f"Minor: {cresult.summary.minor} · "
            f"Threshold: {cresult.threshold or ss.consistency_eval_pass_threshold}"
        )
        if getattr(cresult, "score_breakdown", None):
            bd = {
                k: v
                for k, v in cresult.score_breakdown.items()
                if v and v > 0
            }
            if bd:
                from lop_eval.scoring import CHECKER_LABELS

                rows = [
                    {
                        "Check": CHECKER_LABELS.get(k, k),
                        "Points deducted": round(v, 1),
                        "Issues": cresult.issue_counts_by_checker.get(k, 0),
                        "Weight": cresult.checker_weights_used.get(k, 1.0),
                    }
                    for k, v in sorted(bd.items(), key=lambda x: -x[1])
                ]
                st.dataframe(rows, use_container_width=True, hide_index=True)
        if cresult.issues:
            for issue in cresult.issues:
                marker = _consistency_severity_marker(
                    issue.severity.value
                    if hasattr(issue.severity, "value")
                    else str(issue.severity)
                )
                with st.container(border=True):
                    st.markdown(
                        f"{marker} **{issue.type}** — {issue.description}"
                    )
                    st.caption(f"Offending: {issue.offending_text[:240]}")
                    st.caption(f"Expected: {issue.expected_text_or_rule[:240]}")
        else:
            st.caption("No consistency issues flagged.")

        qc_dl, qc_rr = st.columns([3, 1])
        with qc_dl:
            try:
                qc_json = cresult.model_dump_json(indent=2)
                st.download_button(
                    "Download consistency report (JSON)",
                    data=qc_json.encode("utf-8"),
                    file_name=f"consistency-qc-{ss.run_id}.json",
                    mime="application/json",
                    key="dl_consistency_eval_json",
                )
            except Exception as exc:
                st.warning(f"JSON download failed: {exc}")
        with qc_rr:
            if st.button("Re-run", key="rerun_consistency_eval"):
                ss.consistency_eval_done = False
                ss.consistency_eval_result = None
                _sync_evals_bundle()
                st.rerun()


    # ─── Optional buyer / loss-risk (not consistency eval) ─────────
    with st.expander(
        "Optional — buyer & loss-risk reviews (separate from consistency eval)",
        expanded=False,
    ):
        st.caption(
            "LLM reads from a buyer or red-team angle — not internal consistency QC."
        )
        opt_tab_client, opt_tab_loss = st.tabs(
            ["Buyer perspective", "Loss-risk red team"]
        )
        with opt_tab_client:
                st.markdown(
                    "Plays the role of the company owner / business sponsor "
                    "receiving this LoP and grades the proposal from the "
                    "buyer's perspective — RFP coverage, owner priorities "
                    "(explicit + inferred), per-chapter view, reasonableness "
                    "of timeline / fees / team / approach, top concerns, "
                    "missing items, and recommended changes."
                )

                client_eval_model = _render_model_selector("client_eval")

                if st.button(
                    "Run Client Evaluator",
                    type="primary",
                    key="run_client_eval",
                    disabled=run_disabled,
                ):
                    try:
                        log_agent_start(
                            ss.run_id,
                            "client-evaluator-agent",
                            f"source={ss.client_eval_source} | "
                            f"model={client_eval_model}",
                        )

                        user_msg = _build_client_eval_payload(
                            intake=ss.intake_package,
                            synthesis=ss.synthesis_doc,
                            validation=ss.validation_report,
                            dotdash=ss.dotdash_doc,
                            deck_text=shared_deck_text,
                            proposal_upload_text=shared_proposal_text,
                            proposal_upload_filename=shared_proposal_filename,
                        )

                        with st.spinner(
                            "Running client evaluator — reading the proposal "
                            "end-to-end…"
                        ):
                            result = run_agent(
                                "client-evaluator-agent",
                                user_message=user_msg,
                                files=shared_files_for_agent,
                                use_extended_thinking=True,
                                model=client_eval_model,
                            )

                        report = ClientEvaluationReport.model_validate(result)

                        log_event(
                            ss.run_id,
                            f"[CLIENT EVALUATOR] Verdict: {report.overall_verdict} "
                            f"(score {report.score}/100)",
                        )
                        log_event(
                            ss.run_id,
                            f"[CLIENT EVALUATOR] Takeaway: "
                            f"{report.headline_takeaway[:200]}",
                        )
                        log_event(
                            ss.run_id,
                            f"[CLIENT EVALUATOR] RFP coverage: "
                            f"{sum(1 for x in report.rfp_coverage if x.status == 'covered')} covered, "
                            f"{sum(1 for x in report.rfp_coverage if x.status == 'partial')} partial, "
                            f"{sum(1 for x in report.rfp_coverage if x.status == 'missing')} missing",
                        )
                        for concern in report.top_concerns[:5]:
                            log_event(
                                ss.run_id, f"[CLIENT EVALUATOR]     concern: {concern}"
                            )
                        log_agent_result(
                            ss.run_id, "client-evaluator-agent", result
                        )

                        ss.client_eval_report = report
                        ss.client_eval_done = True
                        _sync_evals_bundle()
                        st.rerun()
                    except Exception as exc:
                        traceback.print_exc()
                        log_event(
                            ss.run_id, f"[CLIENT EVALUATOR AGENT] ERROR: {exc}"
                        )
                        st.error(f"Client evaluator failed: {exc}")

                # Done view: full results + re-run button (clears only this eval).
                if ss.client_eval_done and ss.client_eval_report:
                    report = ss.client_eval_report

                    # Headline strip — verdict pill + score + takeaway.
                    top_left, top_right = st.columns([3, 1])
                    with top_left:
                        _client_eval_verdict_badge(report.overall_verdict)
                        st.markdown(f"_{report.headline_takeaway}_")
                    with top_right:
                        st.metric("Buyer score", f"{report.score} / 100")

                    # Source note: which input the agent actually graded.
                    source_note = (
                        "Evaluated the rendered deck from Step 8."
                        if ss.client_eval_source == "in_app_deck"
                        else f"Evaluated uploaded final deliverable: "
                             f"`{ss.client_eval_upload_name or 'uploaded file'}`."
                    )
                    st.caption(source_note)

                    # RFP coverage and owner priorities — side by side.
                    col_l, col_r = st.columns(2)

                    with col_l:
                        st.markdown("##### RFP coverage")
                        if report.rfp_coverage:
                            for item in report.rfp_coverage:
                                marker = _client_eval_status_marker(item.status)
                                with st.container(border=True):
                                    st.markdown(
                                        f"{marker} **{item.status.upper()}** — "
                                        f"{item.requirement}"
                                    )
                                    if item.evidence:
                                        st.caption(f"Evidence: {item.evidence}")
                                    if item.concern:
                                        st.caption(f"Concern: {item.concern}")
                        else:
                            st.caption(
                                "No RFP requirements were extracted in the "
                                "intake — evaluator skipped this check."
                            )

                    with col_r:
                        st.markdown("##### Owner priorities")
                        if report.owner_priorities:
                            for prio in report.owner_priorities:
                                marker = "🟢" if prio.addressed else "🔴"
                                src_tag = prio.source.replace("_", " ")
                                with st.container(border=True):
                                    st.markdown(
                                        f"{marker} **{src_tag}** — {prio.priority}"
                                    )
                                    if prio.evidence:
                                        st.caption(f"Evidence: {prio.evidence}")
                                    if prio.concern:
                                        st.caption(f"Concern: {prio.concern}")
                        else:
                            st.caption("No owner priorities surfaced.")

                    # Chapter assessment and reasonableness — side by side.
                    chap_col, reason_col = st.columns(2)

                    with chap_col:
                        st.markdown("##### Chapter assessment (owner's view)")
                        if report.chapter_assessment:
                            for ch in report.chapter_assessment:
                                marker = _client_eval_status_marker(ch.verdict)
                                with st.container(border=True):
                                    st.markdown(
                                        f"{marker} **{ch.chapter}** — "
                                        f"{ch.verdict}"
                                    )
                                    if ch.client_view:
                                        st.caption(ch.client_view)
                        else:
                            st.caption("No chapter-level assessment returned.")

                    with reason_col:
                        st.markdown("##### Reasonableness")
                        checks = [
                            ("Timeline", report.timeline_check),
                            ("Fees",     report.fees_check),
                            ("Team",     report.team_check),
                            ("Approach", report.approach_check),
                        ]
                        for label, chk in checks:
                            marker = _client_eval_status_marker(chk.verdict)
                            with st.container(border=True):
                                st.markdown(
                                    f"{marker} **{label}** — {chk.verdict}"
                                )
                                if chk.concern:
                                    st.caption(chk.concern)

                    # Narrative quality.
                    if report.quality_assessment:
                        st.markdown("##### Narrative quality (owner's read)")
                        st.write(report.quality_assessment)

                    # Lists at the bottom.
                    bot_l, bot_m, bot_r = st.columns(3)
                    with bot_l:
                        st.markdown("##### Top concerns")
                        if report.top_concerns:
                            for c in report.top_concerns:
                                st.markdown(f"- {c}")
                        else:
                            st.caption("No concerns surfaced.")
                    with bot_m:
                        st.markdown("##### Missing for the owner")
                        if report.missing_for_owner:
                            for m in report.missing_for_owner:
                                st.markdown(f"- {m}")
                        else:
                            st.caption("Nothing material missing.")
                    with bot_r:
                        st.markdown("##### Recommended changes")
                        if report.recommended_changes:
                            for r_item in report.recommended_changes:
                                st.markdown(f"- {r_item}")
                        else:
                            st.caption("No edits recommended.")

                    # JSON download + re-run button (clears only client_eval).
                    dl_col, rerun_col = st.columns([3, 1])
                    with dl_col:
                        try:
                            report_json = json.dumps(
                                report.model_dump(),
                                indent=2,
                                ensure_ascii=False,
                            )
                            st.download_button(
                                "Download evaluation report (JSON)",
                                data=report_json.encode("utf-8"),
                                file_name=(
                                    f"client-evaluation-{ss.run_id}.json"
                                ),
                                mime="application/json",
                                key="dl_client_eval_json",
                            )
                        except Exception as exc:
                            st.warning(f"JSON download packaging failed: {exc}")
                    with rerun_col:
                        if st.button(
                            "Re-run", key="rerun_client_eval_done"
                        ):
                            ss.client_eval_done = False
                            ss.client_eval_report = None
                            _sync_evals_bundle()
                            st.rerun()

        with opt_tab_loss:
                st.markdown(
                    "An independent critical reviewer asking one question: "
                    "**why would we lose this proposal** (or to whom). "
                    "Produces ranked loss reasons, per-competitor angles, "
                    "a save-or-kill verdict with a 0-100 loss-risk score, "
                    "and specific key improvements tied to each loss reason."
                )

                loss_question_text, loss_framing_hint = _loss_eval_framing_preview(
                    ss.intake_package,
                    ss.loss_eval_competitor_override,
                )
                st.markdown(loss_question_text)

                override_value = st.text_input(
                    "Stress-test against a specific competitor (optional — "
                    "overrides the intake-derived list)",
                    value=ss.loss_eval_competitor_override,
                    key="loss_eval_competitor_override_input",
                    placeholder="e.g. BCG, Bain, Deloitte, Accenture, or leave empty",
                )
                if override_value != ss.loss_eval_competitor_override:
                    ss.loss_eval_competitor_override = override_value
                    # Re-compute the framing once the override changes so the
                    # next render shows the new question above the input. We do
                    # NOT auto-rerun the agent on override changes — the partner
                    # explicitly clicks Run.
                    st.rerun()

                loss_eval_model = _render_model_selector("loss_eval")

                if st.button(
                    "Run Loss-Risk Analysis",
                    type="primary",
                    key="run_loss_eval",
                    disabled=run_disabled,
                ):
                    try:
                        log_agent_start(
                            ss.run_id,
                            "loss-analysis-agent",
                            f"source={ss.client_eval_source} | "
                            f"competitor_override="
                            f"{ss.loss_eval_competitor_override or '(none)'} | "
                            f"model={loss_eval_model}",
                        )

                        user_msg = _build_loss_eval_payload(
                            intake=ss.intake_package,
                            synthesis=ss.synthesis_doc,
                            validation=ss.validation_report,
                            dotdash=ss.dotdash_doc,
                            deck_text=shared_deck_text,
                            proposal_upload_text=shared_proposal_text,
                            proposal_upload_filename=shared_proposal_filename,
                            competitor_override=ss.loss_eval_competitor_override,
                            framing_hint=loss_framing_hint,
                        )

                        with st.spinner(
                            "Running loss-risk analysis — reading the proposal "
                            "with the rival's pitch coach in mind…"
                        ):
                            result = run_agent(
                                "loss-analysis-agent",
                                user_message=user_msg,
                                files=shared_files_for_agent,
                                use_extended_thinking=True,
                                model=loss_eval_model,
                            )

                        report = LossAnalysisReport.model_validate(result)

                        log_event(
                            ss.run_id,
                            f"[LOSS ANALYSIS] Likelihood: {report.loss_likelihood} "
                            f"(risk {report.loss_risk_score}/100) | "
                            f"verdict: {report.save_or_kill_verdict}",
                        )
                        log_event(
                            ss.run_id,
                            f"[LOSS ANALYSIS] Framing: {report.framing_question}",
                        )
                        log_event(
                            ss.run_id,
                            f"[LOSS ANALYSIS] Punchline: {report.punchline[:240]}",
                        )
                        for reason in report.top_loss_reasons[:5]:
                            log_event(
                                ss.run_id,
                                f"[LOSS ANALYSIS]     loss reason "
                                f"[{reason.severity}/{reason.category}]: "
                                f"{reason.reason}",
                            )
                        log_agent_result(
                            ss.run_id, "loss-analysis-agent", result
                        )

                        ss.loss_eval_report = report
                        ss.loss_eval_done = True
                        _sync_evals_bundle()
                        st.rerun()
                    except Exception as exc:
                        traceback.print_exc()
                        log_event(
                            ss.run_id, f"[LOSS ANALYSIS AGENT] ERROR: {exc}"
                        )
                        st.error(f"Loss-risk analysis failed: {exc}")

                # Done view for the loss tab — only renders when a report exists.
                if ss.loss_eval_done and ss.loss_eval_report:
                    loss_report: LossAnalysisReport = ss.loss_eval_report

                    # Framing question + headline strip.
                    st.markdown(f"_{loss_report.framing_question}_")
                    head_l, head_m, head_r = st.columns([2, 1, 1])
                    with head_l:
                        _loss_eval_verdict_badge(loss_report.save_or_kill_verdict)
                        _loss_eval_likelihood_badge(loss_report.loss_likelihood)
                    with head_m:
                        st.metric(
                            "Loss-risk score",
                            f"{loss_report.loss_risk_score} / 100",
                        )
                    with head_r:
                        if loss_report.primary_competitors:
                            st.caption("Analysed against:")
                            st.markdown(
                                ", ".join(
                                    f"`{c}`" for c in loss_report.primary_competitors
                                )
                            )

                    # Source note + partner-voice punchline.
                    source_note = (
                        "Evaluated the rendered deck from Step 8."
                        if ss.client_eval_source == "in_app_deck"
                        else f"Evaluated uploaded final deliverable: "
                             f"`{ss.client_eval_upload_name or 'uploaded file'}`."
                    )
                    st.caption(source_note)
                    if loss_report.punchline:
                        st.info(f"**Partner voice:** {loss_report.punchline}")

                    # Top loss reasons + competitor angles — side by side.
                    loss_l, loss_r = st.columns(2)

                    with loss_l:
                        st.markdown("##### Top loss reasons")
                        if loss_report.top_loss_reasons:
                            for reason in loss_report.top_loss_reasons:
                                marker = _loss_eval_severity_marker(reason.severity)
                                with st.container(border=True):
                                    st.markdown(
                                        f"{marker} **{reason.severity.upper()}** "
                                        f"· `{reason.category}` — {reason.reason}"
                                    )
                                    if reason.evidence:
                                        st.caption(f"Evidence: {reason.evidence}")
                        else:
                            st.caption("No loss reasons surfaced.")

                    with loss_r:
                        st.markdown("##### Competitor angles")
                        if loss_report.competitor_angles:
                            for angle in loss_report.competitor_angles:
                                marker = _loss_eval_severity_marker(angle.severity)
                                with st.container(border=True):
                                    st.markdown(
                                        f"{marker} **{angle.competitor}** — "
                                        f"{angle.competitor_strength}"
                                    )
                                    if angle.where_it_lands:
                                        st.caption(
                                            f"Where it lands: "
                                            f"{angle.where_it_lands}"
                                        )
                                    if angle.model_knowledge_note:
                                        st.caption(
                                            "Note (validate with partner): "
                                            f"{angle.model_knowledge_note}"
                                        )
                        else:
                            st.caption(
                                "No competitor angles — likely a "
                                "non-competitive pursuit."
                            )

                    # Vulnerable chapters as chips.
                    if loss_report.vulnerable_chapters:
                        st.markdown("##### Vulnerable chapters (worst first)")
                        st.markdown(
                            "  ".join(
                                f"`{c}`" for c in loss_report.vulnerable_chapters
                            )
                        )

                    # Key improvements as ordered cards.
                    st.markdown("##### Key improvements (ranked)")
                    if loss_report.key_improvements:
                        for i, imp in enumerate(loss_report.key_improvements, 1):
                            marker = _loss_eval_priority_marker(imp.priority)
                            with st.container(border=True):
                                st.markdown(
                                    f"**{i}. {marker} {imp.priority.upper()}** "
                                    f"· `{imp.linked_chapter}` — {imp.improvement}"
                                )
                                if imp.expected_impact:
                                    st.caption(
                                        f"Closes: {imp.expected_impact}"
                                    )
                    else:
                        st.caption("No improvements proposed.")

                    # JSON download + re-run (clears only loss_eval).
                    dl_col, rerun_col = st.columns([3, 1])
                    with dl_col:
                        try:
                            loss_json = json.dumps(
                                loss_report.model_dump(),
                                indent=2,
                                ensure_ascii=False,
                            )
                            st.download_button(
                                "Download loss-analysis report (JSON)",
                                data=loss_json.encode("utf-8"),
                                file_name=f"loss-analysis-{ss.run_id}.json",
                                mime="application/json",
                                key="dl_loss_eval_json",
                            )
                        except Exception as exc:
                            st.warning(f"JSON download packaging failed: {exc}")
                    with rerun_col:
                        if st.button(
                            "Re-run", key="rerun_loss_eval_done"
                        ):
                            ss.loss_eval_done = False
                            ss.loss_eval_report = None
                            _sync_evals_bundle()
                            st.rerun()

    st.divider()
