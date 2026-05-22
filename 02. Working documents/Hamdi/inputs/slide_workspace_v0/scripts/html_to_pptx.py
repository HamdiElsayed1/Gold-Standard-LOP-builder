"""Compile HTML slides to a native PPTX file.

Renders each HTML slide in a headless browser (Playwright), extracts
element positions and computed styles via extract_positions.js, then
creates native PPTX shapes on the appropriate template layout.

Usage:
    python html_to_pptx.py <input> [-t <template>] [-o <output.pptx>]

    <input>   Path to a single .html file, or a directory of .html files
              (one per slide, sorted alphabetically for slide order).
    -t        Path to the PPTX template (default: pptx-template.pptx in
              the workspace root).
    -o        Output PPTX path (default: output.pptx in the current dir).

Supported data-pptx tags (optional except where noted):
    rect, ellipse, line, textbox, image, table  — auto-detected when omitted
    placeholder  — required (injects text into template placeholder)
    chrome       — required (marks template-replica elements to skip)

The extractor auto-discovers shapes from visible elements (background,
border, images, text) so data-pptx is only needed for overrides.
"""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

from lxml import etree
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.util import Pt, Emu

# Ensure the scripts/ directory is importable
sys.path.insert(0, str(Path(__file__).resolve().parent))

from shared.constants import (
    EXTRACT_JS_PATH,
    FOOTNOTE_BOX_HEIGHT,
    FOOTNOTE_BOX_LEFT,
    FOOTNOTE_BOX_TOP,
    FOOTNOTE_BOX_WIDTH,
    LAYOUT_INDEX,
    PPTX_TAGS,
    SHAPE_TYPE_MAP,
    SOURCE_BOX_HEIGHT,
    SOURCE_BOX_LEFT,
    SOURCE_BOX_TOP,
    SOURCE_BOX_WIDTH,
    TEMPLATE_PATH,
    get_template_path,
)
from shared.emu import px_to_emu_x, px_to_emu_y, CSS_PX_TO_PT
from chart_builder import add_chart as _add_chart, patch_total_labels_xlsx


# ---------------------------------------------------------------------------
# Color / style helpers
# ---------------------------------------------------------------------------

_RGB_RE = re.compile(
    r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)(?:\s*,\s*[\d.]+)?\s*\)"
)
_HEX_RE = re.compile(r"^#?([0-9a-fA-F]{3,8})$")


def _parse_color(css_color: str) -> RGBColor | None:
    """Convert a CSS color string to an RGBColor, or None if transparent."""
    if not css_color:
        return None
    css_color = css_color.strip()
    if css_color in ("transparent", "rgba(0, 0, 0, 0)"):
        return None

    m = _RGB_RE.match(css_color)
    if m:
        return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))

    m = _HEX_RE.match(css_color)
    if m:
        h = m.group(1)
        if len(h) == 3:
            h = h[0] * 2 + h[1] * 2 + h[2] * 2
        elif len(h) in (4, 8):
            h = h[:6]
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    return None


def _parse_font_size(css_size: str) -> float | None:
    """Extract numeric pt size from a CSS fontSize string (e.g. '50px').

    Viewport is 1920 px for a 13.333-inch slide → 144 effective DPI.
    1 CSS-px = 0.5 PPTX-pt at this DPI (CSS_PX_TO_PT).
    """
    if not css_size:
        return None
    m = re.match(r"([\d.]+)\s*px", css_size)
    if m:
        return float(m.group(1)) * CSS_PX_TO_PT
    m = re.match(r"([\d.]+)\s*pt", css_size)
    if m:
        return float(m.group(1))
    return None


def _parse_border_width(css_width: str) -> float:
    m = re.match(r"([\d.]+)\s*px", css_width or "")
    return float(m.group(1)) if m else 0


def _clean_font_family(css_family: str) -> str:
    """Return the first font name, stripping quotes and fallbacks."""
    if not css_family:
        return "Arial"
    first = css_family.split(",")[0].strip().strip("'\"")
    return first or "Arial"


# ---------------------------------------------------------------------------
# Text helpers
# ---------------------------------------------------------------------------

def _apply_run_style(run, style: dict, skip_font: bool = False) -> None:
    """Apply font styling from an extracted run style dict.

    When *skip_font* is True, font name/size are not set so that
    placeholder formatting inherited from the layout is preserved.
    """
    font = run.font

    if not skip_font:
        family = _clean_font_family(style.get("fontFamily", ""))
        if family:
            font.name = family
        size = _parse_font_size(style.get("fontSize", ""))
        if size:
            font.size = Pt(size)

    weight = style.get("fontWeight", "")
    if weight in ("bold", "700", "800", "900"):
        font.bold = True

    if style.get("fontStyle") == "italic":
        font.italic = True

    color = _parse_color(style.get("color", ""))
    if color:
        font.color.rgb = color

    if style.get("superscript"):
        from lxml import etree
        rPr = run._r.get_or_add_rPr()
        rPr.set("baseline", "30000")
    elif style.get("subscript"):
        from lxml import etree
        rPr = run._r.get_or_add_rPr()
        rPr.set("baseline", "-25000")


_BULLET_CHARS = {"\u2022", "\u2013", "\u2014", "\u25cf", "\u25cb", "\u25aa", "\u25a0"}


def _strip_bullet_prefix(text: str) -> tuple[str, bool]:
    """Strip leading bullet character + space. Return (cleaned, was_bullet)."""
    if len(text) >= 2 and text[0] in _BULLET_CHARS and text[1] == " ":
        return text[2:], True
    return text, False


def _set_paragraph_bullet(para, char: str = "\u2022") -> None:
    """Enable bullet formatting on a paragraph with explicit bullet char.

    Textboxes don't inherit the slide master's list styles, so we must
    set the bullet character explicitly.
    """
    from lxml import etree
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    pPr = para._p.get_or_add_pPr()
    pPr.set("marL", "228600")
    pPr.set("indent", "-228600")
    existing = pPr.find(f"{{{ns_a}}}buChar")
    if existing is None:
        buChar = etree.SubElement(pPr, f"{{{ns_a}}}buChar")
        buChar.set("char", char)


def _set_vertical_anchor(text_frame, anchor: str) -> None:
    """Set the vertical text anchor on bodyPr. anchor is 't', 'ctr', or 'b'."""
    from lxml import etree
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    txBody = text_frame._txBody
    bodyPr = txBody.find(f"{{{ns_a}}}bodyPr")
    if bodyPr is None:
        bodyPr = etree.SubElement(txBody, f"{{{ns_a}}}bodyPr")
        txBody.insert(0, bodyPr)
    bodyPr.set("anchor", anchor)


def _infer_valign(shape_desc: dict) -> str | None:
    """Infer PPTX bodyPr anchor (t/ctr/b) from computed flex alignment.

    For column flex: vertical alignment is controlled by justifyContent.
    For row flex (default): vertical alignment is controlled by alignItems.
    """
    styles = shape_desc.get("styles", {})
    fd = styles.get("flexDirection", "") or ""
    jc = styles.get("justifyContent", "") or ""
    ai = styles.get("alignItems", "") or ""
    if "column" in fd:
        if "flex-end" in jc:
            return "b"
        if "center" in jc:
            return "ctr"
    else:
        if "flex-end" in ai:
            return "b"
        if "center" in ai:
            return "ctr"
    return None


def _css_white_space_disables_wrap(styles: dict) -> bool:
    """Return True when CSS white-space requests no automatic wrapping."""
    white_space = (styles.get("whiteSpace", "") or "").strip().lower()
    # CSS values that suppress automatic wrapping of normal spaces.
    return white_space in {"nowrap", "pre"}


def _populate_text_frame(
    text_frame, paragraphs: list[dict], styles: dict,
    skip_font: bool = False,
) -> None:
    """Fill a text_frame with paragraph/run data from extraction."""
    if not paragraphs:
        return

    text_frame.word_wrap = not _css_white_space_disables_wrap(styles)

    alignment = styles.get("textAlign", "")
    from pptx.enum.text import PP_ALIGN
    align_map = {
        "start": PP_ALIGN.LEFT,
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "end": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }

    for i, para_data in enumerate(paragraphs):
        if i == 0:
            para = text_frame.paragraphs[0]
        else:
            para = text_frame.add_paragraph()

        if alignment in align_map:
            para.alignment = align_map[alignment]

        is_bullet = para_data.get("bullet", False)
        first_run_done = False

        for j, run_data in enumerate(para_data.get("runs", [])):
            text = run_data.get("text", "")

            if not first_run_done and is_bullet:
                text, _ = _strip_bullet_prefix(text)
                first_run_done = True

            if j == 0:
                run = para.runs[0] if para.runs else para.add_run()
            else:
                run = para.add_run()
            run.text = text
            _apply_run_style(run, run_data.get("style", {}), skip_font=skip_font)

        if is_bullet:
            _set_paragraph_bullet(para)


# ---------------------------------------------------------------------------
# Shape builders
# ---------------------------------------------------------------------------

def _parse_css_px(value: str) -> float:
    """Parse a CSS pixel value like '20px' and return the numeric part."""
    m = re.match(r"([\d.]+)\s*px", value or "")
    return float(m.group(1)) if m else 0.0


def _set_text_frame_margins(tf, styles: dict) -> None:
    """Set text frame internal margins from extracted CSS padding values."""
    tf.margin_top = Emu(px_to_emu_y(_parse_css_px(styles.get("paddingTop", ""))))
    tf.margin_bottom = Emu(px_to_emu_y(_parse_css_px(styles.get("paddingBottom", ""))))
    tf.margin_left = Emu(px_to_emu_x(_parse_css_px(styles.get("paddingLeft", ""))))
    tf.margin_right = Emu(px_to_emu_x(_parse_css_px(styles.get("paddingRight", ""))))


def _add_rect_or_ellipse(slide, shape_desc: dict, mso_shape) -> None:
    r = shape_desc["rect"]
    left = px_to_emu_x(r["x"])
    top = px_to_emu_y(r["y"])
    width = px_to_emu_x(r["width"])
    height = px_to_emu_y(r["height"])

    shape = slide.shapes.add_shape(mso_shape, left, top, width, height)
    styles = shape_desc.get("styles", {})

    fill_color = _parse_color(styles.get("backgroundColor", ""))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    border_width = _parse_border_width(styles.get("borderWidth", ""))
    border_color = _parse_color(styles.get("borderColor", ""))
    if border_width > 0 and border_color:
        shape.line.width = Pt(border_width * CSS_PX_TO_PT)
        shape.line.color.rgb = border_color
        border_style = styles.get("borderStyle", "solid")
        if border_style in ("dashed", "dotted"):
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            shape.line.dash_style = MSO_LINE_DASH_STYLE.SQUARE_DOT
    else:
        shape.line.fill.background()

    tf = shape.text_frame
    _set_text_frame_margins(tf, styles)

    valign = _infer_valign(shape_desc)
    text_data = shape_desc.get("text", [])
    if text_data:
        _populate_text_frame(tf, text_data, styles)
    if valign:
        _set_vertical_anchor(tf, valign)


def _add_line(slide, shape_desc: dict) -> None:
    endpoints = shape_desc.get("line")
    if not endpoints:
        return

    x1 = px_to_emu_x(endpoints["x1"])
    y1 = px_to_emu_y(endpoints["y1"])
    x2 = px_to_emu_x(endpoints["x2"])
    y2 = px_to_emu_y(endpoints["y2"])

    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2
    )

    styles = shape_desc.get("styles", {})
    stroke_color = _parse_color(styles.get("stroke", ""))
    if stroke_color:
        connector.line.color.rgb = stroke_color

    stroke_w = styles.get("strokeWidth", "")
    w_match = re.match(r"([\d.]+)", stroke_w or "")
    if w_match:
        connector.line.width = Pt(float(w_match.group(1)) * CSS_PX_TO_PT)


def _add_textbox(slide, shape_desc: dict) -> None:
    r = shape_desc["rect"]
    left = px_to_emu_x(r["x"])
    top = px_to_emu_y(r["y"])
    width = px_to_emu_x(r["width"])
    height = px_to_emu_y(r["height"])

    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    styles = shape_desc.get("styles", {})
    _set_text_frame_margins(tf, styles)
    text_data = shape_desc.get("text", [])
    if text_data:
        _populate_text_frame(tf, text_data, styles)
    valign = _infer_valign(shape_desc)
    if valign:
        _set_vertical_anchor(tf, valign)


def _add_image(slide, shape_desc: dict, html_dir: Path) -> None:
    r = shape_desc["rect"]
    left = px_to_emu_x(r["x"])
    top = px_to_emu_y(r["y"])
    width = px_to_emu_x(r["width"])
    height = px_to_emu_y(r["height"])

    src = shape_desc.get("src", "")
    if not src:
        return

    img_path = (html_dir / src).resolve()
    if not img_path.is_file():
        print(f"  Warning: image not found: {img_path}", file=sys.stderr)
        return

    slide.shapes.add_picture(str(img_path), left, top, width, height)


def _set_placeholder(slide, shape_desc: dict) -> None:
    idx_str = shape_desc.get("data", {}).get("ph-idx")
    if idx_str is None:
        return
    idx = int(idx_str)

    try:
        ph = slide.placeholders[idx]
    except KeyError:
        return

    text_data = shape_desc.get("text", [])
    if text_data:
        styles = shape_desc.get("styles", {})
        _populate_text_frame(ph.text_frame, text_data, styles, skip_font=True)


def _add_table(slide, shape_desc: dict) -> None:
    r = shape_desc["rect"]
    left = px_to_emu_x(r["x"])
    top = px_to_emu_y(r["y"])
    width = px_to_emu_x(r["width"])
    height = px_to_emu_y(r["height"])

    table_data = shape_desc.get("tableData", [])
    if not table_data:
        return

    rows = len(table_data)
    cols = max(len(row) for row in table_data) if rows else 0
    if rows == 0 or cols == 0:
        return

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    for ri, row in enumerate(table_data):
        for ci, cell_data in enumerate(row):
            if ci < cols:
                table.cell(ri, ci).text = cell_data.get("text", "")


# ---------------------------------------------------------------------------
# Footnote / Source builders (match example.pptx exactly)
# ---------------------------------------------------------------------------

_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _configure_bottom_zone_bodyPr(bodyPr, anchor: str) -> None:
    """Set bodyPr attributes shared by footnote and source boxes."""
    bodyPr.set("wrap", "square")
    bodyPr.set("lIns", "0")
    bodyPr.set("tIns", "0")
    bodyPr.set("rIns", "0")
    bodyPr.set("bIns", "0")
    bodyPr.set("rtlCol", "0")
    bodyPr.set("anchor", anchor)
    if bodyPr.find(f"{{{_NS_A}}}spAutoFit") is None:
        etree.SubElement(bodyPr, f"{{{_NS_A}}}spAutoFit")


def _populate_bottom_zone_text(text_frame, paragraphs: list[dict]) -> None:
    """Fill text into a footnote/source box at 8pt."""
    if not paragraphs:
        return
    for i, para_data in enumerate(paragraphs):
        if i == 0:
            para = text_frame.paragraphs[0]
        else:
            para = text_frame.add_paragraph()
        for j, run_data in enumerate(para_data.get("runs", [])):
            text = run_data.get("text", "")
            if j == 0:
                run = para.runs[0] if para.runs else para.add_run()
            else:
                run = para.add_run()
            run.text = text
            run.font.size = Pt(8)
            style = run_data.get("style", {})
            family = _clean_font_family(style.get("fontFamily", ""))
            if family:
                run.font.name = family
            weight = style.get("fontWeight", "")
            if weight in ("bold", "700", "800", "900"):
                run.font.bold = True
            if style.get("fontStyle") == "italic":
                run.font.italic = True
            color = _parse_color(style.get("color", ""))
            if color:
                run.font.color.rgb = color


def _add_footnote_box(slide, shape_desc: dict) -> None:
    """Create a footnote textbox matching the example.pptx '4. Footnote' shape."""
    txbox = slide.shapes.add_textbox(
        FOOTNOTE_BOX_LEFT, FOOTNOTE_BOX_TOP,
        FOOTNOTE_BOX_WIDTH, FOOTNOTE_BOX_HEIGHT,
    )
    tf = txbox.text_frame
    tf.word_wrap = True

    bodyPr = tf._txBody.find(f"{{{_NS_A}}}bodyPr")
    _configure_bottom_zone_bodyPr(bodyPr, "b")

    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    spPr = txbox._element.find(f"{{{ns_p}}}spPr")
    if spPr is None:
        spPr = txbox._element.find(f"{{{_NS_A}}}spPr")
    if spPr is not None and spPr.find(f"{{{_NS_A}}}noFill") is None:
        noFill = etree.SubElement(spPr, f"{{{_NS_A}}}noFill")

    text_data = shape_desc.get("text", [])
    _populate_bottom_zone_text(tf, text_data)

    for para in tf.paragraphs:
        pPr = para._p.get_or_add_pPr()
        pPr.set("marL", "203200")
        pPr.set("indent", "-212725")
        pPr.set("algn", "l")
        lnSpc = etree.SubElement(pPr, f"{{{_NS_A}}}lnSpc")
        etree.SubElement(lnSpc, f"{{{_NS_A}}}spcPct").set("val", "100000")
        etree.SubElement(pPr, f"{{{_NS_A}}}buNone")


def _add_source_box(slide, shape_desc: dict) -> None:
    """Create a source textbox matching the example.pptx '5. Source' shape."""
    txbox = slide.shapes.add_textbox(
        SOURCE_BOX_LEFT, SOURCE_BOX_TOP,
        SOURCE_BOX_WIDTH, SOURCE_BOX_HEIGHT,
    )
    tf = txbox.text_frame
    tf.word_wrap = True

    bodyPr = tf._txBody.find(f"{{{_NS_A}}}bodyPr")
    bodyPr.set("vert", "horz")
    _configure_bottom_zone_bodyPr(bodyPr, "t")

    # Replace the empty lstStyle with one matching the example
    old_lstStyle = tf._txBody.find(f"{{{_NS_A}}}lstStyle")
    if old_lstStyle is not None:
        idx = list(tf._txBody).index(old_lstStyle)
        tf._txBody.remove(old_lstStyle)
    else:
        idx = 1  # after bodyPr
    lstStyle = etree.Element(f"{{{_NS_A}}}lstStyle")
    tf._txBody.insert(idx, lstStyle)

    lvl1 = etree.SubElement(lstStyle, f"{{{_NS_A}}}lvl1pPr")
    lvl1.set("indent", "0")
    lnSpc = etree.SubElement(lvl1, f"{{{_NS_A}}}lnSpc")
    etree.SubElement(lnSpc, f"{{{_NS_A}}}spcPct").set("val", "100000")
    spcBef = etree.SubElement(lvl1, f"{{{_NS_A}}}spcBef")
    etree.SubElement(spcBef, f"{{{_NS_A}}}spcPts").set("val", "300")
    spcAft = etree.SubElement(lvl1, f"{{{_NS_A}}}spcAft")
    etree.SubElement(spcAft, f"{{{_NS_A}}}spcPts").set("val", "300")
    buFont = etree.SubElement(lvl1, f"{{{_NS_A}}}buFont")
    buFont.set("typeface", "Segoe UI")
    buChar = etree.SubElement(lvl1, f"{{{_NS_A}}}buChar")
    buChar.set("char", "\u200b")
    defRPr = etree.SubElement(lvl1, f"{{{_NS_A}}}defRPr")
    defRPr.set("sz", "900")

    text_data = shape_desc.get("text", [])
    _populate_bottom_zone_text(tf, text_data)


# ---------------------------------------------------------------------------
# Main compiler
# ---------------------------------------------------------------------------

def html_to_pptx(
    html_paths: list[Path],
    template: Path | None = None,
    output: Path = Path("output.pptx"),
) -> None:
    """Compile one or more HTML slide files to a PPTX presentation.

    Args:
        html_paths: Ordered list of .html files (one per slide).
        template: Path to the .pptx template. If None, uses the default
                  template via get_template_path() (handles OneDrive).
        output: Where to write the resulting .pptx file.
    """
    tpl_path = template if template else get_template_path()
    prs = Presentation(str(tpl_path))
    extract_js = EXTRACT_JS_PATH.read_text(encoding="utf-8")

    with sync_playwright() as pw:
        browser = pw.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": 1920, "height": 1080})

        for html_path in html_paths:
            html_path = Path(html_path).resolve()
            file_url = html_path.as_uri()
            page.goto(file_url, wait_until="networkidle")

            result = page.evaluate(extract_js)
            meta = result.get("meta", {})
            shapes = result.get("shapes", [])

            layout_name = meta.get("layout", "Default")
            layout_idx = LAYOUT_INDEX.get(layout_name, 1)
            slide_layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(slide_layout)

            html_dir = html_path.parent

            for shape_desc in shapes:
                tag = shape_desc.get("tag", "")
                if tag == "chrome":
                    continue

                dispatch = SHAPE_TYPE_MAP.get(tag)
                if dispatch == "skip":
                    continue

                if tag == "rect":
                    _add_rect_or_ellipse(slide, shape_desc, MSO_SHAPE.RECTANGLE)
                elif tag == "ellipse":
                    _add_rect_or_ellipse(slide, shape_desc, MSO_SHAPE.OVAL)
                elif tag == "line":
                    _add_line(slide, shape_desc)
                elif tag == "textbox":
                    _add_textbox(slide, shape_desc)
                elif tag == "image":
                    _add_image(slide, shape_desc, html_dir)
                elif tag == "placeholder":
                    _set_placeholder(slide, shape_desc)
                elif tag == "table":
                    _add_table(slide, shape_desc)
                elif tag == "chart":
                    _add_chart(slide, shape_desc)
                elif tag == "footnote":
                    _add_footnote_box(slide, shape_desc)
                elif tag == "source":
                    _add_source_box(slide, shape_desc)

        browser.close()

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    patch_total_labels_xlsx(output)


def _find_html_files(input_path: Path) -> list[Path]:
    """Resolve *input_path* to a sorted list of .html files."""
    if input_path.is_file():
        return [input_path]
    if input_path.is_dir():
        files = sorted(input_path.glob("*.html"))
        if not files:
            print(f"Error: no .html files found in {input_path}", file=sys.stderr)
            sys.exit(1)
        return files
    print(f"Error: {input_path} is not a file or directory", file=sys.stderr)
    sys.exit(1)


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Compile HTML slides to PPTX via Playwright rendering",
    )
    parser.add_argument(
        "input",
        type=Path,
        help="Path to a .html file or directory of .html files",
    )
    parser.add_argument(
        "-t", "--template",
        type=Path,
        default=TEMPLATE_PATH,
        help=f"PPTX template path (default: {TEMPLATE_PATH})",
    )
    parser.add_argument(
        "-o", "--output",
        type=Path,
        default=Path("output.pptx"),
        help="Output PPTX path (default: output.pptx)",
    )
    args = parser.parse_args()

    html_files = _find_html_files(args.input)
    print(f"Compiling {len(html_files)} slide(s) -> {args.output}")
    html_to_pptx(html_files, args.template, args.output)
    print(f"Done. Output: {args.output}")


if __name__ == "__main__":
    main()
