"""Extract plain text from PPTX/PDF for Background Material indexing."""

from __future__ import annotations

import logging
import re
from pathlib import Path

logger = logging.getLogger(__name__)


def extract_text_from_pptx(path: Path) -> str:
    """Concatenate slide titles and shape text from a PowerPoint file."""
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for snum, slide in enumerate(prs.slides, start=1):
        parts.append(f"\n## Slide {snum}\n")
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text and text.strip():
                parts.append(text.strip())
            # Group shapes
            if getattr(shape, "shapes", None):
                for sub in shape.shapes:
                    st = getattr(sub, "text", None)
                    if st and st.strip():
                        parts.append(st.strip())
    return "\n".join(parts).strip()


def extract_text_from_pdf(path: Path) -> str:
    """Extract text per page from PDF."""
    try:
        from pypdf import PdfReader
    except ImportError:
        logger.warning("pypdf not installed; skipping PDF extraction for %s", path)
        return ""

    reader = PdfReader(str(path))
    parts: list[str] = []
    for i, page in enumerate(reader.pages):
        parts.append(f"\n## Page {i + 1}\n")
        try:
            t = page.extract_text() or ""
        except Exception as exc:  # noqa: BLE001
            logger.warning("PDF page %s extract failed for %s: %s", i + 1, path, exc)
            t = ""
        parts.append(t.strip())
    return "\n".join(parts).strip()


def safe_cache_stem(path: Path) -> str:
    """Filesystem-safe stem for cache filenames."""
    stem = path.stem
    stem = re.sub(r'[<>:"/\\|?*]', "_", stem)
    return stem[:200] if len(stem) > 200 else stem


def ensure_extracted_markdown(
    source_path: Path,
    extracted_root: Path,
    *,
    force: bool = False,
) -> Path | None:
    """
    If ``source_path`` is pptx/pdf, write/update ``extracted_root/<stem>.md`` and return that path.

    Uses source mtime vs cached file mtime for invalidation.
    """
    suffix = source_path.suffix.lower()
    if suffix not in {".pptx", ".pdf"}:
        return None

    extracted_root.mkdir(parents=True, exist_ok=True)
    out = extracted_root / f"{safe_cache_stem(source_path)}.md"

    try:
        src_mtime = source_path.stat().st_mtime
    except OSError:
        return None

    if out.is_file() and not force:
        try:
            if out.stat().st_mtime >= src_mtime:
                return out
        except OSError:
            pass

    if suffix == ".pptx":
        body = extract_text_from_pptx(source_path)
    else:
        body = extract_text_from_pdf(source_path)

    header = (
        f"<!-- extracted from: {source_path.name} -->\n"
        f"<!-- source_uri: {source_path.resolve()} -->\n\n"
    )
    out.write_text(header + (body or "(no text extracted)"), encoding="utf-8")
    return out
